# -*- coding: utf-8 -*-
"""
BookBridge — Business storytelling deck generator.
Builds a single .pptx that walks an audience from problem -> vision ->
users -> the product journey -> every feature module in detail -> the
trust differentiator -> architecture/data -> quality -> roadmap.

Run: python build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ----------------------------------------------------------------------------
# Design system
# ----------------------------------------------------------------------------
EMU = 914400
SW, SH = Inches(13.333), Inches(7.5)  # 16:9

# Palette — warm, trust-driven, "books + community"
INK      = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
INK_SOFT = RGBColor(0x3D, 0x3D, 0x56)
PAPER    = RGBColor(0xFB, 0xF8, 0xF1)   # warm paper
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xE6, 0xDF, 0xD2)
MUTED    = RGBColor(0x76, 0x70, 0x66)

BRAND    = RGBColor(0x2E, 0x6F, 0x5B)   # forest green (sustainability)
BRAND_DK = RGBColor(0x1E, 0x4D, 0x3E)
BRAND_LT = RGBColor(0xE3, 0xEF, 0xE9)
AMBER    = RGBColor(0xE0, 0x8E, 0x2C)   # warm accent
AMBER_LT = RGBColor(0xFA, 0xEC, 0xD6)
CORAL    = RGBColor(0xD9, 0x5D, 0x4A)
BLUE     = RGBColor(0x3B, 0x6E, 0xA5)
PURPLE   = RGBColor(0x6B, 0x4E, 0x9E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

FONT  = "Calibri"
FONTH = "Calibri"

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)

def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def bg(s, color=PAPER):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _set_fill(r, color)
    r.shadow.inherit = False
    _send_back(r)
    return r

def _send_back(shape):
    sp = shape._element
    sp.getparent().remove(sp)
    # insert as first child after non-shape elements
    parent = shape._element  # noop guard
    return shape

def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    r = s.shapes.add_shape(shape, x, y, w, h)
    if color is None:
        r.fill.background()
    else:
        r.fill.solid()
        r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = line_w or Pt(1)
    r.shadow.inherit = False
    return r

def round_rect(s, x, y, w, h, color, line=None, line_w=None, radius=0.08):
    r = rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=line, line_w=line_w)
    try:
        r.adjustments[0] = radius
    except Exception:
        pass
    return r

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(2), line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb

def R(txt, size, color=INK, bold=False, italic=False, font=FONT):
    return (txt, size, color, bold, italic, font)

def P(*runs):
    return list(runs)


# ----------------------------------------------------------------------------
# Decorative components
# ----------------------------------------------------------------------------
def kicker(s, x, y, label, color=BRAND):
    rect(s, x, y + Pt(2), Inches(0.32), Pt(4), color)
    text(s, x + Inches(0.42), y - Pt(6), Inches(6), Inches(0.3),
         [P(R(label.upper(), 12.5, color, bold=True))])

def page_num(s, n, total=18, section=""):
    text(s, Inches(11.7), Inches(7.02), Inches(1.4), Inches(0.3),
         [P(R(f"{n:02d} / {total}", 9, MUTED))], align=PP_ALIGN.RIGHT)
    if section:
        text(s, Inches(0.6), Inches(7.02), Inches(7), Inches(0.3),
             [P(R("BookBridge", 9, MUTED, bold=True), R("   ·   " + section, 9, MUTED))])

def footer_brand(s):
    rect(s, 0, Inches(7.34), SW, Inches(0.16), BRAND)

def chip(s, x, y, w, label, fill, txt_color, size=10.5, h=Inches(0.34)):
    c = round_rect(s, x, y, w, h, fill, radius=0.5)
    text(s, x, y, w, h, [P(R(label, size, txt_color, bold=True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c

def icon_circle(s, x, y, d, glyph, fill, glyph_color=WHITE, size=18):
    rect(s, x, y, d, d, fill, shape=MSO_SHAPE.OVAL)
    text(s, x, y, d, d, [P(R(glyph, size, glyph_color, bold=True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = slide()
bg(s, INK)
# warm band
rect(s, 0, 0, Inches(5.0), SH, BRAND_DK)
rect(s, Inches(4.7), 0, Inches(0.3), SH, BRAND)

# left: motif of stacked "book spines"
spine_y = Inches(1.5)
spine_colors = [AMBER, BRAND, CORAL, BLUE, PURPLE, BRAND_LT]
sx = Inches(0.85)
for i, col in enumerate(spine_colors):
    hh = Inches(3.4) - Inches(0.16) * i
    round_rect(s, sx + Inches(0.62) * i, SH - Inches(0.9) - hh, Inches(0.46), hh, col, radius=0.12)

text(s, Inches(0.85), Inches(0.7), Inches(3.6), Inches(0.6),
     [P(R("ZOOTOPIA  ·  SE 2025.2", 12, AMBER_LT, bold=True))])

# right: title block
text(s, Inches(5.5), Inches(1.55), Inches(7.2), Inches(2.6),
     [P(R("BookBridge", 60, WHITE, bold=True)),
      P(R("Keep books in circulation,", 26, AMBER_LT, font=FONTH)),
      P(R("not in landfills.", 26, AMBER_LT, font=FONTH))],
     line_spacing=1.04)

text(s, Inches(5.5), Inches(4.35), Inches(7.0), Inches(1.4),
     [P(R("A trust-driven community where book lovers ", 15, RGBColor(0xCF,0xD6,0xD2)),
        R("gift, exchange, or sell", 15, WHITE, bold=True),
        R(" pre-owned books at symbolic prices.", 15, RGBColor(0xCF,0xD6,0xD2)))],
     line_spacing=1.25)

# SDG chips
chip(s, Inches(5.5), Inches(5.55), Inches(2.5), "UN SDG 4 · Education", BRAND, WHITE, size=11)
chip(s, Inches(8.15), Inches(5.55), Inches(3.0), "UN SDG 12 · Responsible Use", AMBER, INK, size=11)

text(s, Inches(5.5), Inches(6.45), Inches(7), Inches(0.4),
     [P(R("Business & Product Overview", 13, WHITE, bold=True),
        R("    Team Zootopia — 6 members", 13, RGBColor(0xA9,0xB4,0xAE)))])

# ============================================================================
# SLIDE 2 — THE PROBLEM
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "The problem")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("Good books die on shelves while readers go without", 32, INK, bold=True))])

# three problem cards
probs = [
    ("📚", "Books sit idle", "Households keep shelves of read-once books. They gather dust instead of finding the next reader."),
    ("💸", "Marketplaces miss the point", "General resale platforms push profit and shipping fees — wrong incentives for casual book sharing."),
    ("🤝", "No trust between strangers", "Giving a book to someone you don't know feels risky. There is no signal of who is reliable."),
]
cx = Inches(0.7); cw = Inches(3.95); gap = Inches(0.13)
for i, (g, h, b) in enumerate(probs):
    x = cx + i * (cw + gap)
    round_rect(s, x, Inches(2.15), cw, Inches(2.95), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    rect(s, x, Inches(2.15), Inches(0.12), Inches(2.95), [CORAL, AMBER, BLUE][i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.35), Inches(2.45), Inches(1), Inches(0.6), [P(R(g, 30, INK))])
    text(s, x + Inches(0.35), Inches(3.25), cw - Inches(0.7), Inches(0.5),
         [P(R(h, 17, INK, bold=True))])
    text(s, x + Inches(0.35), Inches(3.75), cw - Inches(0.7), Inches(1.2),
         [P(R(b, 12.5, MUTED))], line_spacing=1.18)

# bottom insight band
round_rect(s, Inches(0.7), Inches(5.45), Inches(11.93), Inches(1.05), BRAND_LT, radius=0.08)
text(s, Inches(1.1), Inches(5.45), Inches(11.2), Inches(1.05),
     [P(R("The gap:  ", 16, BRAND_DK, bold=True),
        R("circulation needs coordination + trust — not a payment processor. ", 16, INK),
        R("That gap is exactly what BookBridge fills.", 16, BRAND_DK, bold=True))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
page_num(s, 2, section="The problem")
footer_brand(s)

# ============================================================================
# SLIDE 3 — VISION & SCOPE
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "Vision & scope")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("A non-profit community, not a marketplace", 32, INK, bold=True))])

# big quote
round_rect(s, Inches(0.7), Inches(1.95), Inches(11.93), Inches(1.5), INK, radius=0.06)
text(s, Inches(1.2), Inches(1.95), Inches(11.0), Inches(1.5),
     [P(R("“", 40, AMBER, bold=True),
        R("Keep books in circulation and out of landfills — through a trust-driven ", 18, WHITE),
        R("community of readers who gift, exchange, or sell at symbolic prices.”", 18, WHITE))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# we are / we are not
col_w = Inches(5.85)
round_rect(s, Inches(0.7), Inches(3.7), col_w, Inches(2.9), CARD, line=LINE, line_w=Pt(1), radius=0.05)
text(s, Inches(1.05), Inches(3.95), col_w-Inches(0.7), Inches(0.4),
     [P(R("✓  We are", 16, BRAND, bold=True))])
are = ["Web-based, mobile-responsive",
       "Non-profit, grant-funded",
       "A coordination tool for user-to-user sharing",
       "A trust-driven community with reputation"]
for i, t in enumerate(are):
    text(s, Inches(1.05), Inches(4.5) + i*Inches(0.48), col_w-Inches(0.7), Inches(0.4),
         [P(R("•  ", 14, BRAND, bold=True), R(t, 14, INK))])

x2 = Inches(0.7) + col_w + Inches(0.23)
round_rect(s, x2, Inches(3.7), col_w, Inches(2.9), CARD, line=LINE, line_w=Pt(1), radius=0.05)
text(s, x2+Inches(0.35), Inches(3.95), col_w-Inches(0.7), Inches(0.4),
     [P(R("✕  We are not", 16, CORAL, bold=True))])
arenot = ["A commercial marketplace",
          "A payment processor",
          "A shipping / logistics provider",
          "A replacement for legal enforcement"]
for i, t in enumerate(arenot):
    text(s, x2+Inches(0.35), Inches(4.5) + i*Inches(0.48), col_w-Inches(0.7), Inches(0.4),
         [P(R("•  ", 14, CORAL, bold=True), R(t, 14, INK_SOFT))])

# price cap callout
chip(s, x2+Inches(0.35), Inches(6.05), Inches(5.0),
     "Sale price capped at 50,000 VND — by design", AMBER_LT, INK, size=11.5)
page_num(s, 3, section="Vision & scope")
footer_brand(s)

# ============================================================================
# SLIDE 4 — WHO IT'S FOR (user classes)
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "Who it's for")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("Four roles, escalating trust & responsibility", 32, INK, bold=True))])

roles = [
    ("Guest", "Browse & search freely. Cannot transact or message.", "occasional", MUTED, RGBColor(0xEF,0xEC,0xE4)),
    ("Registered User", "List, request, gift, exchange, sell, rate & chat.", "weekly–daily", BRAND, BRAND_LT),
    ("Community Moderator", "Curate a sub-community; review in-scope reports.", "frequent", AMBER, AMBER_LT),
    ("System Administrator", "Full dashboard, configuration, analytics, exports.", "daily", PURPLE, RGBColor(0xEC,0xE6,0xF5)),
]
y = Inches(2.1); rh = Inches(1.18); rgap = Inches(0.12)
for i, (name, desc, freq, accent, light) in enumerate(roles):
    ry = y + i*(rh+rgap)
    round_rect(s, Inches(0.7), ry, Inches(11.93), rh, CARD, line=LINE, line_w=Pt(1), radius=0.08)
    rect(s, Inches(0.7), ry, Inches(0.14), rh, accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # level number
    rect(s, Inches(1.05), ry+Inches(0.29), Inches(0.6), Inches(0.6), light, shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.05), ry+Inches(0.29), Inches(0.6), Inches(0.6),
         [P(R(str(i+1), 22, accent, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.95), ry+Inches(0.18), Inches(4.2), Inches(0.5),
         [P(R(name, 18, INK, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.95), ry+Inches(0.62), Inches(6.5), Inches(0.45),
         [P(R(desc, 13, MUTED))], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, Inches(10.3), ry+Inches(0.4), Inches(2.0), freq, light, accent, size=11)
page_num(s, 4, section="Who it's for")
footer_brand(s)

# ============================================================================
# SLIDE 5 — THE JOURNEY (storytelling spine)
# ============================================================================
s = slide()
bg(s, INK)
kicker(s, Inches(0.7), Inches(0.6), "The journey", color=AMBER)
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("One story, end to end", 32, WHITE, bold=True))])
text(s, Inches(0.7), Inches(1.6), Inches(11.5), Inches(0.5),
     [P(R("Meet Alice. Every feature in this deck appears at a step of her journey.", 15, RGBColor(0xC9,0xD0,0xCC)))])

steps = [
    ("1", "Join", "Register, verify email,\nbuild a profile", BRAND),
    ("2", "List", "Post a book with ISBN\nauto-fill & photos", AMBER),
    ("3", "Discover", "Others search, follow\nAlice & see her feed", BLUE),
    ("4", "Transact", "Request → accept → ship\n→ complete, with chat", CORAL),
    ("5", "Trust", "Both rate each other;\nreputation grows", PURPLE),
    ("6", "Belong", "Communities, notifications\nkeep everyone in the loop", BRAND),
]
n = len(steps)
sx = Inches(0.7); total_w = Inches(11.93)
cw = Inches(1.82); gap = (total_w - cw*n) / (n-1)
cy = Inches(2.6)
# connector line
rect(s, sx + cw/2, cy + Inches(0.45), total_w - cw, Pt(2.5), RGBColor(0x4A,0x4A,0x66))
for i, (num, title, body, col) in enumerate(steps):
    x = sx + i*(cw+gap)
    icon_circle(s, x + cw/2 - Inches(0.35), cy + Inches(0.1), Inches(0.7), num, col, WHITE, size=22)
    text(s, x, cy + Inches(0.95), cw, Inches(0.4),
         [P(R(title, 16, WHITE, bold=True))], align=PP_ALIGN.CENTER)
    text(s, x - Inches(0.1), cy + Inches(1.4), cw + Inches(0.2), Inches(1.3),
         [P(R(body, 11.5, RGBColor(0xB8,0xC0,0xBB)))], align=PP_ALIGN.CENTER, line_spacing=1.12)

round_rect(s, Inches(0.7), Inches(5.75), Inches(11.93), Inches(0.95), BRAND_DK, radius=0.1)
text(s, Inches(1.1), Inches(5.75), Inches(11.2), Inches(0.95),
     [P(R("Behind these 6 steps sit ", 15, RGBColor(0xD9,0xE2,0xDD)),
        R("10 functional modules", 15, AMBER, bold=True),
        R(" and ", 15, RGBColor(0xD9,0xE2,0xDD)),
        R("13 data entities", 15, AMBER, bold=True),
        R(" — the rest of the deck details each.", 15, RGBColor(0xD9,0xE2,0xDD)))],
     anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 5, section="The journey")

# ============================================================================
# SLIDE 6 — FEATURE MAP (10 modules overview)
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "Feature map")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("Ten modules, mapped to the SRS", 32, INK, bold=True))])

mods = [
    ("4.1", "Identity & Auth", "Register, verify, login, profile", BRAND),
    ("4.2", "Book Catalog", "Listings, photos, ISBN lookup", AMBER),
    ("4.3", "Search & Discovery", "Full-text, filters, feed", BLUE),
    ("4.4", "Transactions", "8-state workflow + ratings", CORAL),
    ("4.5", "Reputation & Trust", "Scoring + anti-gaming", PURPLE),
    ("4.6", "Social Connectivity", "Follow, feed, notifications", BLUE),
    ("4.7", "Communities", "University / location / genre", BRAND),
    ("4.8", "In-app Messaging", "1-to-1 chat, real-time", CORAL),
    ("4.9", "Moderation & Reports", "Queue + audited actions", PURPLE),
    ("4.10", "Admin Dashboard", "Stats + grant reporting", AMBER),
]
cols = 5; cw = Inches(2.3); ch = Inches(1.95); gx = Inches(0.12); gy = Inches(0.18)
x0 = Inches(0.7); y0 = Inches(2.05)
for i, (sec, name, desc, col) in enumerate(mods):
    r = i // cols; c = i % cols
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    round_rect(s, x, y, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=0.06)
    rect(s, x, y, cw, Inches(0.1), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    chip(s, x+Inches(0.22), y+Inches(0.28), Inches(0.7), sec, col, WHITE, size=10, h=Inches(0.3))
    text(s, x+Inches(0.22), y+Inches(0.72), cw-Inches(0.44), Inches(0.7),
         [P(R(name, 14.5, INK, bold=True))], line_spacing=1.0)
    text(s, x+Inches(0.22), y+Inches(1.38), cw-Inches(0.44), Inches(0.5),
         [P(R(desc, 11, MUTED))], line_spacing=1.05)
page_num(s, 6, section="Feature map")
footer_brand(s)


# ----------------------------------------------------------------------------
# Reusable feature-detail slide builder
# ----------------------------------------------------------------------------
def feature_slide(num, sec_tag, title, lead, accent, light, features, prove, total=18):
    s = slide()
    bg(s)
    # accent header band
    rect(s, 0, 0, SW, Inches(1.55), accent)
    rect(s, 0, Inches(1.55), SW, Inches(0.06), light)
    chip(s, Inches(0.7), Inches(0.42), Inches(1.5), "SRS " + sec_tag, WHITE, accent, size=12, h=Inches(0.36))
    text(s, Inches(2.4), Inches(0.32), Inches(10.4), Inches(0.7),
         [P(R(title, 28, WHITE, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(2.42), Inches(1.0), Inches(10.4), Inches(0.4),
         [P(R(lead, 13.5, RGBColor(0xF2,0xF2,0xF2)))], anchor=MSO_ANCHOR.MIDDLE)

    # feature rows (left, larger column)
    fx = Inches(0.7); fw = Inches(8.0); fy = Inches(2.0)
    for i, (h, b) in enumerate(features):
        ry = fy + i*Inches(0.92)
        icon_circle(s, fx, ry+Inches(0.06), Inches(0.4), "✓", light, accent, size=14)
        text(s, fx+Inches(0.6), ry-Inches(0.02), fw-Inches(0.6), Inches(0.4),
             [P(R(h, 15, INK, bold=True))])
        text(s, fx+Inches(0.6), ry+Inches(0.34), fw-Inches(0.6), Inches(0.55),
             [P(R(b, 12, MUTED))], line_spacing=1.1)

    # right rail — "what proves it"
    rx = Inches(9.0); rw = Inches(3.63)
    round_rect(s, rx, fy, rw, Inches(4.55), light, radius=0.05)
    text(s, rx+Inches(0.3), fy+Inches(0.25), rw-Inches(0.6), Inches(0.4),
         [P(R(prove[0], 13.5, accent, bold=True))])
    for i, item in enumerate(prove[1]):
        text(s, rx+Inches(0.3), fy+Inches(0.8)+i*Inches(0.62), rw-Inches(0.6), Inches(0.6),
             [P(R("›  ", 13, accent, bold=True), R(item, 11.5, INK))], line_spacing=1.05)
    page_num(s, num, total=total, section=title)
    footer_brand(s)
    return s

# ============================================================================
# SLIDE 7 — Identity & Auth (4.1)
# ============================================================================
feature_slide(
    7, "§4.1", "Identity & Profile",
    "Every other module depends on knowing who the user is — securely.",
    BRAND, BRAND_LT,
    [("Email registration & verification",
      "Sign up with email + strong password; verify within a 72-hour token window before the account activates."),
     ("Secure sessions",
      "Argon2id password hashing, AES-encrypted httpOnly cookies, 30-minute idle timeout, multi-device support."),
     ("Rate-limited, constant-time login",
      "5 attempts / 15 min per IP; dummy-hash verification on unknown users defeats timing attacks."),
     ("Rich, editable profiles",
      "Display name, avatar, bio, preferred genres and district-level location — never exact coordinates."),
     ("Password reset & account deletion",
      "1-hour reset tokens; deletion anonymises transaction history (right to be forgotten).")],
    ("Trust signals", ["Argon2id (OWASP-aligned)",
                        "iron-session encrypted cookies",
                        "Roles: Guest · User · Mod · Admin",
                        "Suspended users can't log in",
                        "Tier shown across the app"]),
)

# ============================================================================
# SLIDE 8 — Book Catalog (4.2)
# ============================================================================
feature_slide(
    8, "§4.2", "Book Catalog",
    "Listing a book takes under a minute — and stays non-commercial by design.",
    AMBER, AMBER_LT,
    [("Nine-field rich listings",
      "Title, author, ISBN, genre, condition, 20–2000-char description, 1–5 photos (≤5 MB), type and price."),
     ("ISBN auto-fill",
      "Paste an ISBN → Open Library populates metadata instantly; manual entry remains for books with no ISBN."),
     ("Price cap enforced server-side",
      "SELL listings above the 50,000 VND cap are rejected — keeping the platform symbolic, not commercial."),
     ("Photo pipeline",
      "Parallel upload, resized to WebP ≤1024px to save bandwidth; stored locally in dev, S3/R2 in production."),
     ("Safe edits & soft delete",
      "Editing is blocked once a transaction is accepted; deletes are soft (status REMOVED) to preserve history.")],
    ("Three listing types", ["🎁  Gift — free",
                             "🔄  Exchange — book-for-book",
                             "🏷️  Sell — capped price",
                             "Owner tier shown on every card",
                             "Optional community scoping"]),
)

# ============================================================================
# SLIDE 9 — Search & Discovery (4.3 + 4.6 social)
# ============================================================================
feature_slide(
    9, "§4.3 · 4.6", "Search & Discovery",
    "Readers find the right book fast — and stay connected to people they trust.",
    BLUE, RGBColor(0xE4,0xEC,0xF3),
    [("Full-text relevance search",
      "Ranks across title, author, ISBN & description with PostgreSQL tsvector + GIN index — sub-second at scale."),
     ("Combinable filters",
      "Genre, condition, type, max price, distance radius and sub-community — intersected with AND logic."),
     ("Personalised feed",
      "Listings from followed users and joined communities, newest first — materialised via fan-out on write."),
     ("Follow / unfollow graph",
      "One-directional follows with maintained counts; new listings from followed users appear in real time."),
     ("Related listings",
      "Detail pages surface same-author / same-genre / same-community books to deepen discovery.")],
    ("Performance target", ["< 2 s search at 95th pct",
                            "GIN index: 200 ms → 5 ms",
                            "Cursor pagination (not offset)",
                            "District centroids for distance",
                            "Real-time feed via SSE"]),
)

# ============================================================================
# SLIDE 10 — Transactions (4.4)  — the centrepiece, gets a custom layout
# ============================================================================
s = slide()
bg(s)
rect(s, 0, 0, SW, Inches(1.55), CORAL)
rect(s, 0, Inches(1.55), SW, Inches(0.06), RGBColor(0xF3,0xD9,0xD4))
chip(s, Inches(0.7), Inches(0.42), Inches(1.5), "SRS §4.4", WHITE, CORAL, size=12, h=Inches(0.36))
text(s, Inches(2.4), Inches(0.32), Inches(10.4), Inches(0.7),
     [P(R("Transaction Workflow", 28, WHITE, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(2.42), Inches(1.0), Inches(10.4), Inches(0.4),
     [P(R("The heart of BookBridge — a rigorously tested 8-state machine.", 13.5, WHITE))], anchor=MSO_ANCHOR.MIDDLE)

# state machine ribbon
states = ["PENDING", "ACCEPTED", "IN_DELIVERY", "COMPLETED", "RATED"]
sx = Inches(0.7); sy = Inches(2.05); bw = Inches(1.95); bh = Inches(0.7); arrow_w = Inches(0.42)
for i, st in enumerate(states):
    x = sx + i*(bw+arrow_w)
    col = BRAND if st in ("COMPLETED","RATED") else INK
    round_rect(s, x, sy, bw, bh, col, radius=0.18)
    text(s, x, sy, bw, bh, [P(R(st, 11.5, WHITE, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(states)-1:
        text(s, x+bw, sy, arrow_w, bh, [P(R("→", 20, CORAL, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# side states
text(s, Inches(0.7), sy+Inches(0.85), Inches(11.9), Inches(0.4),
     [P(R("Branches:  ", 12, CORAL, bold=True),
        R("DECLINED · CANCELLED · WAITLISTED · DISPUTED", 12, INK_SOFT, bold=True),
        R("   — every illegal transition is rejected and unit-tested.", 12, MUTED))])

# feature columns
feats = [
    ("Three transaction types", "Gift, Exchange (book-for-book) and Sell (capped) — type is locked at request time."),
    ("Multi-requester waitlist", "Many can request one book; owner accepts one, the rest auto-move to Waitlisted, promoted if it falls through."),
    ("Delivery & confirmation", "In-person (map hint) or postal (tracking #). 14-day reminder, auto-complete at 21 days if unconfirmed."),
    ("Two-way ratings", "After completion, both parties rate 1–5 ★ + comment — feeding the reputation engine."),
    ("Full audit trail", "Every transition writes a TransactionEvent row — event-sourced, replayable history."),
    ("Dispute path", "Either party can dispute → routed to a moderator who resolves to Completed or Cancelled."),
]
fx = Inches(0.7); fw = Inches(5.85); fy = Inches(3.55)
for i, (h, b) in enumerate(feats):
    c = i % 2; r = i // 2
    x = fx + c*(fw+Inches(0.23)); y = fy + r*Inches(1.05)
    text(s, x, y, fw, Inches(0.35), [P(R("●  ", 12, CORAL, bold=True), R(h, 14, INK, bold=True))])
    text(s, x+Inches(0.25), y+Inches(0.32), fw-Inches(0.25), Inches(0.65),
         [P(R(b, 11.5, MUTED))], line_spacing=1.1)
page_num(s, 10, section="Transaction Workflow")
footer_brand(s)

# ============================================================================
# SLIDE 11 — Messaging (4.8)
# ============================================================================
feature_slide(
    11, "§4.8", "In-app Messaging",
    "Coordinate the hand-off without leaving the platform or sharing personal contacts.",
    CORAL, RGBColor(0xF3,0xD9,0xD4),
    [("1-to-1 conversations",
      "A conversation is the unordered pair of two users, optionally tied to a specific transaction thread."),
     ("Real-time delivery",
      "Server-Sent Events push new messages instantly; the client posts normally and receives live updates."),
     ("Transaction-scoped chat",
      "Embedded directly in the transaction detail page, so context (which book, which step) is never lost."),
     ("Read receipts & inbox",
      "Per-message read state and a unified inbox sorted by latest activity."),
     ("Reportable content",
      "Any message can be reported to moderation — abuse never hides inside private threads.")],
    ("Why SSE, not WebSocket?", ["HTTP-based, auto-reconnect",
                                 "One-way push fits messaging",
                                 "Simpler to operate & deploy",
                                 "Long-poll fallback if blocked",
                                 "One connection per user-tab"]),
)

# ============================================================================
# SLIDE 12 — Reputation & Trust (4.5)  — the differentiator
# ============================================================================
s = slide()
bg(s, INK)
kicker(s, Inches(0.7), Inches(0.55), "The differentiator", color=AMBER)
text(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.9),
     [P(R("Reputation & Trust — what sets us apart", 30, WHITE, bold=True))])
text(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.5),
     [P(R("Facebook Marketplace optimises for sales. BookBridge optimises for ", 14, RGBColor(0xC9,0xD0,0xCC)),
        R("trustworthy circulation.", 14, AMBER, bold=True))])

# tier ladder
tiers = [("0–19", "New Member", MUTED),
         ("20–49", "Active Sharer", BLUE),
         ("50–79", "Trusted Contributor", BRAND),
         ("80–100", "Community Champion", AMBER)]
tx = Inches(0.7); tw = Inches(2.9); th = Inches(1.35); tgap = Inches(0.12)
for i, (rng, name, col) in enumerate(tiers):
    x = tx + i*(tw+tgap)
    round_rect(s, x, Inches(2.3), tw, th, RGBColor(0x24,0x24,0x40), radius=0.08)
    rect(s, x, Inches(2.3), tw, Inches(0.12), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, Inches(2.55), tw, Inches(0.4), [P(R(rng, 20, col, bold=True))], align=PP_ALIGN.CENTER)
    text(s, x, Inches(3.05), tw, Inches(0.5), [P(R(name, 13.5, WHITE, bold=True))], align=PP_ALIGN.CENTER, line_spacing=1.0)

# how it works + anti-gaming
round_rect(s, Inches(0.7), Inches(3.95), Inches(5.85), Inches(2.55), RGBColor(0x24,0x24,0x40), radius=0.05)
text(s, Inches(1.05), Inches(4.2), Inches(5.2), Inches(0.4), [P(R("How the score moves", 15, AMBER, bold=True))])
moves = ["Completed transaction  →  +10 each party",
         "Rating received  →  ± by star value",
         "Cancellation  →  −3 to cancelling party",
         "Report upheld  →  −15 to target",
         "Inactivity  →  −1 every 30 days (decay)"]
for i, m in enumerate(moves):
    text(s, Inches(1.05), Inches(4.65)+i*Inches(0.37), Inches(5.2), Inches(0.35),
         [P(R(m, 12, RGBColor(0xDD,0xE4,0xE0)))])

round_rect(s, Inches(6.78), Inches(3.95), Inches(5.85), Inches(2.55), BRAND_DK, radius=0.05)
text(s, Inches(7.13), Inches(4.2), Inches(5.2), Inches(0.4), [P(R("Anti-gaming detection", 15, AMBER, bold=True))])
text(s, Inches(7.13), Inches(4.65), Inches(5.2), Inches(1.8),
     [P(R("Reciprocal-only pairs", 12.5, WHITE, bold=True),
        R(" (two accounts only trading with each other) and ", 12.5, RGBColor(0xD9,0xE2,0xDD)),
        R("zero-unique-counterparty", 12.5, WHITE, bold=True),
        R(" accounts are flagged automatically for moderator review — never auto-banned, to avoid punishing real friends.", 12.5, RGBColor(0xD9,0xE2,0xDD)))],
     line_spacing=1.22)
text(s, Inches(7.13), Inches(6.05), Inches(5.2), Inches(0.4),
     [P(R("Source of truth: ReputationEvent log; score is denormalised for fast reads.", 10.5, RGBColor(0xA9,0xC2,0xB6), italic=True))], line_spacing=1.1)
page_num(s, 12, section="Reputation & Trust")

# ============================================================================
# SLIDE 13 — Moderation (4.9)
# ============================================================================
feature_slide(
    13, "§4.9", "Moderation & Reporting",
    "A safe community needs human judgement, backed by a full audit trail.",
    PURPLE, RGBColor(0xEC,0xE6,0xF5),
    [("Report anything",
      "Users can report a user, listing, transaction or message with a reason and details."),
     ("Scoped moderator queue",
      "Community moderators see only in-scope reports; admins see everything. Pending items surface first."),
     ("Graduated actions",
      "Warn, remove listing, suspend user, restore, and resolve/reject disputes — proportionate responses."),
     ("Suspend vs delete",
      "Suspension preserves data for appeal; deletion anonymises — both are deliberate, reversible-where-fair choices."),
     ("Compliant audit log",
      "Every ModerationAction is recorded — who, on whom, why — for appeals and grant compliance.")],
    ("Cross-module effects", ["Upheld report → −15 reputation",
                              "Remove → listing status REMOVED",
                              "Suspend → login blocked",
                              "Target user is notified",
                              "Dispute resolution closes txn"]),
)

# ============================================================================
# SLIDE 14 — Communities, Notifications & Admin (4.7 + 4.6 + 4.10)
# ============================================================================
s = slide()
bg(s)
rect(s, 0, 0, SW, Inches(1.55), BRAND)
rect(s, 0, Inches(1.55), SW, Inches(0.06), BRAND_LT)
chip(s, Inches(0.7), Inches(0.42), Inches(2.2), "SRS §4.7 · 4.6 · 4.10", WHITE, BRAND, size=11, h=Inches(0.36))
text(s, Inches(3.1), Inches(0.32), Inches(9.6), Inches(0.7),
     [P(R("Community, Notifications & Admin", 27, WHITE, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(3.12), Inches(1.0), Inches(9.6), Inches(0.4),
     [P(R("The connective tissue that keeps everyone informed and the platform accountable.", 13, WHITE))], anchor=MSO_ANCHOR.MIDDLE)

cols3 = [
    ("Communities", BRAND, [
        "University / location / genre scopes",
        "Join / leave; max 20 per user",
        "Listings can be scoped to a community",
        "Posts, likes & comments with points",
    ]),
    ("Notifications", AMBER, [
        "9 event kinds across all modules",
        "Real-time via SSE + long-poll fallback",
        "Email: immediate / daily digest / off",
        "Central dispatcher routes every event",
    ]),
    ("Admin & Ops", PURPLE, [
        "Live stats: users, txns, books circulated",
        "Grant-sponsor CSV export (injection-safe)",
        "GitHub Actions CI on every PR",
        "Vercel + managed Postgres deploy",
    ]),
]
cx = Inches(0.7); cw = Inches(3.9); gap = Inches(0.13); cy = Inches(2.05)
for i, (title, col, items) in enumerate(cols3):
    x = cx + i*(cw+gap)
    round_rect(s, x, cy, cw, Inches(4.4), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    rect(s, x, cy, cw, Inches(0.7), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, cy+Inches(0.5), cw, Inches(0.2), col)  # square off bottom of header
    text(s, x, cy, cw, Inches(0.7), [P(R(title, 17, WHITE, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for j, it in enumerate(items):
        iy = cy + Inches(0.95) + j*Inches(0.82)
        icon_circle(s, x+Inches(0.3), iy, Inches(0.36), "›", BRAND_LT if col==BRAND else (AMBER_LT if col==AMBER else RGBColor(0xEC,0xE6,0xF5)), col, size=14)
        text(s, x+Inches(0.8), iy-Inches(0.05), cw-Inches(1.0), Inches(0.75),
             [P(R(it, 12, INK))], line_spacing=1.08, anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 14, section="Community · Notifications · Admin")
footer_brand(s)

# ============================================================================
# SLIDE 15 — ARCHITECTURE
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "Under the hood")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("One language, front to back", 32, INK, bold=True))])
text(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.5),
     [P(R("TypeScript everywhere keeps a 6-person team — most new to web — productive and aligned.", 14, MUTED))])

# layered diagram
layers = [
    ("Browser", "React + Tailwind · mobile-responsive · WCAG 2.1 AA", BLUE),
    ("Next.js 14 (App Router)", "Pages + API routes — validate input, call into domain modules", BRAND),
    ("Domain modules  (src/server)", "auth · listings · search · transactions · reputation · moderation · communities · notifications · admin", AMBER),
    ("Prisma 5 ORM", "The only layer that talks to the database — type-safe queries", PURPLE),
    ("PostgreSQL 16", "13 entities · ~30 indexes · full-text search", INK),
]
ly = Inches(2.25); lh = Inches(0.82); lgap = Inches(0.12); lx = Inches(1.4); lw = Inches(10.5)
for i, (name, desc, col) in enumerate(layers):
    y = ly + i*(lh+lgap)
    round_rect(s, lx, y, lw, lh, CARD, line=LINE, line_w=Pt(1.2), radius=0.08)
    rect(s, lx, y, Inches(0.16), lh, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, lx+Inches(0.45), y, Inches(3.6), lh, [P(R(name, 15, INK, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, lx+Inches(4.2), y, lw-Inches(4.5), lh, [P(R(desc, 11.5, MUTED))], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    if i < len(layers)-1:
        text(s, lx+lw/2-Inches(0.2), y+lh-Inches(0.02), Inches(0.4), lgap+Inches(0.12),
             [P(R("▼", 11, LINE))], align=PP_ALIGN.CENTER)

# side note: events
text(s, Inches(1.4), Inches(6.95), Inches(10.5), Inches(0.4),
     [P(R("Modules stay decoupled via in-process domain events: ", 11.5, INK_SOFT),
        R("txn.completed → reputation +10, listing→COMPLETED, notify both parties.", 11.5, BRAND, bold=True))])
page_num(s, 15, section="Architecture")
footer_brand(s)

# ============================================================================
# SLIDE 16 — DATA MODEL
# ============================================================================
s = slide()
bg(s, INK)
kicker(s, Inches(0.7), Inches(0.6), "Data model", color=AMBER)
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("13 entities, clear ownership", 32, WHITE, bold=True))])

groups = [
    ("Identity", BRAND, ["User", "Session", "+ tokens"]),
    ("Catalog", AMBER, ["Listing", "ListingPhoto"]),
    ("Discovery", BLUE, ["Follow", "FeedItem"]),
    ("Transactions", CORAL, ["Transaction", "TxnEvent", "Rating", "Conversation", "Message"]),
    ("Trust", PURPLE, ["ReputationEvent", "Report", "ModerationAction"]),
    ("Community", BRAND, ["Community", "Membership", "Notification"]),
]
gx = Inches(0.7); gw = Inches(1.92); ggap = Inches(0.115); gy = Inches(2.15)
for i, (title, col, ents) in enumerate(groups):
    x = gx + i*(gw+ggap)
    round_rect(s, x, gy, gw, Inches(3.6), RGBColor(0x24,0x24,0x40), radius=0.05)
    rect(s, x, gy, gw, Inches(0.6), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, gy+Inches(0.42), gw, Inches(0.18), col)
    text(s, x, gy, gw, Inches(0.6), [P(R(title, 13, WHITE, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for j, e in enumerate(ents):
        ey = gy + Inches(0.78) + j*Inches(0.5)
        round_rect(s, x+Inches(0.18), ey, gw-Inches(0.36), Inches(0.42), RGBColor(0x32,0x32,0x52), radius=0.2)
        text(s, x+Inches(0.18), ey, gw-Inches(0.36), Inches(0.42),
             [P(R(e, 11, RGBColor(0xE4,0xE9,0xE6), bold=(not e.startswith("+"))))],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# bottom facts
facts = [("onDelete rules", "Cascade for owned data, SetNull for soft links"),
         ("Denormalised score", "User.reputationScore for fast reads"),
         ("~30 indexes", "Search, dashboards & inbox stay fast")]
fy = Inches(6.05); fw = Inches(3.85)
for i, (h, b) in enumerate(facts):
    x = Inches(0.7) + i*(fw+Inches(0.19))
    text(s, x, fy, fw, Inches(0.4), [P(R(h, 13, AMBER, bold=True))])
    text(s, x, fy+Inches(0.38), fw, Inches(0.6), [P(R(b, 11.5, RGBColor(0xC9,0xD0,0xCC)))], line_spacing=1.1)
page_num(s, 16, section="Data model")

# ============================================================================
# SLIDE 17 — QUALITY & NON-FUNCTIONAL
# ============================================================================
s = slide()
bg(s)
kicker(s, Inches(0.7), Inches(0.6), "Built to standard")
text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(1.0),
     [P(R("Quality is a feature", 32, INK, bold=True))])

metrics = [
    ("< 3 s", "Page load over 4G", BLUE),
    ("< 2 s", "Search (95th pct)", BRAND),
    ("1,000+", "Concurrent users", AMBER),
    ("99.5%", "Monthly uptime", PURPLE),
    ("≥ 70%", "Test coverage", CORAL),
    ("AA", "WCAG 2.1 accessibility", BRAND),
    ("VI / EN", "Internationalisation", BLUE),
    ("District", "Location precision only", AMBER),
]
cols = 4; cw = Inches(2.86); ch = Inches(1.75); gx = Inches(0.16); gy = Inches(0.2)
x0 = Inches(0.7); y0 = Inches(2.05)
for i, (big, lab, col) in enumerate(metrics):
    r = i//cols; c = i%cols
    x = x0 + c*(cw+gx); y = y0 + r*(ch+gy)
    round_rect(s, x, y, cw, ch, CARD, line=LINE, line_w=Pt(1), radius=0.06)
    rect(s, x, y, Inches(0.12), ch, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+Inches(0.35), y+Inches(0.3), cw-Inches(0.5), Inches(0.8),
         [P(R(big, 34, col, bold=True))])
    text(s, x+Inches(0.37), y+Inches(1.12), cw-Inches(0.55), Inches(0.5),
         [P(R(lab, 13, INK_SOFT))], line_spacing=1.0)

round_rect(s, Inches(0.7), Inches(6.05), Inches(11.93), Inches(0.95), BRAND_LT, radius=0.1)
text(s, Inches(1.1), Inches(6.05), Inches(11.2), Inches(0.95),
     [P(R("Verified by ", 14, INK),
        R("Vitest unit + integration tests", 14, BRAND_DK, bold=True),
        R(" (state machine, scoring & anti-gaming run DB-free), Playwright E2E for the 3 core flows, and a green CI gate on every PR.", 14, INK))],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
page_num(s, 17, section="Quality")
footer_brand(s)

# ============================================================================
# SLIDE 18 — CLOSE / ROADMAP
# ============================================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(0.3), SH, BRAND)
text(s, Inches(0.9), Inches(0.85), Inches(11.5), Inches(1.2),
     [P(R("BookBridge keeps books moving —", 34, WHITE, bold=True)),
      P(R("and builds the trust that makes sharing work.", 34, AMBER, bold=True))],
     line_spacing=1.05)

# recap pillars
pillars = [
    ("Circulation", "Gift, exchange or sell at symbolic prices — books reach the next reader."),
    ("Trust", "Reputation, ratings & anti-gaming make strangers safe to deal with."),
    ("Community", "Sub-communities, follows & notifications turn users into a network."),
]
px = Inches(0.9); pw = Inches(3.78); py = Inches(2.85)
for i, (h, b) in enumerate(pillars):
    x = px + i*(pw+Inches(0.19))
    round_rect(s, x, py, pw, Inches(1.95), RGBColor(0x24,0x24,0x40), radius=0.06)
    rect(s, x, py, pw, Inches(0.1), [BRAND, AMBER, BLUE][i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+Inches(0.3), py+Inches(0.3), pw-Inches(0.6), Inches(0.5),
         [P(R(h, 19, [BRAND_LT, AMBER, RGBColor(0x9D,0xC0,0xE8)][i], bold=True))])
    text(s, x+Inches(0.3), py+Inches(0.85), pw-Inches(0.6), Inches(1.0),
         [P(R(b, 12.5, RGBColor(0xD9,0xE2,0xDD)))], line_spacing=1.18)

# roadmap line
text(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4),
     [P(R("WHAT'S NEXT", 12.5, AMBER, bold=True))])
nexts = ["Background fan-out queue for hot users",
         "Trading-ring detection (≥3 accounts)",
         "E2E-encrypted messaging",
         "Monitoring: Prometheus + Grafana"]
for i, t in enumerate(nexts):
    c = i % 2; r = i // 2
    text(s, Inches(0.9) + c*Inches(5.9), Inches(5.55) + r*Inches(0.45), Inches(5.7), Inches(0.4),
         [P(R("→  ", 13, BRAND, bold=True), R(t, 13.5, WHITE))])

round_rect(s, Inches(0.9), Inches(6.55), Inches(11.53), Inches(0.62), BRAND_DK, radius=0.2)
text(s, Inches(0.9), Inches(6.55), Inches(11.53), Inches(0.62),
     [P(R("Team Zootopia  ·  Software Engineering 2025.2  ·  Next.js + Prisma + PostgreSQL", 13, WHITE, bold=True))],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

out = "BookBridge_Storytelling.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")

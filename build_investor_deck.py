# -*- coding: utf-8 -*-
"""
BookBridge — Investor pitch deck (12 slides, ~7 minutes).

Storytelling edition. Each slide carries ONE short, logical paragraph in a
human voice — built to be spoken from, in front of a board. The arc:

  1  Title / hook
  2  The shelf            (open on a relatable moment)
  3  Why we don't share   (the tension)
  4  The real gap         (the insight — why now)
  5  Meet BookBridge      (the turn / solution)
  6  Listing              (a book's journey begins)
  7  Search & discovery   (a reader finds it)
  8  Chat & the hand-off  (two strangers, safely)
  9  Community            (they belong now)
  10 Trust — the moat     (why it compounds)
  11 Model & scale        (how it lasts)
  12 Close & the ask      (the vision)

Run: python build_investor_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Design system
# ----------------------------------------------------------------------------
SW, SH = Inches(13.333), Inches(7.5)  # 16:9

INK      = RGBColor(0x1A, 0x1A, 0x2E)
INK_SOFT = RGBColor(0x3D, 0x3D, 0x56)
PAPER    = RGBColor(0xFB, 0xF8, 0xF1)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xE6, 0xDF, 0xD2)
MUTED    = RGBColor(0x76, 0x70, 0x66)

BRAND    = RGBColor(0x2E, 0x6F, 0x5B)   # forest green
BRAND_DK = RGBColor(0x1E, 0x4D, 0x3E)
BRAND_LT = RGBColor(0xE3, 0xEF, 0xE9)
AMBER    = RGBColor(0xE0, 0x8E, 0x2C)
AMBER_LT = RGBColor(0xFA, 0xEC, 0xD6)
CORAL    = RGBColor(0xD9, 0x5D, 0x4A)
CORAL_LT = RGBColor(0xF3, 0xD9, 0xD4)
BLUE     = RGBColor(0x3B, 0x6E, 0xA5)
BLUE_LT  = RGBColor(0xE4, 0xEC, 0xF3)
PURPLE   = RGBColor(0x6B, 0x4E, 0x9E)
PURPLE_LT= RGBColor(0xEC, 0xE6, 0xF5)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_DK = RGBColor(0x24, 0x24, 0x40)
CREAM    = RGBColor(0xCF, 0xD6, 0xD2)

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]
TOTAL = 12

# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color=PAPER):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r

def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=None):
    r = s.shapes.add_shape(shape, x, y, w, h)
    if color is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = line_w or Pt(1)
    r.shadow.inherit = False
    return r

def round_rect(s, x, y, w, h, color, line=None, line_w=None, radius=0.08):
    r = rect(s, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=line, line_w=line_w)
    try:
        r.adjustments[0] = radius
    except Exception:
        pass
    return r

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=Pt(2), line_spacing=1.0, wrap=True):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = font
    return tb

def R(txt, size, color=INK, bold=False, italic=False, font=FONT):
    return (txt, size, color, bold, italic, font)

def P(*runs):
    return list(runs)

def page_num(s, n, section="", light=False):
    col = RGBColor(0xA9, 0xB4, 0xAE) if light else MUTED
    text(s, Inches(11.7), Inches(7.04), Inches(1.4), Inches(0.3),
         [P(R(f"{n:02d} / {TOTAL}", 9, col))], align=PP_ALIGN.RIGHT)
    if section:
        text(s, Inches(0.85), Inches(7.04), Inches(8), Inches(0.3),
             [P(R("BookBridge", 9, col, bold=True), R("   ·   " + section, 9, col))])

def footer_brand(s):
    rect(s, 0, Inches(7.34), SW, Inches(0.16), BRAND)

def chip(s, x, y, w, label, fill, txt_color, size=10.5, h=Inches(0.34)):
    round_rect(s, x, y, w, h, fill, radius=0.5)
    text(s, x, y, w, h, [P(R(label, size, txt_color, bold=True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def icon_circle(s, x, y, d, glyph, fill, glyph_color=WHITE, size=18):
    rect(s, x, y, d, d, fill, shape=MSO_SHAPE.OVAL)
    text(s, x, y, d, d, [P(R(glyph, size, glyph_color, bold=True))],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ----------------------------------------------------------------------------
# Story-slide template — the workhorse of this deck
#   A calm editorial page: chapter label, headline, one narrative paragraph,
#   and an optional motif drawn on the right by `motif(s)`.
# ----------------------------------------------------------------------------
def story_slide(num, chapter, headline, paragraph, accent=BRAND,
                motif=None, pull=None, text_w=Inches(8.7), head_size=30):
    s = slide()
    bg(s)
    rect(s, 0, 0, Inches(0.26), SH, accent)               # left spine
    text(s, Inches(0.95), Inches(0.78), Inches(9), Inches(0.4),
         [P(R(chapter.upper(), 12.5, accent, bold=True))])
    text(s, Inches(0.95), Inches(1.28), Inches(11.3), Inches(1.4),
         [P(R(headline, head_size, INK, bold=True))], line_spacing=1.05)
    text(s, Inches(0.95), Inches(2.85), text_w, Inches(3.2),
         paragraph, line_spacing=1.42, space_after=Pt(6))
    if pull:
        text(s, Inches(0.95), Inches(6.25), Inches(11.4), Inches(0.7),
             [P(R(pull, 15, accent, bold=True, italic=True))], line_spacing=1.1)
    if motif:
        motif(s)
    page_num(s, num, section=chapter, light=False)
    footer_brand(s)
    return s


# ============================================================================
# SLIDE 1 — TITLE / HOOK
# ============================================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(5.0), SH, BRAND_DK)
rect(s, Inches(4.7), 0, Inches(0.3), SH, BRAND)

spine_colors = [AMBER, BRAND, CORAL, BLUE, PURPLE, BRAND_LT]
sx = Inches(0.85)
for i, col in enumerate(spine_colors):
    hh = Inches(3.4) - Inches(0.16) * i
    round_rect(s, sx + Inches(0.62) * i, SH - Inches(0.9) - hh, Inches(0.46), hh, col, radius=0.12)

text(s, Inches(0.85), Inches(0.7), Inches(3.6), Inches(0.6),
     [P(R("INVESTOR BRIEF  ·  2026", 12, AMBER_LT, bold=True))])

text(s, Inches(5.5), Inches(1.55), Inches(7.2), Inches(2.0),
     [P(R("BookBridge", 62, WHITE, bold=True)),
      P(R("Where good books find", 25, AMBER_LT)),
      P(R("their next reader.", 25, AMBER_LT))],
     line_spacing=1.05)

text(s, Inches(5.5), Inches(4.35), Inches(7.0), Inches(1.2),
     [P(R("A campus community where readers ", 15, CREAM),
        R("gift, swap, or sell", 15, WHITE, bold=True),
        R(" the books they're done with — and trust the strangers they hand them to.", 15, CREAM))],
     line_spacing=1.3)

chip(s, Inches(5.5), Inches(5.75), Inches(2.5), "UN SDG 4 · Education", BRAND, WHITE, size=11)
chip(s, Inches(8.15), Inches(5.75), Inches(3.05), "UN SDG 12 · Responsible Use", AMBER, INK, size=11)

text(s, Inches(5.5), Inches(6.55), Inches(7), Inches(0.4),
     [P(R("Seeking grant & seed partners", 13, WHITE, bold=True),
        R("    ·    Team Zootopia", 13, RGBColor(0xA9,0xB4,0xAE)))])
page_num(s, 1, light=True)

# ============================================================================
# SLIDE 2 — THE SHELF  (open the story)
# ============================================================================
def motif_shelf(s):
    bx = Inches(9.75); by = Inches(2.95); bw = Inches(2.85)
    round_rect(s, bx, by, bw, Inches(3.0), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    cols = [MUTED, MUTED, AMBER, MUTED, MUTED, MUTED]
    for r in range(2):
        for c in range(3):
            i = r * 3 + c
            x = bx + Inches(0.3) + c * Inches(0.78)
            y = by + Inches(0.35) + r * Inches(1.2)
            round_rect(s, x, y, Inches(0.6), Inches(0.95), cols[i], radius=0.1)
        rect(s, bx + Inches(0.22), by + Inches(1.42) + r * Inches(1.2), bw - Inches(0.44), Pt(3), LINE)
    text(s, bx, by + Inches(2.55), bw, Inches(0.4),
         [P(R("one read, then forgotten", 11.5, MUTED, italic=True))],
         align=PP_ALIGN.CENTER)

story_slide(
    2, "Chapter 01 · The shelf",
    "Every shelf has a graveyard.",
    [P(R("Picture the last textbook you bought. You read it once, maybe twice, then "
         "it joined the quiet row of books behind you — too good to throw out, too "
         "forgotten to open again. ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("Now multiply that by every student on a single campus. Thousands of books "
         "sit frozen on shelves while the next reader pays full price for a copy that "
         "already exists three doors down.", 16.5, INK_SOFT, bold=False))],
    accent=CORAL, motif=motif_shelf, text_w=Inches(8.3),
    pull="The value is there. It's just stuck.",
)

# ============================================================================
# SLIDE 3 — WHY WE DON'T SHARE  (the tension)
# ============================================================================
def motif_barriers(s):
    items = [("Awkward & risky", "Handing a book to a stranger feels uneasy.", CORAL),
             ("Marketplaces sell", "Built for profit and shipping, not sharing.", AMBER),
             ("No way to trust", "No signal for who's actually reliable.", BLUE)]
    bx = Inches(9.55); by = Inches(2.95); bw = Inches(3.05)
    for i, (h, b, col) in enumerate(items):
        y = by + i * Inches(1.02)
        round_rect(s, bx, y, bw, Inches(0.9), CARD, line=LINE, line_w=Pt(1), radius=0.08)
        rect(s, bx, y, Inches(0.1), Inches(0.9), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        text(s, bx + Inches(0.28), y + Inches(0.13), bw - Inches(0.4), Inches(0.3),
             [P(R(h, 13, INK, bold=True))])
        text(s, bx + Inches(0.28), y + Inches(0.45), bw - Inches(0.4), Inches(0.4),
             [P(R(b, 10.5, MUTED))], line_spacing=1.05)

story_slide(
    3, "Chapter 02 · The friction",
    "So why don't we just pass them on?",
    [P(R("Because giving a book to someone you've never met is awkward, and a little "
         "risky. Will they actually show up? Are they who they say they are? ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("The usual options don't help. Marketplaces are built to ", 16.5, INK_SOFT),
       R("sell", 16.5, INK, bold=True),
       R(" — they push haggling, shipping and profit. Libraries have waitlists. So the "
         "book stays put, not because nobody wants it, but because there's no easy, "
         "trusted way to move it.", 16.5, INK_SOFT))],
    accent=AMBER, motif=motif_barriers, text_w=Inches(8.2),
)

# ============================================================================
# SLIDE 4 — THE REAL GAP  (the insight / why now)
# ============================================================================
def motif_now(s):
    bx = Inches(9.55); by = Inches(2.95); bw = Inches(3.05)
    round_rect(s, bx, by, bw, Inches(3.0), BRAND_LT, radius=0.06)
    text(s, bx + Inches(0.3), by + Inches(0.28), bw - Inches(0.6), Inches(0.4),
         [P(R("WHY NOW", 12, BRAND, bold=True))])
    tw = [("Reuse is a habit", "Sustainability is now a buying behaviour, not a sacrifice."),
          ("Community > counter", "Gen Z would rather belong than transact."),
          ("Campuses are dense", "High-trust networks where one seeds the next.")]
    for i, (h, b) in enumerate(tw):
        y = by + Inches(0.8) + i * Inches(0.72)
        text(s, bx + Inches(0.3), y, bw - Inches(0.6), Inches(0.3),
             [P(R("›  ", 13, BRAND, bold=True), R(h, 12.5, INK, bold=True))])
        text(s, bx + Inches(0.55), y + Inches(0.27), bw - Inches(0.8), Inches(0.4),
             [P(R(b, 10.5, INK_SOFT))], line_spacing=1.05)

story_slide(
    4, "Chapter 03 · The insight",
    "The missing piece was never payment. It was trust.",
    [P(R("We stopped thinking about how to take a cut of a sale, and started asking a "
         "simpler question: what would make two strangers comfortable handing a book "
         "between them? ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("The answer is a place where you can see who's reliable, agree on a hand-off, "
         "and feel safe doing it — especially inside a community you already belong to, "
         "like your campus. ", 16.5, INK_SOFT),
       R("Nobody owns that gap yet, and the moment to fill it is now.", 16.5, INK, bold=True))],
    accent=BRAND, motif=motif_now, text_w=Inches(8.2),
)

# ============================================================================
# SLIDE 5 — MEET BOOKBRIDGE  (the turn / solution)
# ============================================================================
s = slide()
bg(s)
rect(s, 0, 0, Inches(0.26), SH, BRAND)
text(s, Inches(0.95), Inches(0.78), Inches(9), Inches(0.4),
     [P(R("CHAPTER 04 · THE TURN", 12.5, BRAND, bold=True))])
text(s, Inches(0.95), Inches(1.28), Inches(11.4), Inches(0.9),
     [P(R("Meet BookBridge — one trusted place to keep a book moving.", 28, INK, bold=True))],
     line_spacing=1.05)
text(s, Inches(0.95), Inches(2.45), Inches(11.4), Inches(1.1),
     [P(R("Anyone can list a book in under a minute and choose how to pass it on. "
         "That single choice — give, swap, or sell — is the whole idea, wrapped in a "
         "community that already knows who to trust.", 16.5, INK_SOFT))],
     line_spacing=1.38)

modes = [
    ("🎁", "Gift", "Give it away free.\nPure circulation.", BRAND, BRAND_LT),
    ("🔄", "Exchange", "Trade book for book.\nNo money moves.", AMBER, AMBER_LT),
    ("🏷️", "Sell", "A symbolic price,\ncapped at 50,000 VND.", CORAL, CORAL_LT),
]
cx = Inches(0.95); cw = Inches(3.78); gap = Inches(0.18); cy = Inches(3.95)
for i, (g, h, b, col, light) in enumerate(modes):
    x = cx + i * (cw + gap)
    round_rect(s, x, cy, cw, Inches(1.85), light, radius=0.07)
    icon_circle(s, x + Inches(0.32), cy + Inches(0.32), Inches(0.8), g, WHITE, col, size=24)
    text(s, x + Inches(1.3), cy + Inches(0.42), cw - Inches(1.5), Inches(0.6),
         [P(R(h, 21, col, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + Inches(0.36), cy + Inches(1.2), cw - Inches(0.7), Inches(0.6),
         [P(R(b, 12.5, INK_SOFT))], line_spacing=1.12)

text(s, Inches(0.95), Inches(6.15), Inches(11.4), Inches(0.7),
     [P(R("The 50,000 VND cap is enforced in the code itself", 15, BRAND, bold=True, italic=True),
        R(" — which keeps BookBridge non-commercial, grant-eligible, and free of payment liability.",
          15, INK_SOFT, italic=True))],
     line_spacing=1.15)
page_num(s, 5, section="Chapter 04 · The turn")
footer_brand(s)

# ============================================================================
# SLIDE 6 — LISTING  (a book's journey begins)
# ============================================================================
def motif_listing(s):
    bx = Inches(9.55); by = Inches(2.95); bw = Inches(3.05)
    round_rect(s, bx, by, bw, Inches(3.0), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    round_rect(s, bx + Inches(0.3), by + Inches(0.3), Inches(0.95), Inches(1.35), AMBER_LT, radius=0.06)
    text(s, bx + Inches(0.3), by + Inches(0.78), Inches(0.95), Inches(0.4),
         [P(R("cover", 10, AMBER, italic=True))], align=PP_ALIGN.CENTER)
    lines = [("Title · Author", INK), ("auto-filled from ISBN", MUTED),
             ("Condition · Genre", INK), ("Gift / Swap / Sell", BRAND)]
    for i, (t, c) in enumerate(lines):
        y = by + Inches(0.34) + i * Inches(0.34)
        rect(s, bx + Inches(1.42), y + Inches(0.05), Inches(1.3), Pt(2), LINE)
        text(s, bx + Inches(1.42), y - Inches(0.05), bw - Inches(1.6), Inches(0.3),
             [P(R(t, 11, c, bold=(i in (0, 3))))])
    chip(s, bx + Inches(0.3), by + Inches(2.35), bw - Inches(0.6), "Posted in under 60 seconds",
         BRAND, WHITE, size=11, h=Inches(0.42))

story_slide(
    6, "Chapter 05 · The listing",
    "It starts with one photo and an ISBN.",
    [P(R("A book's journey begins the moment it leaves the shelf and becomes a listing. "
         "Paste an ISBN and the title, author and cover fill themselves in; snap a "
         "couple of photos, pick how you want to pass it on, and you're done — under a "
         "minute, start to finish. ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("The easier it is to list, the faster shelves empty", 16.5, INK, bold=True),
       R(" — and the more there is for the next reader to discover.", 16.5, INK_SOFT))],
    accent=AMBER, motif=motif_listing, text_w=Inches(8.2),
)

# ============================================================================
# SLIDE 7 — SEARCH & DISCOVERY  (a reader finds it)
# ============================================================================
def motif_search(s):
    bx = Inches(9.55); by = Inches(2.95); bw = Inches(3.05)
    round_rect(s, bx, by, bw, Inches(3.0), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    round_rect(s, bx + Inches(0.28), by + Inches(0.3), bw - Inches(0.56), Inches(0.5), BLUE_LT, radius=0.3)
    text(s, bx + Inches(0.5), by + Inches(0.3), bw - Inches(0.8), Inches(0.5),
         [P(R("🔍  calculus, good condition", 11.5, BLUE))], anchor=MSO_ANCHOR.MIDDLE)
    res = ["Calculus · Stewart", "Linear Algebra · Lay", "Physics · Halliday"]
    for i, t in enumerate(res):
        y = by + Inches(1.0) + i * Inches(0.5)
        round_rect(s, bx + Inches(0.28), y, bw - Inches(0.56), Inches(0.4), PAPER, radius=0.1)
        text(s, bx + Inches(0.45), y, bw - Inches(0.8), Inches(0.4),
             [P(R(t, 11, INK))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx, by + Inches(2.55), bw, Inches(0.4),
         [P(R("ranked results in under a second", 11, MUTED, italic=True))],
         align=PP_ALIGN.CENTER)

story_slide(
    7, "Chapter 06 · The search",
    "And somewhere, a reader is looking for exactly that.",
    [P(R("Search returns results in under a second — ranked across title, author and "
         "description, and narrowed by genre, condition, price, or simply how close the "
         "book is. ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("Between searches, a personal feed keeps surfacing fresh listings from the "
         "people and communities you follow. ", 16.5, INK_SOFT),
       R("Discovery is what brings readers back day after day", 16.5, INK, bold=True),
       R(" — and what keeps every book in circulation.", 16.5, INK_SOFT))],
    accent=BLUE, motif=motif_search, text_w=Inches(8.2),
)

# ============================================================================
# SLIDE 8 — CHAT & THE HAND-OFF  (two strangers, safely)
# ============================================================================
s = slide()
bg(s)
rect(s, 0, 0, Inches(0.26), SH, CORAL)
text(s, Inches(0.95), Inches(0.78), Inches(9), Inches(0.4),
     [P(R("CHAPTER 07 · THE HAND-OFF", 12.5, CORAL, bold=True))])
text(s, Inches(0.95), Inches(1.28), Inches(11.4), Inches(0.9),
     [P(R("Two strangers agree to meet — safely.", 30, INK, bold=True))], line_spacing=1.05)
text(s, Inches(0.95), Inches(2.35), Inches(11.5), Inches(1.5),
     [P(R("When someone wants your book, a chat opens right inside the deal — no phone "
         "numbers, no leaving the platform. A clear, well-tested flow walks both sides "
         "from request to hand-off to a final rating, with a waitlist if several people "
         "want the same title, and a moderator on call if anything ever goes wrong.",
         16.5, INK_SOFT))],
     line_spacing=1.38)

states = ["REQUEST", "ACCEPT", "MEET", "COMPLETE", "RATE"]
sx = Inches(0.95); sy = Inches(4.35); bw = Inches(1.9); bh = Inches(0.66); aw = Inches(0.4)
for i, st in enumerate(states):
    x = sx + i * (bw + aw)
    col = BRAND if st in ("COMPLETE", "RATE") else INK
    round_rect(s, x, sy, bw, bh, col, radius=0.2)
    text(s, x, sy, bw, bh, [P(R(st, 12, WHITE, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(states) - 1:
        text(s, x + bw, sy, aw, bh, [P(R("→", 19, CORAL, bold=True))], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

text(s, Inches(0.95), Inches(5.35), Inches(11.4), Inches(0.4),
     [P(R("Every awkward case is handled: ", 13, CORAL, bold=True),
        R("decline, cancel, no-show and dispute are all caught, queued, or escalated — never left hanging.",
          13, INK_SOFT))], line_spacing=1.15)
text(s, Inches(0.95), Inches(6.25), Inches(11.4), Inches(0.6),
     [P(R("The hand-off is the moment trust is earned — and the rating that follows feeds everything on the next slide.",
          15, CORAL, bold=True, italic=True))], line_spacing=1.15)
page_num(s, 8, section="Chapter 07 · The hand-off")
footer_brand(s)

# ============================================================================
# SLIDE 9 — COMMUNITY  (they belong now)
# ============================================================================
def motif_community(s):
    bx = Inches(9.55); by = Inches(2.95); bw = Inches(3.05)
    round_rect(s, bx, by, bw, Inches(3.0), CARD, line=LINE, line_w=Pt(1), radius=0.05)
    groups = [("🎓", "Your university", BRAND_LT, BRAND),
              ("📍", "Your city", BLUE_LT, BLUE),
              ("📖", "Your genres", AMBER_LT, AMBER)]
    for i, (g, t, light, col) in enumerate(groups):
        y = by + Inches(0.32) + i * Inches(0.62)
        round_rect(s, bx + Inches(0.28), y, bw - Inches(0.56), Inches(0.5), light, radius=0.15)
        text(s, bx + Inches(0.45), y, Inches(0.5), Inches(0.5), [P(R(g, 16, col))], anchor=MSO_ANCHOR.MIDDLE)
        text(s, bx + Inches(1.0), y, bw - Inches(1.2), Inches(0.5),
             [P(R(t, 12.5, INK, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx + Inches(0.28), by + Inches(2.3), bw - Inches(0.56), Inches(0.55),
         [P(R("Join → feed fills instantly with that group's books.", 11, MUTED, italic=True))],
         line_spacing=1.1)

story_slide(
    9, "Chapter 08 · The belonging",
    "One good exchange, and they're part of something.",
    [P(R("A reader doesn't just complete a deal — they join the communities that fit "
         "them: their university, their city, their favourite genres. Joining instantly "
         "fills their feed with books from that group, so there's never a cold, empty "
         "start. ", 16.5, INK_SOFT)),
     P(R("", 6, INK_SOFT)),
     P(R("Follows and shared spaces turn one-off swaps into a network that grows itself",
         16.5, INK, bold=True),
       R(" — one campus quietly inviting the next.", 16.5, INK_SOFT))],
    accent=BRAND, motif=motif_community, text_w=Inches(8.2),
)

# ============================================================================
# SLIDE 10 — TRUST, THE MOAT  (why it compounds)
# ============================================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(0.26), SH, AMBER)
text(s, Inches(0.95), Inches(0.72), Inches(9), Inches(0.4),
     [P(R("CHAPTER 09 · THE MOAT", 12.5, AMBER, bold=True))])
text(s, Inches(0.95), Inches(1.2), Inches(11.4), Inches(0.9),
     [P(R("Why this gets harder to copy every day.", 30, WHITE, bold=True))], line_spacing=1.05)
text(s, Inches(0.95), Inches(2.2), Inches(11.6), Inches(1.4),
     [P(R("Every completed deal and honest rating builds a reputation that's ", 16.5, CREAM),
        R("earned one exchange at a time", 16.5, AMBER, bold=True),
        R(" — from New Member to Community Champion. A competitor can copy our features in "
          "a weekend, but not years of accumulated trust. Quiet checks flag suspicious "
          "patterns for a human to review, so we protect the network without punishing "
          "real friendships.", 16.5, CREAM))],
     line_spacing=1.36)

tiers = [("0–19", "New Member", MUTED),
         ("20–49", "Active Sharer", BLUE),
         ("50–79", "Trusted Contributor", BRAND),
         ("80–100", "Community Champion", AMBER)]
tx = Inches(0.95); tw = Inches(2.83); th = Inches(1.25); tg = Inches(0.13); ty = Inches(4.35)
for i, (rng, name, col) in enumerate(tiers):
    x = tx + i * (tw + tg)
    round_rect(s, x, ty, tw, th, PANEL_DK, radius=0.08)
    rect(s, x, ty, tw, Inches(0.12), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, ty + Inches(0.28), tw, Inches(0.4), [P(R(rng, 19, col, bold=True))], align=PP_ALIGN.CENTER)
    text(s, x, ty + Inches(0.74), tw, Inches(0.45), [P(R(name, 13, WHITE, bold=True))], align=PP_ALIGN.CENTER)

text(s, Inches(0.95), Inches(6.05), Inches(11.4), Inches(0.6),
     [P(R("The longer BookBridge runs, the more trustworthy it becomes — and the wider the moat.",
          15, AMBER, bold=True, italic=True))], line_spacing=1.1)
page_num(s, 10, section="Chapter 09 · The moat", light=True)

# ============================================================================
# SLIDE 11 — MODEL & SCALE  (how it lasts)
# ============================================================================
s = slide()
bg(s)
rect(s, 0, 0, Inches(0.26), SH, BRAND)
text(s, Inches(0.95), Inches(0.78), Inches(9), Inches(0.4),
     [P(R("CHAPTER 10 · HOW IT LASTS", 12.5, BRAND, bold=True))])
text(s, Inches(0.95), Inches(1.28), Inches(11.4), Inches(0.9),
     [P(R("Built to last on very little.", 30, INK, bold=True))], line_spacing=1.05)
text(s, Inches(0.95), Inches(2.3), Inches(11.5), Inches(1.5),
     [P(R("BookBridge runs on a deliberately lean stack — one database, no payment "
         "rails, no message broker — so it can hold a thousand readers at once while "
         "staying cheap to operate. Because it's non-commercial by design, it qualifies "
         "for education and sustainability grants, and reports its impact to sponsors at "
         "the click of a button.", 16.5, INK_SOFT))],
     line_spacing=1.36)

metrics = [
    ("1,000+", "concurrent users, one DB", AMBER),
    ("99.5%", "uptime target", PURPLE),
    ("< 3 s", "page load on 4G", BLUE),
    ("≥ 70%", "automated test coverage", BRAND),
]
mx = Inches(0.95); mw = Inches(2.83); mh = Inches(1.35); mg = Inches(0.13); my = Inches(4.35)
for i, (big, lab, col) in enumerate(metrics):
    x = mx + i * (mw + mg)
    round_rect(s, x, my, mw, mh, CARD, line=LINE, line_w=Pt(1), radius=0.06)
    rect(s, x, my, Inches(0.12), mh, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x + Inches(0.34), my + Inches(0.24), mw - Inches(0.45), Inches(0.7),
         [P(R(big, 29, col, bold=True))])
    text(s, x + Inches(0.36), my + Inches(0.9), mw - Inches(0.5), Inches(0.4),
         [P(R(lab, 11.5, INK_SOFT))], line_spacing=1.0)

text(s, Inches(0.95), Inches(6.05), Inches(11.4), Inches(0.6),
     [P(R("Low burn, clear mission, real headroom", 15, BRAND, bold=True, italic=True),
        R(" — the runway goes to growth, not infrastructure.", 15, INK_SOFT, italic=True))],
     line_spacing=1.1)
page_num(s, 11, section="Chapter 10 · How it lasts")
footer_brand(s)

# ============================================================================
# SLIDE 12 — CLOSE & THE ASK
# ============================================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(0.3), SH, BRAND)
text(s, Inches(0.95), Inches(0.85), Inches(11.5), Inches(1.4),
     [P(R("Keep books moving —", 33, WHITE, bold=True)),
      P(R("and build the trust that makes sharing work.", 33, AMBER, bold=True))],
     line_spacing=1.08)

text(s, Inches(0.95), Inches(2.65), Inches(11.4), Inches(0.9),
     [P(R("We've built the whole platform — listings, search, chat, communities and a "
         "working trust engine. It's live today. What we need now is the runway and the "
         "partners to put it in front of the first campuses.", 15.5, CREAM))],
     line_spacing=1.34)

text(s, Inches(0.95), Inches(4.2), Inches(11.5), Inches(0.4),
     [P(R("THE ASK", 12.5, AMBER, bold=True))])
asks = ["Grant & seed partners to fund the next 12 months",
        "Pilot partnerships with 3 universities",
        "Advisors in community growth and trust & safety",
        "A platform already live and ready to scale"]
for i, t in enumerate(asks):
    c = i % 2; r = i // 2
    text(s, Inches(0.95) + c * Inches(5.9), Inches(4.6) + r * Inches(0.5), Inches(5.7), Inches(0.4),
         [P(R("→  ", 14, BRAND_LT, bold=True), R(t, 13.5, WHITE))])

round_rect(s, Inches(0.95), Inches(6.05), Inches(11.45), Inches(0.7), BRAND_DK, radius=0.18)
text(s, Inches(0.95), Inches(6.05), Inches(11.45), Inches(0.7),
     [P(R("Team Zootopia   ·   Let's keep books in circulation, not in landfills.", 13.5, WHITE, bold=True))],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 12, light=True)

import os
out = "BookBridge_Investor_Deck.pptx"
try:
    prs.save(out)
except PermissionError:
    out = "BookBridge_Investor_Deck_v2.pptx"
    prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")

"""Generate BookBridge_Pitch.pptx using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = r"C:\Users\ADMIN\Desktop\SE\SE_book_bridge\docs\BookBridge_Pitch.pptx"

# Colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_GRAY = RGBColor(0x33, 0x41, 0x55)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SLATE_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
INDIGO_SOFT = RGBColor(0x63, 0x66, 0xF1)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]


def add_rect(slide, left, top, width, height, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    return shp


def set_text(shape, text, size=18, bold=False, color=BODY_GRAY, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # clear default run
        for r in list(p.runs):
            r.text = ""
        run = p.add_run() if not p.runs else p.runs[0]
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tf


def add_text(slide, left, top, width, height, text, **kwargs):
    tb = slide.shapes.add_textbox(left, top, width, height)
    set_text(tb, text, **kwargs)
    return tb


def add_multitext(slide, left, top, width, height, runs, anchor=MSO_ANCHOR.TOP):
    """runs: list of paragraphs, each is list of (text, size, bold, color, italic)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        for j, item in enumerate(para):
            text, size, bold, color, italic = item
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return tb


def add_slide_number(slide, n, total, dark=False):
    color = SLATE_LIGHT if not dark else SLATE_LIGHT
    add_text(slide, Inches(12.6), Inches(7.15), Inches(0.6), Inches(0.3),
             f"{n} / {total}", size=10, color=color, align=PP_ALIGN.RIGHT)


def add_arrow(slide, x1, y1, x2, y2, color=INDIGO, weight=2.25):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # add arrow head
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


TOTAL = 10

# =============== SLIDE 1 — Cover ===============
s1 = prs.slides.add_slide(blank_layout)
add_rect(s1, 0, 0, SW, SH, fill=NAVY)
add_text(s1, Inches(0.5), Inches(0.4), Inches(6), Inches(0.3),
         "SE Capstone 2025.2 · Team Zootopia",
         size=11, color=SLATE_LIGHT)
add_text(s1, Inches(0.5), Inches(2.9), Inches(12.33), Inches(1.4),
         "BookBridge", size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.6),
         "A trust-driven platform that keeps books in circulation",
         size=20, color=SLATE_LIGHT, align=PP_ALIGN.CENTER)
add_text(s1, Inches(0.5), Inches(7.0), Inches(10), Inches(0.3),
         "Next.js · PostgreSQL · Prisma · Tailwind · Vercel",
         size=11, color=INDIGO_SOFT)

# =============== SLIDE 2 — The Problem ===============
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, SW, SH, fill=WHITE)
add_text(s2, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8),
         "The Problem", size=36, bold=True, color=NAVY)

problems = [
    ("Books sit on shelves",
     "Students buy textbooks for one semester. When the course ends, the books stop moving."),
    ("No safe coordination layer",
     "Social media groups lack structure. Commercial marketplaces are built for profit, not community."),
    ("Trust is missing",
     "Without accountability, strangers will not exchange. Without reputation, every transaction carries risk."),
]
top = Inches(1.8)
for label, desc in problems:
    # left accent line
    add_rect(s2, Inches(0.6), top, Inches(0.06), Inches(1.3), fill=INDIGO)
    add_multitext(s2, Inches(0.85), top, Inches(11.5), Inches(1.3), [
        [(label, 22, True, NAVY, False)],
        [(desc, 18, False, BODY_GRAY, False)],
    ])
    top += Inches(1.6)

add_slide_number(s2, 2, TOTAL)

# =============== SLIDE 3 — What BookBridge Does ===============
s3 = prs.slides.add_slide(blank_layout)
add_rect(s3, 0, 0, SW, SH, fill=WHITE)
add_text(s3, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
         "What BookBridge Does", size=36, bold=True, color=NAVY)
add_text(s3, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.2),
         "BookBridge is a web platform where registered users list pre-owned "
         "books, request them from other users, coordinate the exchange, and "
         "earn reputation as they complete transactions.",
         size=18, color=BODY_GRAY)

cards = [
    ("GIFT", EMERALD, "Pass a book to someone who needs it at zero cost."),
    ("EXCHANGE", INDIGO, "Swap one book for another with a matched reader."),
    ("SELL", AMBER, "Recover symbolic costs. Price capped at 50,000 VND by the platform."),
]
card_w = Inches(3.9)
card_h = Inches(2.6)
gap = Inches(0.25)
total_w = card_w * 3 + gap * 2
start_left = (SW - total_w) // 2
ctop = Inches(4.1)
for i, (label, col, body) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    rect = add_rect(s3, left, ctop, card_w, card_h, fill=LIGHT_BG, line=BORDER,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect.adjustments[0] = 0.08
    add_text(s3, left + Inches(0.3), ctop + Inches(0.3), card_w - Inches(0.6), Inches(0.6),
             label, size=22, bold=True, color=col)
    add_text(s3, left + Inches(0.3), ctop + Inches(1.0), card_w - Inches(0.6), card_h - Inches(1.2),
             body, size=18, color=BODY_GRAY)

add_slide_number(s3, 3, TOTAL)

# =============== SLIDE 4 — Six Modules ===============
s4 = prs.slides.add_slide(blank_layout)
add_rect(s4, 0, 0, SW, SH, fill=WHITE)
add_text(s4, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 1", size=11, bold=True, color=INDIGO)
add_text(s4, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
         "Six Modules, One Platform", size=36, bold=True, color=NAVY)

modules = [
    ("1", "Identity and Profile", "auth, sessions, user profiles"),
    ("2", "Book Catalog", "listings, photos, ISBN lookup"),
    ("3", "Search and Discovery", "full text search, feed, follows"),
    ("4", "Transactions and Messaging", "state machine, ratings, chat"),
    ("5", "Trust and Safety", "reputation engine, reports, moderation"),
    ("6", "Community and Ops", "communities, notifications, admin, CI/CD"),
]
row_top = Inches(1.85)
row_h = Inches(0.65)
for i, (num, label, desc) in enumerate(modules):
    rt = row_top + row_h * i
    add_text(s4, Inches(0.7), rt, Inches(0.6), row_h,
             num, size=26, bold=True, color=INDIGO, anchor=MSO_ANCHOR.MIDDLE)
    add_multitext(s4, Inches(1.4), rt, Inches(11.3), row_h, [
        [(label, 18, True, NAVY, False),
         ("   —   ", 18, False, SLATE_LIGHT, False),
         (desc, 18, False, BODY_GRAY, False)],
    ], anchor=MSO_ANCHOR.MIDDLE)

add_text(s4, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
         "One shared Prisma schema. Zero handoff bottlenecks.",
         size=14, italic=True, color=SLATE_LIGHT)

add_slide_number(s4, 4, TOTAL)

# =============== SLIDE 5 — Eight Transaction States ===============
s5 = prs.slides.add_slide(blank_layout)
add_rect(s5, 0, 0, SW, SH, fill=WHITE)
add_text(s5, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 2", size=11, bold=True, color=INDIGO)
add_text(s5, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
         "Eight Transaction States", size=36, bold=True, color=NAVY)


def state_box(slide, left, top, w, h, text, fill=INDIGO, txt_color=WHITE):
    rect = add_rect(slide, left, top, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect.adjustments[0] = 0.25
    set_text(rect, text, size=14, bold=True, color=txt_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return rect


# Top row: PENDING -> ACCEPTED -> IN DELIVERY -> COMPLETED
box_w = Inches(2.2)
box_h = Inches(0.8)
arrow_gap = Inches(0.45)
top_y = Inches(2.6)

# Compute positions
positions = []
left = Inches(0.7)
for i in range(4):
    positions.append(left)
    left = left + box_w + arrow_gap

labels_top = ["PENDING", "ACCEPTED", "IN DELIVERY", "COMPLETED"]
fills_top = [INDIGO, INDIGO, INDIGO, EMERALD]
for pos, lbl, fc in zip(positions, labels_top, fills_top):
    state_box(s5, pos, top_y, box_w, box_h, lbl, fill=fc)

# arrows between top row
for i in range(3):
    x1 = positions[i] + box_w
    x2 = positions[i + 1]
    y = top_y + box_h // 2
    add_arrow(s5, x1, y, x2, y)

# Bottom row offshoots
bot_y = Inches(4.4)
# Below PENDING: DECLINED
state_box(s5, positions[0], bot_y, box_w, box_h, "DECLINED", fill=SLATE_LIGHT)
add_arrow(s5, positions[0] + box_w // 2, top_y + box_h,
          positions[0] + box_w // 2, bot_y)
# Below ACCEPTED: CANCELLED
state_box(s5, positions[1], bot_y, box_w, box_h, "CANCELLED", fill=SLATE_LIGHT)
add_arrow(s5, positions[1] + box_w // 2, top_y + box_h,
          positions[1] + box_w // 2, bot_y)
# Below IN DELIVERY: DISPUTED -> COMPLETED or CANCELLED
state_box(s5, positions[2], bot_y, box_w, box_h, "DISPUTED", fill=AMBER)
add_arrow(s5, positions[2] + box_w // 2, top_y + box_h,
          positions[2] + box_w // 2, bot_y)

# DISPUTED -> COMPLETED or CANCELLED (note text)
add_text(s5, positions[3] - Inches(0.2), bot_y, Inches(2.6), box_h,
         "→  COMPLETED or CANCELLED",
         size=13, bold=True, color=BODY_GRAY, anchor=MSO_ANCHOR.MIDDLE)
add_arrow(s5, positions[2] + box_w, bot_y + box_h // 2,
          positions[3] - Inches(0.2), bot_y + box_h // 2)

# Caption
add_text(s5, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.2),
         "Every transition is logged to an immutable audit table. "
         "The state machine is a pure function with 13 passing unit tests.",
         size=16, color=BODY_GRAY, italic=True)

add_slide_number(s5, 5, TOTAL)

# =============== SLIDE 6 — Four Reputation Tiers ===============
s6 = prs.slides.add_slide(blank_layout)
add_rect(s6, 0, 0, SW, SH, fill=WHITE)
add_text(s6, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 3", size=11, bold=True, color=INDIGO)
add_text(s6, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
         "Four Reputation Tiers", size=36, bold=True, color=NAVY)

tiers = [
    ("0 to 19", "New Member", SLATE_LIGHT),
    ("20 to 49", "Active Sharer", EMERALD),
    ("50 to 79", "Trusted Contributor", INDIGO),
    ("80 to 100", "Community Champion", AMBER),
]
tcw = Inches(2.85)
tch = Inches(2.6)
tgap = Inches(0.2)
ttotal = tcw * 4 + tgap * 3
tstart = (SW - ttotal) // 2
ttop = Inches(2.0)
for i, (rng, name, col) in enumerate(tiers):
    left = tstart + (tcw + tgap) * i
    rect = add_rect(s6, left, ttop, tcw, tch, fill=LIGHT_BG, line=BORDER,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect.adjustments[0] = 0.06
    # color bar
    add_rect(s6, left + Inches(0.2), ttop + Inches(0.3),
             tcw - Inches(0.4), Inches(0.18), fill=col)
    add_text(s6, left + Inches(0.2), ttop + Inches(0.65),
             tcw - Inches(0.4), Inches(0.5),
             rng, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s6, left + Inches(0.2), ttop + Inches(1.3),
             tcw - Inches(0.4), Inches(1.2),
             name, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

add_text(s6, Inches(0.6), Inches(4.95), Inches(12.1), Inches(0.5),
         "Score events: +10 transaction completed, +/- by star rating, "
         "-3 cancellation, -15 report upheld, -1 per 30 inactive days.",
         size=15, color=BODY_GRAY)
add_text(s6, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.5),
         "Anti-gaming detection flags reciprocal trading pairs automatically "
         "for moderator review.",
         size=15, color=BODY_GRAY)

add_slide_number(s6, 6, TOTAL)

# =============== SLIDE 7 — Price Cap ===============
s7 = prs.slides.add_slide(blank_layout)
add_rect(s7, 0, 0, SW, SH, fill=WHITE)
add_text(s7, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 4", size=11, bold=True, color=INDIGO)
add_text(s7, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
         "Three Exchange Types. One Hard Rule.",
         size=36, bold=True, color=NAVY)

add_text(s7, Inches(0.6), Inches(1.9), Inches(12.1), Inches(1.0),
         "50,000 VND", size=48, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
add_text(s7, Inches(0.6), Inches(2.95), Inches(12.1), Inches(0.5),
         "Maximum sell price, enforced server-side on every listing creation and edit.",
         size=16, color=BODY_GRAY, align=PP_ALIGN.CENTER, italic=True)

# Two columns
col_w = Inches(5.6)
col_h = Inches(3.0)
col_gap = Inches(0.4)
col_top = Inches(3.8)
col_left1 = Inches(0.7)
col_left2 = col_left1 + col_w + col_gap

add_rect(s7, col_left1, col_top, col_w, col_h, fill=LIGHT_BG, line=BORDER,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_rect(s7, col_left2, col_top, col_w, col_h, fill=LIGHT_BG, line=BORDER,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)

add_text(s7, col_left1 + Inches(0.3), col_top + Inches(0.2),
         col_w - Inches(0.6), Inches(0.5),
         "What this enforces", size=18, bold=True, color=INDIGO)
left_items = [
    "Non-commercial mission",
    "Community over commerce",
    "Grant funding alignment",
    "SDG 12 Responsible Consumption",
]
add_text(s7, col_left1 + Inches(0.3), col_top + Inches(0.85),
         col_w - Inches(0.6), col_h - Inches(1.0),
         "\n".join(left_items), size=16, color=BODY_GRAY)

add_text(s7, col_left2 + Inches(0.3), col_top + Inches(0.2),
         col_w - Inches(0.6), Inches(0.5),
         "What this enables", size=18, bold=True, color=EMERALD)
right_items = [
    "Shipping cost recovery",
    "Symbolic fair exchange",
    "Trust between strangers",
    "Sustainable book circulation",
]
add_text(s7, col_left2 + Inches(0.3), col_top + Inches(0.85),
         col_w - Inches(0.6), col_h - Inches(1.0),
         "\n".join(right_items), size=16, color=BODY_GRAY)

add_slide_number(s7, 7, TOTAL)

# =============== SLIDE 8 — Real-Time SSE ===============
s8 = prs.slides.add_slide(blank_layout)
add_rect(s8, 0, 0, SW, SH, fill=WHITE)
add_text(s8, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 5", size=11, bold=True, color=INDIGO)
add_text(s8, Inches(0.6), Inches(0.75), Inches(12), Inches(0.8),
         "Real-Time in a Single Runtime", size=36, bold=True, color=NAVY)

sse_cards = [
    ("/api/notifications/stream", "Live notification bell updates"),
    ("/api/conversations/[id]/stream", "Per-conversation message push"),
    ("/api/feed/stream", "Personalized activity feed"),
]
sw_card = Inches(3.9)
sh_card = Inches(2.4)
sg = Inches(0.25)
stotal = sw_card * 3 + sg * 2
sstart = (SW - stotal) // 2
stop = Inches(2.1)
for i, (label, body) in enumerate(sse_cards):
    left = sstart + (sw_card + sg) * i
    rect = add_rect(s8, left, stop, sw_card, sh_card, fill=LIGHT_BG, line=BORDER,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect.adjustments[0] = 0.06
    add_text(s8, left + Inches(0.25), stop + Inches(0.35),
             sw_card - Inches(0.5), Inches(0.7),
             label, size=15, bold=True, color=INDIGO, font="Consolas")
    add_text(s8, left + Inches(0.25), stop + Inches(1.2),
             sw_card - Inches(0.5), sh_card - Inches(1.4),
             body, size=18, color=BODY_GRAY)

add_text(s8, Inches(0.6), Inches(5.0), Inches(12.1), Inches(0.6),
         "No WebSockets. No broker. No additional service. "
         "HTTP-native Server-Sent Events with automatic reconnect.",
         size=17, color=BODY_GRAY)
add_text(s8, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.6),
         "The entire product deploys as one Next.js application.",
         size=17, color=BODY_GRAY, italic=True)

add_slide_number(s8, 8, TOTAL)

# =============== SLIDE 9 — Scale ===============
s9 = prs.slides.add_slide(blank_layout)
add_rect(s9, 0, 0, SW, SH, fill=WHITE)
add_text(s9, Inches(0.6), Inches(0.4), Inches(3), Inches(0.3),
         "Figure 6", size=11, bold=True, color=INDIGO)
add_text(s9, Inches(0.6), Inches(0.75), Inches(12.5), Inches(0.8),
         "1,000 Concurrent Users on Standard Infrastructure",
         size=32, bold=True, color=NAVY)

metrics = [
    ("< 2s", "Search response at 95th percentile"),
    ("99.5%", "Monthly uptime target"),
    ("< 3s", "Initial page load over 4G"),
]
mw = Inches(3.9)
mh = Inches(2.0)
mg = Inches(0.25)
mtotal = mw * 3 + mg * 2
mstart = (SW - mtotal) // 2
mtop = Inches(2.0)
for i, (big, small) in enumerate(metrics):
    left = mstart + (mw + mg) * i
    rect = add_rect(s9, left, mtop, mw, mh, fill=LIGHT_BG, line=BORDER,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect.adjustments[0] = 0.07
    add_text(s9, left, mtop + Inches(0.2), mw, Inches(0.9),
             big, size=40, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
    add_text(s9, left + Inches(0.2), mtop + Inches(1.2),
             mw - Inches(0.4), Inches(0.7),
             small, size=15, color=BODY_GRAY, align=PP_ALIGN.CENTER)

add_text(s9, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.5),
         "How these are met:", size=18, bold=True, color=NAVY)

bullets = [
    "Reputation score denormalized onto user row — no aggregate queries on listing cards",
    "GIN indexed PostgreSQL full-text search",
    "Cursor-based feed pagination",
    "Single Vercel deployment plus managed Postgres",
]
add_text(s9, Inches(0.6), Inches(5.0), Inches(12.1), Inches(2.0),
         "\n".join(bullets), size=16, color=BODY_GRAY)

add_slide_number(s9, 9, TOTAL)

# =============== SLIDE 10 — What Was Built ===============
s10 = prs.slides.add_slide(blank_layout)
add_rect(s10, 0, 0, SW, SH, fill=NAVY)
add_text(s10, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8),
         "What Was Built", size=36, bold=True, color=WHITE)

product_items = [
    "23 pages across 6 feature areas",
    "50+ API routes",
    "13 database models",
    "Real-time SSE across 3 channels",
    "Formal 8-state transaction machine",
    "Reputation engine with anti-gaming",
    "Community sub-groups",
    "Admin dashboard with grant CSV export",
]
engineering_items = [
    "Next.js 14 App Router, TypeScript end to end",
    "PostgreSQL 16, Prisma 5, 30+ indexes",
    "bcryptjs auth, iron-session cookies",
    "Vitest unit tests, 70% coverage target",
    "GitHub Actions CI with hosted PostgreSQL",
    "Vercel deploy, auto-deploy on push to main",
]

col_w10 = Inches(5.9)
col_top10 = Inches(1.7)
col_h10 = Inches(4.7)

add_text(s10, Inches(0.7), col_top10, col_w10, Inches(0.5),
         "Product", size=22, bold=True, color=INDIGO_SOFT)
add_text(s10, Inches(0.7), col_top10 + Inches(0.6), col_w10, col_h10 - Inches(0.6),
         "\n".join(product_items), size=16, color=WHITE)

add_text(s10, Inches(7.0), col_top10, col_w10, Inches(0.5),
         "Engineering", size=22, bold=True, color=EMERALD)
add_text(s10, Inches(7.0), col_top10 + Inches(0.6), col_w10, col_h10 - Inches(0.6),
         "\n".join(engineering_items), size=16, color=WHITE)

add_text(s10, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         "Built by 6 people, most building a web application for the first time, in 8 weeks.",
         size=16, italic=True, color=SLATE_LIGHT, align=PP_ALIGN.CENTER)

add_text(s10, Inches(12.6), Inches(7.15), Inches(0.6), Inches(0.3),
         f"10 / {TOTAL}", size=10, color=SLATE_LIGHT, align=PP_ALIGN.RIGHT)


# Slide numbers for slide 1 already omitted (cover)
# Slide 1 has no number per design

prs.save(OUT)
print("OK", OUT)

"""
BookBridge_Pitch.pptx generator
Run: python gen_pitch.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
BG        = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NEARWHITE = RGBColor(0xF1, 0xF5, 0xF9)
INDIGO    = RGBColor(0x63, 0x66, 0xF1)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
GREY_DIM  = RGBColor(0x1E, 0x29, 0x3B)   # card background
GREY_ROW  = RGBColor(0x1E, 0x29, 0x3B)   # alt table row
INDIGO_DK = RGBColor(0x45, 0x48, 0xC8)   # table header

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a fully blank slide (no placeholders)."""
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    # Set background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_name="Calibri", font_size=18,
                bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True,
                italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_rect(slide, left, top, width, height,
             fill_color=GREY_DIM, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height,
                     fill_color=GREY_DIM, line_color=INDIGO, line_width=Pt(1.5)):
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width     = line_width
    # set rounding
    shape.adjustments[0] = 0.05
    return shape


def label_in_rect(slide, text, left, top, width, height,
                  fill=GREY_DIM, border=INDIGO,
                  font_size=13, color=WHITE, bold=False,
                  align=PP_ALIGN.CENTER):
    add_rounded_rect(slide, left, top, width, height, fill_color=fill, line_color=border)
    add_textbox(slide, text, left, top, width, height,
                font_size=font_size, color=color, bold=bold,
                align=align, wrap=True)


def tag(slide, text):
    """Small tag in top-right corner."""
    add_textbox(slide, text,
                Inches(10.8), Inches(0.08), Inches(2.4), Inches(0.35),
                font_size=9, color=INDIGO, align=PP_ALIGN.RIGHT, italic=True)


def slide_heading(slide, text, top=Inches(0.28)):
    add_textbox(slide, text,
                Inches(0.4), top, Inches(12.0), Inches(0.55),
                font_size=26, bold=True, color=INDIGO, align=PP_ALIGN.LEFT)


def slide_subheading(slide, text, top=Inches(0.88)):
    add_textbox(slide, text,
                Inches(0.4), top, Inches(12.0), Inches(0.38),
                font_size=14, bold=False, color=NEARWHITE, align=PP_ALIGN.LEFT, italic=True)


def divider(slide, top):
    rect = slide.shapes.add_shape(1,
        Inches(0.4), top, Inches(12.53), Inches(0.03))
    rect.fill.solid()
    rect.fill.fore_color.rgb = INDIGO
    rect.line.fill.background()


def add_table(slide, rows, cols, left, top, width, height,
              headers=None, col_widths=None):
    """Add a table with indigo header row and alternating body rows."""
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def cell_fill(cell, color):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color

    def cell_text(cell, text, bold=False, font_size=12, color=WHITE,
                  align=PP_ALIGN.LEFT):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            if r == 0 and headers:
                cell_fill(cell, INDIGO_DK)
                cell_text(cell, headers[c] if c < len(headers) else "",
                          bold=True, font_size=12)
            else:
                alt = (r % 2 == 0)
                cell_fill(cell, RGBColor(0x16, 0x20, 0x30) if alt else GREY_ROW)

    return tbl


def fill_table_row(tbl, row_idx, values, font_size=11, color=WHITE):
    for c, val in enumerate(values):
        cell = tbl.cell(row_idx, c)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(val)
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color


def arrow_text(slide, left, top, width, height, text="v", color=INDIGO):
    add_textbox(slide, text, left, top, width, height,
                font_size=16, color=color, align=PP_ALIGN.CENTER, bold=True)


# ---------------------------------------------------------------------------
# SLIDE 1: Cover
# ---------------------------------------------------------------------------

def slide_01(prs):
    slide = blank_slide(prs)
    tag(slide, "BookBridge | Cover")

    # Large title
    add_textbox(slide, "BookBridge",
                Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.4),
                font_size=72, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide,
                "A trust-driven book sharing platform for campus communities",
                Inches(1.2), Inches(3.35), Inches(10.9), Inches(0.7),
                font_size=22, color=NEARWHITE, align=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide,
                "SE 2025.2  |  Team Zootopia  |  6 members",
                Inches(1.2), Inches(4.1), Inches(10.9), Inches(0.5),
                font_size=14, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

    # Bottom accent bar
    rect = slide.shapes.add_shape(1,
        Inches(0), Inches(7.1), Inches(13.33), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = INDIGO
    rect.line.fill.background()

    return slide


# ---------------------------------------------------------------------------
# SLIDE 2: The Problem
# ---------------------------------------------------------------------------

def slide_02(prs):
    slide = blank_slide(prs)
    tag(slide, "Problem")
    slide_heading(slide, "Books are expensive. Sharing is broken.")
    divider(slide, Inches(0.88))

    cards = [
        ("Books bought once,\nshelved forever", INDIGO),
        ("Facebook Marketplace\nhas no trust layer", AMBER),
        ("Library waitlists\nare long", GREEN),
    ]
    cx = Inches(0.5)
    for title, accent in cards:
        add_rounded_rect(slide, cx, Inches(1.15), Inches(3.9), Inches(5.5),
                         fill_color=GREY_DIM, line_color=accent, line_width=Pt(2))
        add_textbox(slide, title,
                    cx + Inches(0.15), Inches(2.8), Inches(3.6), Inches(1.5),
                    font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += Inches(4.2)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 3: The Solution
# ---------------------------------------------------------------------------

def slide_03(prs):
    slide = blank_slide(prs)
    tag(slide, "Solution")
    slide_heading(slide, "Three exchange modes. One trusted platform.")
    divider(slide, Inches(0.88))

    # Table
    headers = ["Mode", "Description"]
    col_widths = [Inches(2.5), Inches(8.5)]
    tbl = add_table(slide, 4, 2,
                    Inches(1.5), Inches(1.1),
                    Inches(10.33), Inches(3.5),
                    headers=headers, col_widths=col_widths)

    rows_data = [
        ("Gift",     "Owner gives the book away for free"),
        ("Exchange", "Owner trades for another book"),
        ("Sell",     "Price capped at 50,000 VND (non-commercial)"),
    ]
    for i, (mode, desc) in enumerate(rows_data):
        fill_table_row(tbl, i + 1, [mode, desc], font_size=13)
        # colour the mode cell
        c = tbl.cell(i + 1, 0)
        clr = [GREEN, INDIGO, AMBER][i]
        for run in c.text_frame.paragraphs[0].runs:
            run.font.color.rgb = clr
            run.font.bold = True

    # Note below
    add_textbox(slide,
                "The 50,000 VND cap is enforced server-side. Not advisory.",
                Inches(1.5), Inches(4.85), Inches(10.33), Inches(0.5),
                font_size=13, color=AMBER, align=PP_ALIGN.CENTER, italic=True, bold=True)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 4: Six Modules
# ---------------------------------------------------------------------------

def slide_04(prs):
    slide = blank_slide(prs)
    tag(slide, "Architecture: Modules")
    slide_heading(slide, "Six modules. One platform.")
    divider(slide, Inches(0.88))

    # Flow diagram
    bw, bh = Inches(2.8), Inches(0.55)
    gap = Inches(0.15)

    # Row 1: two boxes
    r1y = Inches(1.05)
    boxes_r1 = [
        (Inches(1.0),  "Identity and Profile", INDIGO),
        (Inches(4.53), "Book Catalog",          INDIGO),
    ]
    for bx, label, col in boxes_r1:
        add_rounded_rect(slide, bx, r1y, bw, bh, fill_color=GREY_DIM, line_color=col)
        add_textbox(slide, label, bx, r1y, bw, bh,
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Arrow down
    arrow_text(slide, Inches(3.3), r1y + bh, Inches(1.5), Inches(0.4))

    # Row 2: centre
    r2y = r1y + bh + Inches(0.42)
    add_rounded_rect(slide, Inches(2.5), r2y, bw, bh, fill_color=GREY_DIM, line_color=GREEN)
    add_textbox(slide, "Discovery and Feed", Inches(2.5), r2y, bw, bh,
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    arrow_text(slide, Inches(3.3), r2y + bh, Inches(1.5), Inches(0.4))

    # Row 3: two boxes
    r3y = r2y + bh + Inches(0.42)
    boxes_r3 = [
        (Inches(0.5),  "Transactions and Messaging", AMBER),
        (Inches(4.0),  "Trust and Safety",           RGBColor(0xEF, 0x44, 0x44)),
    ]
    for bx, label, col in boxes_r3:
        add_rounded_rect(slide, bx, r3y, Inches(3.2), bh, fill_color=GREY_DIM, line_color=col)
        add_textbox(slide, label, bx, r3y, Inches(3.2), bh,
                    font_size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    arrow_text(slide, Inches(3.3), r3y + bh, Inches(1.5), Inches(0.4))

    # Row 4: Community and Ops
    r4y = r3y + bh + Inches(0.42)
    add_rounded_rect(slide, Inches(2.0), r4y, Inches(3.8), bh,
                     fill_color=GREY_DIM, line_color=AMBER)
    add_textbox(slide, "Community and Ops", Inches(2.0), r4y, Inches(3.8), bh,
                font_size=12, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Table on the right side
    headers = ["Module", "Owner", "Core responsibility"]
    col_widths = [Inches(2.3), Inches(1.2), Inches(3.6)]
    tbl = add_table(slide, 7, 3,
                    Inches(7.0), Inches(1.05),
                    Inches(6.0), Inches(5.8),
                    headers=headers, col_widths=col_widths)

    module_data = [
        ("Identity and Profile",        "Person 1", "Auth, sessions, user profiles"),
        ("Book Catalog",                "Person 2", "Listings CRUD, photos, ISBN lookup"),
        ("Discovery and Feed",          "Person 3", "Full-text search, feed, follow graph"),
        ("Transactions and Messaging",  "Person 4", "State machine, chat, ratings"),
        ("Trust and Safety",            "Person 5", "Reputation engine, reports, moderation"),
        ("Community and Ops",           "Person 6", "Communities, notifications, admin, CI/CD"),
    ]
    for i, row in enumerate(module_data):
        fill_table_row(tbl, i + 1, row, font_size=10)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 5: Transaction State Machine
# ---------------------------------------------------------------------------

def slide_05(prs):
    slide = blank_slide(prs)
    tag(slide, "Transactions: State Machine")
    slide_heading(slide, "Eight states. Twelve transitions. Full audit trail.")
    divider(slide, Inches(0.88))

    # State machine text art (left half)
    sm_text = (
        "          [ PENDING ]\n"
        "    /         |         \\\n"
        " accept    decline    cancel\n"
        "   |           |          |\n"
        "[ACCEPTED] [DECLINED] [CANCELLED]\n"
        "   |\n"
        " ship\n"
        "   |\n"
        "[IN_DELIVERY]\n"
        "  /         \\\n"
        "complete   dispute\n"
        "  |              |\n"
        "[COMPLETED] [DISPUTED]\n"
        "               /      \\\n"
        "           resolve   reject\n"
        "             |            |\n"
        "        [COMPLETED] [CANCELLED]"
    )
    add_rounded_rect(slide, Inches(0.4), Inches(1.05), Inches(6.3), Inches(5.3),
                     fill_color=GREY_DIM, line_color=INDIGO)
    add_textbox(slide, sm_text,
                Inches(0.5), Inches(1.1), Inches(6.2), Inches(5.2),
                font_size=11, color=NEARWHITE, align=PP_ALIGN.LEFT, wrap=False)

    # Annotation
    add_textbox(slide,
                "14 days: reminder  |  21 days: auto-complete (cron)",
                Inches(0.4), Inches(6.45), Inches(6.3), Inches(0.4),
                font_size=10, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

    # Bullets right side
    bullets = [
        "Multiple requesters queued chronologically. One accepted, rest waitlisted.",
        "Every transition writes one row to TransactionEvent. Nothing deleted.",
        "Pure state machine logic: 13 unit tests, zero database dependencies.",
    ]
    by = Inches(1.6)
    for b in bullets:
        add_rounded_rect(slide, Inches(7.1), by, Inches(5.9), Inches(1.1),
                         fill_color=GREY_DIM, line_color=INDIGO)
        add_textbox(slide, b,
                    Inches(7.25), by + Inches(0.08), Inches(5.65), Inches(0.95),
                    font_size=12, color=NEARWHITE, wrap=True)
        by += Inches(1.25)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 6: Reputation Tiers
# ---------------------------------------------------------------------------

def slide_06(prs):
    slide = blank_slide(prs)
    tag(slide, "Trust: Reputation Tiers")
    slide_heading(slide, "Four tiers. Earned, not assigned.")
    divider(slide, Inches(0.88))

    # Tier bar
    tiers = [
        ("0 to 19\nNew Member",         RGBColor(0x64, 0x74, 0x8B)),
        ("20 to 49\nActive Sharer",      RGBColor(0x38, 0x8E, 0xFF)),
        ("50 to 79\nTrusted Contributor",GREEN),
        ("80 to 100\nCommunity Champion",AMBER),
    ]
    tx = Inches(0.4)
    tw = Inches(3.1)
    for label, color in tiers:
        add_rect(slide, tx, Inches(1.1), tw, Inches(1.0), fill_color=color)
        add_textbox(slide, label,
                    tx, Inches(1.1), tw, Inches(1.0),
                    font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
        tx += tw + Inches(0.05)

    # Score change table
    headers = ["Event", "Score Change"]
    col_widths = [Inches(6.5), Inches(2.5)]
    tbl = add_table(slide, 7, 2,
                    Inches(0.4), Inches(2.3),
                    Inches(9.0), Inches(3.6),
                    headers=headers, col_widths=col_widths)

    score_data = [
        ("Transaction completed",          "+10 both parties"),
        ("5-star rating received",         "+5"),
        ("1-star rating received",         "-3"),
        ("Report upheld against user",     "-15"),
        ("Cancellation (cancelling party)","-3"),
        ("30 days inactive (cron)",        "-1"),
    ]
    for i, row in enumerate(score_data):
        fill_table_row(tbl, i + 1, row, font_size=11)
        # colour delta
        delta_cell = tbl.cell(i + 1, 1)
        val = row[1]
        clr = GREEN if val.startswith("+") else RGBColor(0xEF, 0x44, 0x44)
        for run in delta_cell.text_frame.paragraphs[0].runs:
            run.font.color.rgb = clr
            run.font.bold = True

    # Bottom note
    add_textbox(slide,
                "Anti-gaming: reciprocal-only pairs and zero-unique-counterparty accounts flagged automatically. Human moderator reviews every flag.",
                Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.6),
                font_size=11, color=AMBER, italic=True, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 7: Price Cap
# ---------------------------------------------------------------------------

def slide_07(prs):
    slide = blank_slide(prs)
    tag(slide, "Solution: Price Cap")
    slide_heading(slide, "50,000 VND cap. Three enforcement layers.")
    divider(slide, Inches(0.88))

    layers = [
        ("Browser form",
         "Client disables submit if price exceeds cap"),
        ("API route: POST /api/listings",
         "Zod schema rejects price above env var"),
        ("Service layer: listings/service.ts",
         "Second check before Prisma write"),
        ("Database",
         "askingPriceVnd stored for SELL type only"),
        ("Grant reporting",
         "Admin CSV shows full transaction price breakdown"),
    ]

    y = Inches(1.1)
    bh = Inches(0.72)
    gap = Inches(0.18)

    for i, (label, desc) in enumerate(layers):
        col = INDIGO if i % 2 == 0 else GREEN
        add_rounded_rect(slide, Inches(1.5), y, Inches(10.0), bh,
                         fill_color=GREY_DIM, line_color=col)
        full = f"{label}  —  {desc}"
        add_textbox(slide, full, Inches(1.65), y, Inches(9.8), bh,
                    font_size=13, color=WHITE, bold=(i % 2 == 0))
        if i < len(layers) - 1:
            arrow_text(slide, Inches(6.0), y + bh, Inches(1.5), gap + Inches(0.05))
        y += bh + gap + Inches(0.06)

    # Two business reasons
    reasons = [
        "Keeps the platform eligible for educational grants (non-commercial mission).",
        "Configurable via SALE_PRICE_CAP_VND environment variable. No code change needed.",
    ]
    bx = Inches(0.5)
    for reason in reasons:
        add_rounded_rect(slide, bx, Inches(6.1), Inches(6.0), Inches(0.65),
                         fill_color=GREY_DIM, line_color=AMBER)
        add_textbox(slide, reason, bx + Inches(0.1), Inches(6.1),
                    Inches(5.8), Inches(0.65),
                    font_size=11, color=NEARWHITE, italic=True)
        bx += Inches(6.4)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 8: Real-time Without Extra Services
# ---------------------------------------------------------------------------

def slide_08(prs):
    slide = blank_slide(prs)
    tag(slide, "Architecture: Real-time")
    slide_heading(slide, "Live notifications and chat. No WebSocket server.")
    divider(slide, Inches(0.88))

    # Left column: Notifications SSE
    add_textbox(slide, "Notifications SSE",
                Inches(0.4), Inches(1.05), Inches(6.0), Inches(0.45),
                font_size=14, bold=True, color=INDIGO)

    notif_steps = [
        "Client opens /api/notifications/stream",
        "Server pushes new Notification rows as they arrive",
        "Fallback: long-polling via GET /api/notifications?wait=1&after=cursor",
    ]
    ny = Inches(1.55)
    for step in notif_steps:
        add_rounded_rect(slide, Inches(0.4), ny, Inches(5.9), Inches(0.65),
                         fill_color=GREY_DIM, line_color=INDIGO)
        add_textbox(slide, step, Inches(0.5), ny, Inches(5.8), Inches(0.65),
                    font_size=11, color=NEARWHITE)
        ny += Inches(0.75)

    # Right column: Messaging SSE
    add_textbox(slide, "Messaging SSE",
                Inches(7.0), Inches(1.05), Inches(6.0), Inches(0.45),
                font_size=14, bold=True, color=GREEN)

    msg_steps = [
        "Client opens /api/conversations/[id]/stream",
        "Server pushes new Message rows as they arrive",
        "User sends message via regular POST request",
    ]
    my = Inches(1.55)
    for step in msg_steps:
        add_rounded_rect(slide, Inches(7.0), my, Inches(5.9), Inches(0.65),
                         fill_color=GREY_DIM, line_color=GREEN)
        add_textbox(slide, step, Inches(7.1), my, Inches(5.8), Inches(0.65),
                    font_size=11, color=NEARWHITE)
        my += Inches(0.75)

    # Central note
    add_rounded_rect(slide, Inches(2.5), Inches(4.0), Inches(8.33), Inches(0.75),
                     fill_color=RGBColor(0x1E, 0x1B, 0x4B), line_color=INDIGO, line_width=Pt(2))
    add_textbox(slide,
                "One-way push via HTTP. Native reconnect. No broker required.",
                Inches(2.6), Inches(4.0), Inches(8.13), Inches(0.75),
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom note
    add_textbox(slide,
                "Vercel free tier: 60-second function timeout, client reconnects automatically. Pro tier: maxDuration=300.",
                Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.5),
                font_size=11, color=AMBER, italic=True, align=PP_ALIGN.CENTER)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 9: Infrastructure
# ---------------------------------------------------------------------------

def slide_09(prs):
    slide = blank_slide(prs)
    tag(slide, "Infrastructure")
    slide_heading(slide, "Standard cloud. No proprietary dependencies.")
    divider(slide, Inches(0.88))

    # Deployment flow (left)
    flow = [
        ("GitHub push to main", INDIGO),
        ("GitHub Actions CI:\nnpm ci, prisma migrate deploy, tsc, vitest, next build", INDIGO),
        ("Vercel:\npages, API routes, SSE streams, cron jobs", GREEN),
        ("Managed PostgreSQL (Supabase / Neon / Vercel Postgres)", AMBER),
        ("S3-compatible bucket (Cloudflare R2 / AWS S3)", AMBER),
    ]
    fy = Inches(1.05)
    fh = Inches(0.8)
    for i, (label, col) in enumerate(flow):
        add_rounded_rect(slide, Inches(0.4), fy, Inches(7.3), fh,
                         fill_color=GREY_DIM, line_color=col)
        add_textbox(slide, label, Inches(0.5), fy, Inches(7.1), fh,
                    font_size=11, color=WHITE)
        if i < len(flow) - 1:
            arrow_text(slide, Inches(3.5), fy + fh, Inches(1.5), Inches(0.28))
        fy += fh + Inches(0.3)

    # Cron jobs table (right)
    headers = ["Endpoint", "Frequency", "Action"]
    col_widths = [Inches(2.8), Inches(1.3), Inches(2.3)]
    tbl = add_table(slide, 5, 3,
                    Inches(8.0), Inches(1.05),
                    Inches(5.0), Inches(3.2),
                    headers=headers, col_widths=col_widths)

    cron_data = [
        ("/api/cron/transactions",           "daily",       "14-day reminder, 21-day auto-complete"),
        ("/api/cron/reputation",             "daily",       "Time decay, anti-gaming scan"),
        ("/api/cron/notification-digest",    "daily",       "Email digest for DAILY preference users"),
        ("/api/cron/notification-immediate", "every 5 min", "Immediate email delivery"),
    ]
    for i, row in enumerate(cron_data):
        fill_table_row(tbl, i + 1, row, font_size=9)

    # Performance grid
    add_textbox(slide, "Performance Targets",
                Inches(8.0), Inches(4.5), Inches(5.0), Inches(0.4),
                font_size=12, bold=True, color=INDIGO)

    perf = [
        ("Page load < 3s on 4G",       "Search p95 < 2s"),
        ("1,000+ concurrent users",     "99.5% uptime"),
    ]
    py = Inches(4.95)
    for row in perf:
        px = Inches(8.0)
        for cell_text in row:
            add_rounded_rect(slide, px, py, Inches(2.35), Inches(0.6),
                             fill_color=GREY_DIM, line_color=GREEN)
            add_textbox(slide, cell_text, px, py, Inches(2.35), Inches(0.6),
                        font_size=11, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
            px += Inches(2.45)
        py += Inches(0.7)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 10: Module — Identity and Profile
# ---------------------------------------------------------------------------

def slide_10(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 1: Identity and Profile")
    slide_heading(slide, "Module 1: Identity and Profile")
    slide_subheading(slide, "Every other module depends on this one.")
    divider(slide, Inches(1.3))

    # Flows
    flows = [
        ("Register flow:",
         "Register  ->  Argon2id hash  ->  Email verification token (72h TTL)  ->  ACTIVE",
         INDIGO),
        ("Login flow:",
         "Login  ->  iron-session cookie (signed, AES-encrypted, httpOnly, sameSite=lax)  ->  30-minute sliding session",
         GREEN),
        ("Password reset flow:",
         "Password reset  ->  Token (1h TTL)  ->  new password",
         AMBER),
    ]
    fy = Inches(1.45)
    for label, text, col in flows:
        add_textbox(slide, label,
                    Inches(0.4), fy, Inches(12.0), Inches(0.3),
                    font_size=11, color=col, bold=True)
        add_rounded_rect(slide, Inches(0.4), fy + Inches(0.3), Inches(12.5), Inches(0.55),
                         fill_color=GREY_DIM, line_color=col)
        add_textbox(slide, text,
                    Inches(0.5), fy + Inches(0.3), Inches(12.3), Inches(0.55),
                    font_size=11, color=NEARWHITE)
        fy += Inches(1.05)

    # Public interface
    add_textbox(slide, "Public Interface",
                Inches(0.4), fy + Inches(0.1), Inches(6.0), Inches(0.4),
                font_size=12, bold=True, color=INDIGO)
    iface = (
        "getCurrentUser(): Promise<User | null>\n"
        "requireUser(): Promise<User>\n"
        "requireRole(...roles): Promise<User>"
    )
    add_rounded_rect(slide, Inches(0.4), fy + Inches(0.5), Inches(6.0), Inches(1.1),
                     fill_color=RGBColor(0x0B, 0x10, 0x1E), line_color=INDIGO)
    add_textbox(slide, iface,
                Inches(0.5), fy + Inches(0.5), Inches(5.9), Inches(1.1),
                font_size=11, color=GREEN, wrap=False)

    # Security bullets
    sec_bullets = [
        ("Argon2id:", "memory-hard, GPU-resistant, PHC 2015 winner"),
        ("Constant-time login:", "dummy hash run even when email not found"),
        ("Rate limit:", "5 attempts per 15 minutes per IP"),
    ]
    bx = Inches(7.0)
    by = fy + Inches(0.1)
    for label, desc in sec_bullets:
        add_rounded_rect(slide, bx, by, Inches(5.9), Inches(0.6),
                         fill_color=GREY_DIM, line_color=GREEN)
        full_text = f"{label}  {desc}"
        add_textbox(slide, full_text, bx + Inches(0.1), by, Inches(5.7), Inches(0.6),
                    font_size=11, color=NEARWHITE)
        by += Inches(0.7)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 11: Module — Book Catalog
# ---------------------------------------------------------------------------

def slide_11(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 2: Book Catalog")
    slide_heading(slide, "Module 2: Book Catalog")
    slide_subheading(slide, "The inventory layer.")
    divider(slide, Inches(1.3))

    # Creation flow
    create_steps = [
        "Create listing (title, author, ISBN, condition, 1-5 photos, type, price)",
        "ISBN auto-fill via Open Library API (fallback: manual entry)",
        "Price cap check (SELL only)",
        "Photos normalised to WebP 1024px",
        "ACTIVE listing",
    ]
    sy = Inches(1.45)
    for i, step in enumerate(create_steps):
        col = INDIGO if i == 0 else (GREEN if i == len(create_steps)-1 else NEARWHITE)
        add_rounded_rect(slide, Inches(0.4), sy, Inches(6.5), Inches(0.55),
                         fill_color=GREY_DIM, line_color=col)
        add_textbox(slide, step, Inches(0.5), sy, Inches(6.3), Inches(0.55),
                    font_size=11, color=WHITE)
        if i < len(create_steps) - 1:
            arrow_text(slide, Inches(3.0), sy + Inches(0.55), Inches(1.0), Inches(0.22))
        sy += Inches(0.55 + 0.22)

    # Status flow
    add_textbox(slide, "Listing status transitions:",
                Inches(7.1), Inches(1.45), Inches(5.9), Inches(0.35),
                font_size=12, bold=True, color=INDIGO)

    status_flow = [
        ("ACTIVE",      INDIGO),
        ("RESERVED",    AMBER),
        ("COMPLETED",   GREEN),
    ]
    sx = Inches(7.1)
    for j, (s, c) in enumerate(status_flow):
        add_rounded_rect(slide, sx, Inches(1.85), Inches(1.7), Inches(0.5),
                         fill_color=GREY_DIM, line_color=c)
        add_textbox(slide, s, sx, Inches(1.85), Inches(1.7), Inches(0.5),
                    font_size=11, bold=True, color=c, align=PP_ALIGN.CENTER)
        if j < len(status_flow) - 1:
            arrow_text(slide, sx + Inches(1.7), Inches(1.85), Inches(0.4), Inches(0.5), ">")
        sx += Inches(2.1)

    # UNAVAILABLE and REMOVED branches
    add_textbox(slide, "ACTIVE  ->  UNAVAILABLE (hidden by owner)",
                Inches(7.1), Inches(2.55), Inches(5.9), Inches(0.4),
                font_size=11, color=NEARWHITE)
    add_textbox(slide, "ACTIVE  ->  REMOVED (deleted or moderated, soft delete)",
                Inches(7.1), Inches(2.95), Inches(5.9), Inches(0.4),
                font_size=11, color=NEARWHITE)

    # Rules
    rules = [
        "Edit blocked when transaction is ACCEPTED or IN_DELIVERY",
        "Pending requesters notified if listing details change while PENDING",
        "Soft delete preserves audit log for transaction history",
    ]
    ry = Inches(4.5)
    for rule in rules:
        add_rounded_rect(slide, Inches(0.4), ry, Inches(12.5), Inches(0.55),
                         fill_color=GREY_DIM, line_color=INDIGO)
        add_textbox(slide, rule, Inches(0.55), ry, Inches(12.3), Inches(0.55),
                    font_size=11, color=NEARWHITE)
        ry += Inches(0.65)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 12: Module — Discovery and Feed
# ---------------------------------------------------------------------------

def slide_12(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 3: Discovery and Feed")
    slide_heading(slide, "Module 3: Discovery and Feed")
    slide_subheading(slide, "How readers find books and follow each other.")
    divider(slide, Inches(1.3))

    # Left: Search
    add_textbox(slide, "Search",
                Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=14, bold=True, color=INDIGO)
    search_items = [
        "Query  ->  Zod param validation  ->  PostgreSQL full-text search (tsvector + GIN index)",
        "ts_rank weights: title A > author B > description C",
        "Filters: genre AND condition AND transactionType AND maxPrice",
        "Cursor-based pagination (not offset)",
        "Benchmark: GIN index drops search from 200ms to 5ms on 1,000 listings",
    ]
    sy = Inches(1.9)
    for item in search_items:
        col = GREEN if "Benchmark" in item else NEARWHITE
        add_rounded_rect(slide, Inches(0.4), sy, Inches(6.1), Inches(0.6),
                         fill_color=GREY_DIM, line_color=INDIGO)
        add_textbox(slide, item, Inches(0.5), sy, Inches(6.0), Inches(0.6),
                    font_size=10, color=col)
        sy += Inches(0.68)

    # Right: Feed
    add_textbox(slide, "Feed",
                Inches(7.0), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=14, bold=True, color=GREEN)
    feed_items = [
        "User A follows User B",
        "User B publishes listing",
        "Fan-out writes FeedItem for User A",
        "User A opens feed  ->  sees new listing in real time",
        "FeedItem index: (userId, createdAt) for fast paginated reads",
    ]
    fy2 = Inches(1.9)
    for item in feed_items:
        add_rounded_rect(slide, Inches(7.0), fy2, Inches(6.0), Inches(0.6),
                         fill_color=GREY_DIM, line_color=GREEN)
        add_textbox(slide, item, Inches(7.1), fy2, Inches(5.9), Inches(0.6),
                    font_size=10, color=NEARWHITE)
        fy2 += Inches(0.68)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 13: Module — Transactions and Messaging
# ---------------------------------------------------------------------------

def slide_13(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 4: Transactions and Messaging")
    slide_heading(slide, "Module 4: Transactions and Messaging")
    slide_subheading(slide, "The most complex module. The core value.")
    divider(slide, Inches(1.3))

    # Left: Transaction flow
    add_textbox(slide, "Transaction flow",
                Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=13, bold=True, color=AMBER)
    txn_steps = [
        "Request  ->  PENDING",
        "Owner accepts one  ->  ACCEPTED, rest  ->  WAITLISTED",
        "Ship  ->  IN_DELIVERY",
        "Confirm  ->  COMPLETED  ->  both rate  ->  ReputationEvent emitted",
        "Dispute  ->  DISPUTED  ->  moderator  ->  COMPLETED or CANCELLED",
    ]
    ty = Inches(1.9)
    for step in txn_steps:
        add_rounded_rect(slide, Inches(0.4), ty, Inches(6.1), Inches(0.62),
                         fill_color=GREY_DIM, line_color=AMBER)
        add_textbox(slide, step, Inches(0.5), ty, Inches(6.0), Inches(0.62),
                    font_size=10, color=NEARWHITE)
        ty += Inches(0.72)

    # Right: Messaging
    add_textbox(slide, "Messaging",
                Inches(7.0), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=13, bold=True, color=GREEN)
    msg_steps = [
        "Transaction accepted  ->  Conversation created",
        "POST message  ->  database write",
        "SSE stream pushes to other client in real time",
        "Message tied to transaction ID for context",
    ]
    my = Inches(1.9)
    for step in msg_steps:
        add_rounded_rect(slide, Inches(7.0), my, Inches(6.0), Inches(0.62),
                         fill_color=GREY_DIM, line_color=GREEN)
        add_textbox(slide, step, Inches(7.1), my, Inches(5.9), Inches(0.62),
                    font_size=10, color=NEARWHITE)
        my += Inches(0.72)

    # Architecture note
    note = "Architecture note: State machine is a pure function. Takes state + action, returns next state or throws. 13 tests, zero DB dependencies."
    add_rounded_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.75),
                     fill_color=RGBColor(0x1E, 0x1B, 0x4B), line_color=INDIGO)
    add_textbox(slide, note, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.75),
                font_size=12, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 14: Module — Trust and Safety
# ---------------------------------------------------------------------------

def slide_14(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 5: Trust and Safety")
    slide_heading(slide, "Module 5: Trust and Safety")
    slide_subheading(slide, "What makes BookBridge different from Facebook Marketplace.")
    divider(slide, Inches(1.3))

    # Left: Reputation
    add_textbox(slide, "Reputation",
                Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=13, bold=True, color=GREEN)
    rep_steps = [
        "Transaction event emitted",
        "ReputationEvent written (atomic with transaction)",
        "User.reputationScore updated (denormalised cache)",
        "Tier recalculated",
        "Notification if tier changed",
    ]
    ry = Inches(1.9)
    for step in rep_steps:
        add_rounded_rect(slide, Inches(0.4), ry, Inches(6.1), Inches(0.58),
                         fill_color=GREY_DIM, line_color=GREEN)
        add_textbox(slide, step, Inches(0.5), ry, Inches(6.0), Inches(0.58),
                    font_size=11, color=NEARWHITE)
        ry += Inches(0.67)

    # Right: Moderation
    add_textbox(slide, "Moderation",
                Inches(7.0), Inches(1.45), Inches(6.0), Inches(0.4),
                font_size=13, bold=True, color=RGBColor(0xEF, 0x44, 0x44))
    mod_steps = [
        "User files report (user / listing / transaction / message)",
        "Report row PENDING",
        "Moderator opens queue (community-scoped or admin global)",
        "Action: WARN, REMOVE_LISTING, SUSPEND_USER, RESTORE,\nRESOLVE_DISPUTE, REJECT_DISPUTE",
        "ModerationAction row written (full audit trail)",
    ]
    my = Inches(1.9)
    for step in mod_steps:
        add_rounded_rect(slide, Inches(7.0), my, Inches(6.0), Inches(0.65),
                         fill_color=GREY_DIM, line_color=RGBColor(0xEF, 0x44, 0x44))
        add_textbox(slide, step, Inches(7.1), my, Inches(5.9), Inches(0.65),
                    font_size=10, color=NEARWHITE)
        my += Inches(0.75)

    # Bottom note
    add_rounded_rect(slide, Inches(0.4), Inches(5.75), Inches(12.5), Inches(0.65),
                     fill_color=GREY_DIM, line_color=AMBER)
    add_textbox(slide,
                "Anti-gaming cron: scan for reciprocal-only pairs and zero-unique-counterparty accounts. Flag only. Human reviews every case.",
                Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.65),
                font_size=11, color=AMBER, italic=True, align=PP_ALIGN.CENTER)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 15: Module — Community and Ops
# ---------------------------------------------------------------------------

def slide_15(prs):
    slide = blank_slide(prs)
    tag(slide, "Module 6: Community and Ops")
    slide_heading(slide, "Module 6: Community and Ops")
    slide_subheading(slide, "The connective tissue. Everything routes through here.")
    divider(slide, Inches(1.3))

    sections = [
        {
            "title": "Communities",
            "color": INDIGO,
            "x": Inches(0.3),
            "items": [
                "Scope: UNIVERSITY / LOCATION / GENRE",
                "Max 20 communities per user (server-enforced)",
                "Listings can be scoped to one community",
            ],
        },
        {
            "title": "Notification Dispatcher",
            "color": GREEN,
            "x": Inches(4.6),
            "items": [
                "Subscribes to events from all 5 other modules",
                "Routes to in-app Notification row + email (based on user preference)",
                "Email preferences: IMMEDIATE (5 min cron), DAILY (digest), OFF",
            ],
        },
        {
            "title": "Admin Dashboard",
            "color": AMBER,
            "x": Inches(8.9),
            "items": [
                "Active users, completed transactions, books circulated",
                "Pending moderation reports",
                "Grant report CSV export (date range filter, OWASP CSV-injection safe)",
                "CI/CD: GitHub Actions on every push, auto-deploy to Vercel on green",
            ],
        },
    ]

    for sec in sections:
        add_textbox(slide, sec["title"],
                    sec["x"], Inches(1.45), Inches(4.1), Inches(0.4),
                    font_size=13, bold=True, color=sec["color"])
        iy = Inches(1.9)
        for item in sec["items"]:
            add_rounded_rect(slide, sec["x"], iy, Inches(4.0), Inches(0.72),
                             fill_color=GREY_DIM, line_color=sec["color"])
            add_textbox(slide, item, sec["x"] + Inches(0.1), iy, Inches(3.9), Inches(0.72),
                        font_size=10, color=NEARWHITE)
            iy += Inches(0.82)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 16: Data Model
# ---------------------------------------------------------------------------

def slide_16(prs):
    slide = blank_slide(prs)
    tag(slide, "Data Model")
    slide_heading(slide, "13 entities. One source of truth per concept.")
    divider(slide, Inches(0.88))

    # Entity relationship as text-art boxes
    # Central: User
    add_rounded_rect(slide, Inches(5.7), Inches(1.1), Inches(2.0), Inches(0.55),
                     fill_color=INDIGO, line_color=WHITE)
    add_textbox(slide, "User", Inches(5.7), Inches(1.1), Inches(2.0), Inches(0.55),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Connected entities
    connected = [
        ("Session",             Inches(0.3),  Inches(0.9)),
        ("Listing",             Inches(0.3),  Inches(2.0)),
        ("Follow",              Inches(0.3),  Inches(3.1)),
        ("FeedItem",            Inches(0.3),  Inches(4.2)),
        ("ReputationEvent",     Inches(3.0),  Inches(4.8)),
        ("Report",              Inches(8.8),  Inches(0.9)),
        ("CommunityMembership", Inches(8.8),  Inches(2.0)),
        ("Notification",        Inches(8.8),  Inches(3.1)),
    ]
    for ename, ex, ey in connected:
        add_rounded_rect(slide, ex, ey, Inches(2.2), Inches(0.5),
                         fill_color=GREY_DIM, line_color=INDIGO)
        add_textbox(slide, ename, ex, ey, Inches(2.2), Inches(0.5),
                    font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    # Listing sub-entities
    listing_subs = [
        ("ListingPhoto",    Inches(0.3),  Inches(5.4)),
        ("Transaction",     Inches(2.9),  Inches(5.4)),
    ]
    for ename, ex, ey in listing_subs:
        add_rounded_rect(slide, ex, ey, Inches(2.2), Inches(0.5),
                         fill_color=GREY_DIM, line_color=AMBER)
        add_textbox(slide, ename, ex, ey, Inches(2.2), Inches(0.5),
                    font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    # Transaction sub-entities
    txn_subs = [
        ("TransactionEvent", Inches(5.5), Inches(5.4)),
        ("Rating",           Inches(7.9), Inches(5.4)),
        ("Conversation",     Inches(10.3),Inches(5.4)),
    ]
    for ename, ex, ey in txn_subs:
        add_rounded_rect(slide, ex, ey, Inches(2.2), Inches(0.5),
                         fill_color=GREY_DIM, line_color=GREEN)
        add_textbox(slide, ename, ex, ey, Inches(2.2), Inches(0.5),
                    font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    # Message and ModerationAction
    add_rounded_rect(slide, Inches(10.3), Inches(4.2), Inches(2.2), Inches(0.5),
                     fill_color=GREY_DIM, line_color=GREEN)
    add_textbox(slide, "Message", Inches(10.3), Inches(4.2), Inches(2.2), Inches(0.5),
                font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(8.8), Inches(4.2), Inches(2.2), Inches(0.5),
                     fill_color=GREY_DIM, line_color=RGBColor(0xEF, 0x44, 0x44))
    add_textbox(slide, "ModerationAction", Inches(8.8), Inches(4.2), Inches(2.2), Inches(0.5),
                font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(8.8), Inches(5.4), Inches(2.2), Inches(0.5),
                     fill_color=GREY_DIM, line_color=INDIGO)
    add_textbox(slide, "Community", Inches(8.8), Inches(5.4), Inches(2.2), Inches(0.5),
                font_size=10, color=NEARWHITE, align=PP_ALIGN.CENTER)

    # Design decisions
    decisions = [
        "User.reputationScore is a denormalised cache. ReputationEvent is the source of truth.",
        "TransactionEvent is append-only. Every state transition preserved permanently.",
        "FeedItem.payload and Notification.payload are JSON. New event types need no migration.",
        "onDelete: Cascade for owner-scoped data. onDelete: SetNull for cross-module soft links.",
    ]
    dy = Inches(1.1)
    for dec in decisions:
        add_textbox(slide, "  •  " + dec,
                    Inches(5.5), dy, Inches(7.6), Inches(0.45),
                    font_size=10, color=NEARWHITE)
        dy += Inches(0.5)

    return slide


# ---------------------------------------------------------------------------
# SLIDE 17: Closing
# ---------------------------------------------------------------------------

def slide_17(prs):
    slide = blank_slide(prs)
    tag(slide, "Closing")

    # Big title
    add_textbox(slide, "BookBridge",
                Inches(1.2), Inches(0.7), Inches(10.9), Inches(1.0),
                font_size=56, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # Three cards
    cards = [
        ("Trust at every step",
         "Reputation tiers, anti-gaming, moderation queue, full audit trail",
         INDIGO),
        ("Non-commercial by design",
         "50,000 VND cap, no payment processing, grant-eligible",
         GREEN),
        ("Standard infrastructure",
         "Next.js + PostgreSQL + Vercel, deployable anywhere Node 20 runs",
         AMBER),
    ]
    cx = Inches(0.5)
    for title, desc, col in cards:
        add_rounded_rect(slide, cx, Inches(2.0), Inches(4.0), Inches(3.5),
                         fill_color=GREY_DIM, line_color=col, line_width=Pt(2.5))
        add_textbox(slide, title,
                    cx + Inches(0.1), Inches(2.4), Inches(3.8), Inches(0.6),
                    font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    cx + Inches(0.15), Inches(3.1), Inches(3.7), Inches(1.8),
                    font_size=12, color=NEARWHITE, align=PP_ALIGN.CENTER, wrap=True)
        cx += Inches(4.3)

    # Final line
    add_textbox(slide,
                "Aligned with UN SDG 4 (Quality Education) and SDG 12 (Responsible Consumption)",
                Inches(1.2), Inches(5.9), Inches(10.9), Inches(0.6),
                font_size=13, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

    # Bottom bar
    rect = slide.shapes.add_shape(1,
        Inches(0), Inches(7.1), Inches(13.33), Inches(0.4))
    rect.fill.solid()
    rect.fill.fore_color.rgb = INDIGO
    rect.line.fill.background()

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    builders = [
        slide_01, slide_02, slide_03, slide_04, slide_05,
        slide_06, slide_07, slide_08, slide_09, slide_10,
        slide_11, slide_12, slide_13, slide_14, slide_15,
        slide_16, slide_17,
    ]

    for fn in builders:
        fn(prs)

    out = r"C:\Users\ADMIN\Desktop\SE\SE_book_bridge\docs\BookBridge_Pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x0F, 0x0F, 0x0F)
ACCENT = RGBColor(0x1A, 0x56, 0xDB)   # strong blue
ACCENT2 = RGBColor(0xE8, 0xF0, 0xFE)  # light blue tint
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x1F, 0x29, 0x37)
LIGHT_LINE = RGBColor(0xE5, 0xE7, 0xEB)

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def textbox(slide, text, l, t, w, h, font_size=Pt(12), bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_text_para(tf, text, font_size=Pt(11), bold=False, color=BLACK,
                  align=PP_ALIGN.LEFT, space_before=Pt(0), italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def slide_header(slide, title, subtitle=None):
    # top accent bar
    rect(slide, 0, 0, 13.33, 0.08, ACCENT)
    # title
    textbox(slide, title, 0.6, 0.25, 12, 0.55,
            font_size=Pt(26), bold=True, color=DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, subtitle, 0.6, 0.78, 12, 0.35,
                font_size=Pt(13), color=GRAY, align=PP_ALIGN.LEFT, italic=True)
    # thin divider
    rect(slide, 0.6, 1.05, 11.73, 0.025, ACCENT)

def footer(slide, page_num, total=12):
    rect(slide, 0, 7.2, 13.33, 0.3, DARK)
    textbox(slide, "BookBridge  ·  Confidential Business Overview", 0.4, 7.22, 9, 0.26,
            font_size=Pt(8), color=RGBColor(0xB0,0xB8,0xC8), align=PP_ALIGN.LEFT)
    textbox(slide, f"{page_num} / {total}", 11.8, 7.22, 1.2, 0.26,
            font_size=Pt(8), color=RGBColor(0xB0,0xB8,0xC8), align=PP_ALIGN.RIGHT)

def icon_card(slide, icon, label, desc, l, t, w=2.8, h=1.5):
    rect(slide, l, t, w, h, ACCENT2, LIGHT_LINE, Pt(1))
    textbox(slide, icon, l+0.12, t+0.1, 0.45, 0.45, font_size=Pt(20), color=ACCENT)
    textbox(slide, label, l+0.6, t+0.12, w-0.7, 0.3,
            font_size=Pt(11), bold=True, color=DARK)
    textbox(slide, desc, l+0.12, t+0.5, w-0.2, 0.9,
            font_size=Pt(9.5), color=GRAY, wrap=True)

# ─────────────────────────────────────────────
# SLIDE 1 — Cover
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, 13.33, 3.5, ACCENT)
# decorative circle top-right
shape = s.shapes.add_shape(9, Inches(10.5), Inches(-1.2), Inches(4), Inches(4))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x2D, 0x6A, 0xDF)
shape.line.fill.background()

textbox(s, "📚", 0.7, 0.55, 1, 0.9, font_size=Pt(48))
textbox(s, "BookBridge", 1.7, 0.6, 10, 1.0,
        font_size=Pt(48), bold=True, color=WHITE)
textbox(s, "Share books. Build trust. Keep stories moving.", 1.7, 1.65, 9, 0.5,
        font_size=Pt(17), color=RGBColor(0xBF,0xD7,0xFF), italic=True)
textbox(s, "Business Overview  ·  Team Zootopia", 0.7, 2.75, 9, 0.4,
        font_size=Pt(11), color=RGBColor(0xBF,0xD7,0xFF))

rect(s, 0, 3.5, 13.33, 0.04, ACCENT2)

textbox(s, "A community-driven platform for second-hand book sharing among students and local readers.\nBuilt on trust, reputation, and real connections — not transaction fees.",
        0.7, 3.75, 11.5, 0.9, font_size=Pt(13), color=DARK, wrap=True)

# KPI strip
kpis = [("4", "Transaction Types"), ("8+", "Core Features"), ("5", "User Roles"), ("∞", "Stories Shared")]
for i, (num, lbl) in enumerate(kpis):
    x = 0.7 + i * 3.0
    rect(s, x, 4.9, 2.7, 1.15, ACCENT2, LIGHT_LINE, Pt(1))
    textbox(s, num, x+0.15, 4.95, 2.4, 0.55,
            font_size=Pt(30), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    textbox(s, lbl, x+0.1, 5.48, 2.5, 0.4,
            font_size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER)

footer(s, 1)

# ─────────────────────────────────────────────
# SLIDE 2 — Problem & Opportunity
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "The Problem We Solve", "Why BookBridge exists")

cols = [
    ("📦", "Books Go to Waste",
     "Millions of textbooks sit unused after each semester. Students have no easy, trusted channel to pass them on."),
    ("💸", "High Cost of Learning",
     "New textbooks are prohibitively expensive. Students need affordable access — gift, exchange, or low-cost resale."),
    ("🤝", "No Trust Infrastructure",
     "Generic marketplaces lack the reputation layer needed for peer-to-peer book sharing within academic communities."),
    ("🏘️", "Community Disconnect",
     "University and neighborhood reading communities have no dedicated space to share, discuss, and discover books together."),
]
for i, (icon, title, body) in enumerate(cols):
    x = 0.5 + i * 3.1
    rect(s, x, 1.25, 2.9, 2.5, ACCENT2, LIGHT_LINE, Pt(1))
    textbox(s, icon, x+0.2, 1.35, 0.55, 0.55, font_size=Pt(22), color=ACCENT)
    textbox(s, title, x+0.15, 1.92, 2.6, 0.38, font_size=Pt(11), bold=True, color=DARK)
    textbox(s, body, x+0.15, 2.32, 2.65, 1.3, font_size=Pt(9.5), color=GRAY, wrap=True)

rect(s, 0.5, 4.1, 12.3, 0.9, RGBColor(0xEF,0xF6,0xFF), ACCENT, Pt(1.5))
textbox(s, "💡  BookBridge is the answer: a trusted social network for book circulation — scoped to universities, locations, and reading communities.",
        0.75, 4.18, 11.8, 0.7, font_size=Pt(11.5), color=DARK, bold=True, wrap=True)

footer(s, 2)

# ─────────────────────────────────────────────
# SLIDE 3 — Product Overview
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Product Overview", "What BookBridge is — at a glance")

features = [
    ("📖", "Book Listings", "Gift · Exchange · Sell with condition ratings, photos, and community scoping"),
    ("🔍", "Discovery", "Full-text search by title, author, ISBN + personalized feed via social follows"),
    ("🔄", "Transactions", "End-to-end state machine: request → accept → deliver → complete → rate"),
    ("💬", "Messaging", "Real-time 1-to-1 chat, optionally tied to a transaction thread"),
    ("⭐", "Reputation", "Atomic score events — positive for completions & contributions, negative for disputes"),
    ("🏘️", "Communities", "University, location, and genre sub-communities with posts, likes, and leaderboards"),
    ("🔔", "Notifications", "In-app & email alerts for listings, transactions, messages, and community activity"),
    ("🛡️", "Moderation", "Report queue with structured actions: warn, remove, suspend, resolve dispute"),
]
for i, (icon, title, body) in enumerate(features):
    row, col = divmod(i, 4)
    x = 0.4 + col * 3.22
    y = 1.25 + row * 2.1
    rect(s, x, y, 3.05, 1.85, WHITE, LIGHT_LINE, Pt(1))
    rect(s, x, y, 3.05, 0.38, ACCENT2)
    textbox(s, f"{icon}  {title}", x+0.12, y+0.06, 2.8, 0.3, font_size=Pt(10.5), bold=True, color=DARK)
    textbox(s, body, x+0.12, y+0.44, 2.82, 1.3, font_size=Pt(9), color=GRAY, wrap=True)

footer(s, 3)

# ─────────────────────────────────────────────
# SLIDE 4 — Core Transaction Workflow
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Core Transaction Workflow", "How a book moves from one reader to the next")

steps = [
    ("1", "List", "Seller creates a listing: type, condition, photos, price (if selling)"),
    ("2", "Discover", "Buyer finds via search, feed, or community browse"),
    ("3", "Request", "Buyer sends transaction request → status: Pending"),
    ("4", "Accept", "Seller accepts (or waitlists/declines) → listing Reserved"),
    ("5", "Chat", "Auto-opened conversation thread for coordination"),
    ("6", "Deliver", "In-person or postal; seller marks shipped + tracking"),
    ("7", "Complete", "Buyer confirms receipt → status: Completed"),
    ("8", "Rate", "Mutual star ratings update both parties' reputation scores"),
]

arrow_y = 2.8
for i, (num, title, body) in enumerate(steps):
    x = 0.35 + i * 1.585
    rect(s, x, 1.2, 1.38, 0.52, ACCENT)
    textbox(s, f"{num}. {title}", x+0.07, 1.25, 1.25, 0.42,
            font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 7:
        textbox(s, "▶", x+1.38, 1.35, 0.2, 0.3, font_size=Pt(10), color=ACCENT)
    rect(s, x, 1.85, 1.38, 1.7, ACCENT2, LIGHT_LINE, Pt(0.75))
    textbox(s, body, x+0.07, 1.9, 1.26, 1.6, font_size=Pt(8.5), color=DARK, wrap=True)

rect(s, 0.35, 3.75, 12.5, 0.04, LIGHT_LINE)

# dispute note
rect(s, 0.35, 3.95, 5.9, 0.9, RGBColor(0xFF,0xF7,0xED), RGBColor(0xFB,0xBF,0x24), Pt(1))
textbox(s, "⚠️  Dispute Flow", 0.55, 3.98, 3, 0.3, font_size=Pt(10), bold=True, color=RGBColor(0x92,0x4E,0x00))
textbox(s, "Either party can raise a dispute at any point. A moderator reviews and resolves or rejects via a structured ModerationAction.",
        0.55, 4.28, 5.6, 0.5, font_size=Pt(9), color=DARK, wrap=True)

rect(s, 6.6, 3.95, 6.1, 0.9, RGBColor(0xF0,0xFD,0xF4), RGBColor(0x34,0xD3,0x99), Pt(1))
textbox(s, "✅  Reputation Impact", 6.8, 3.98, 4, 0.3, font_size=Pt(10), bold=True, color=RGBColor(0x06,0x5F,0x46))
textbox(s, "Completion awards positive reputation delta to both buyer and seller. Disputes and cancellations carry negative deltas.",
        6.8, 4.28, 5.8, 0.5, font_size=Pt(9), color=DARK, wrap=True)

footer(s, 4)

# ─────────────────────────────────────────────
# SLIDE 5 — Trust & Reputation System
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Trust & Reputation System", "The engine that makes peer-to-peer sharing safe")

# left column — how it works
rect(s, 0.5, 1.2, 5.8, 4.7, WHITE, LIGHT_LINE, Pt(1))
textbox(s, "How It Works", 0.75, 1.28, 5, 0.35, font_size=Pt(12), bold=True, color=DARK)
events = [
    ("+", "Transaction completed", "ACCENT"),
    ("+", "Rating received (1–5 stars)", "ACCENT"),
    ("+", "Community contribution (mod-awarded)", "ACCENT"),
    ("−", "Dispute upheld against user", "RED"),
    ("−", "Cancellation penalty", "RED"),
    ("−", "Periodic time decay", "RED"),
    ("−", "Report upheld", "RED"),
]
for i, (sign, label, c) in enumerate(events):
    y = 1.72 + i * 0.5
    clr = ACCENT if c == "ACCENT" else RGBColor(0xDC,0x26,0x26)
    textbox(s, sign, 0.75, y, 0.35, 0.38, font_size=Pt(14), bold=True, color=clr)
    textbox(s, label, 1.15, y+0.04, 4.9, 0.35, font_size=Pt(10), color=DARK)

# right column — tiers
rect(s, 6.8, 1.2, 5.95, 4.7, WHITE, LIGHT_LINE, Pt(1))
textbox(s, "Reputation Tiers", 7.05, 1.28, 5, 0.35, font_size=Pt(12), bold=True, color=DARK)
tiers = [
    ("🌱", "New", "Starting tier — building first impressions"),
    ("🔵", "Trusted", "Consistent completions, positive ratings"),
    ("⭐", "Reliable", "Strong track record, low dispute rate"),
    ("🏆", "Veteran", "Top tier — community pillar, unlocks trust signals"),
]
for i, (icon, tier, desc) in enumerate(tiers):
    y = 1.78 + i * 1.0
    rect(s, 7.05, y, 5.45, 0.82, ACCENT2, LIGHT_LINE, Pt(0.75))
    textbox(s, icon, 7.15, y+0.12, 0.45, 0.5, font_size=Pt(18))
    textbox(s, tier, 7.65, y+0.1, 1.5, 0.3, font_size=Pt(11), bold=True, color=DARK)
    textbox(s, desc, 7.65, y+0.4, 4.7, 0.35, font_size=Pt(9), color=GRAY, wrap=True)

footer(s, 5)

# ─────────────────────────────────────────────
# SLIDE 6 — Communities Feature
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Communities", "Scoped sub-networks for universities, locations, and genres")

# three community types
types = [
    ("🎓", "University", "HUST, UET, RMIT…\nClosed academic circles for textbook sharing within a campus"),
    ("📍", "Location", "District, city, or neighbourhood groups — perfect for local book swaps"),
    ("📚", "Genre / Interest", "Non-fiction, Sci-fi, Vietnamese Literature — passion-driven reading circles"),
]
for i, (icon, t, d) in enumerate(types):
    x = 0.5 + i * 4.2
    rect(s, x, 1.2, 3.9, 1.8, ACCENT2, LIGHT_LINE, Pt(1))
    textbox(s, icon, x+0.2, 1.3, 0.6, 0.6, font_size=Pt(26), color=ACCENT)
    textbox(s, t, x+0.9, 1.35, 2.8, 0.38, font_size=Pt(13), bold=True, color=DARK)
    textbox(s, d, x+0.2, 1.75, 3.55, 1.1, font_size=Pt(9.5), color=GRAY, wrap=True)

# features grid
feats = [
    ("Public & Private modes", "Private communities require an invite code — perfect for closed reading clubs"),
    ("Community Posts", "Members post discussions; posts can be liked, commented on, and pinned by moderators"),
    ("Community Listings", "Listings scoped to a community — only visible to members"),
    ("Points Leaderboard", "Per-member communityPoints earned through likes, comments, and activity"),
    ("Moderator Hierarchy", "Owner → Moderator → Member with promote/demote controls"),
    ("Notifications", "All members notified of new posts; post authors alerted on likes and comments"),
]
for i, (title, desc) in enumerate(feats):
    row, col = divmod(i, 3)
    x = 0.5 + col * 4.2
    y = 3.3 + row * 1.35
    rect(s, x, y, 3.9, 1.2, WHITE, LIGHT_LINE, Pt(1))
    textbox(s, title, x+0.15, y+0.1, 3.6, 0.32, font_size=Pt(10), bold=True, color=ACCENT)
    textbox(s, desc, x+0.15, y+0.44, 3.65, 0.68, font_size=Pt(9), color=GRAY, wrap=True)

footer(s, 6)

# ─────────────────────────────────────────────
# SLIDE 7 — Discovery & Social Feed
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Discovery & Social Feed", "How readers find the right books")

left_items = [
    ("🔍 Full-Text Search",
     "Search by title, author, or ISBN. Filter by genre, condition, and transaction type (Gift / Exchange / Sell)."),
    ("🌐 Explore Page",
     "Public browse mode — no login required. Great for acquisition and first impressions."),
    ("📰 Personalized Feed",
     "Follow other users. Every new listing they publish is fan-fanned into your feed in real time via Server-Sent Events."),
    ("🏘️ Community Browse",
     "Within a community, members see listings and posts scoped only to that group — reducing noise, increasing relevance."),
]
for i, (title, body) in enumerate(left_items):
    y = 1.25 + i * 1.35
    rect(s, 0.5, y, 6.2, 1.2, WHITE, LIGHT_LINE, Pt(1))
    rect(s, 0.5, y, 0.06, 1.2, ACCENT)
    textbox(s, title, 0.72, y+0.1, 5.8, 0.32, font_size=Pt(10.5), bold=True, color=DARK)
    textbox(s, body, 0.72, y+0.44, 5.85, 0.7, font_size=Pt(9.5), color=GRAY, wrap=True)

# right side — social graph
rect(s, 7.2, 1.2, 5.5, 5.2, ACCENT2, LIGHT_LINE, Pt(1))
textbox(s, "Social Graph", 7.45, 1.28, 5, 0.35, font_size=Pt(12), bold=True, color=DARK)
social = [
    "Follow / unfollow any user",
    "Follower + following counts on profile",
    "Feed populated on publish (fan-out write)",
    "Real-time delivery via SSE",
    "New listing notifications to followers",
    "Profiles show reputation tier, bio, preferred genres, and district",
]
for i, item in enumerate(social):
    y = 1.78 + i * 0.65
    textbox(s, "•  " + item, 7.45, y, 5.1, 0.55, font_size=Pt(10), color=DARK, wrap=True)

footer(s, 7)

# ─────────────────────────────────────────────
# SLIDE 8 — Notifications & Messaging
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Notifications & Messaging", "Keeping every participant in the loop")

notif_kinds = [
    ("📦", "New listing from a followed user"),
    ("🔄", "Transaction status changed"),
    ("💬", "New direct message received"),
    ("📢", "Community announcement"),
    ("⭐", "Reputation tier changed"),
    ("🛡️", "Moderation action taken"),
    ("📝", "New post in your community"),
    ("❤️", "Someone liked your post"),
    ("🗨️", "New comment on your post"),
]
textbox(s, "9 Notification Kinds", 0.5, 1.18, 6.5, 0.38, font_size=Pt(11), bold=True, color=DARK)
for i, (icon, label) in enumerate(notif_kinds):
    row, col = divmod(i, 3)
    x = 0.5 + col * 2.1
    y = 1.65 + row * 0.58
    rect(s, x, y, 1.97, 0.48, WHITE, LIGHT_LINE, Pt(0.75))
    textbox(s, f"{icon} {label}", x+0.1, y+0.06, 1.8, 0.38, font_size=Pt(8.5), color=DARK, wrap=True)

rect(s, 0.5, 3.45, 6.3, 0.04, LIGHT_LINE)

textbox(s, "Email Preferences", 0.5, 3.6, 6, 0.32, font_size=Pt(11), bold=True, color=DARK)
for label, desc, x in [
    ("Immediate", "Each event triggers an email", 0.5),
    ("Daily Digest", "Bundled summary once per day", 2.7),
    ("Off", "In-app only", 4.9),
]:
    rect(s, x, 4.0, 1.95, 0.8, ACCENT2, LIGHT_LINE, Pt(0.75))
    textbox(s, label, x+0.1, 4.06, 1.78, 0.28, font_size=Pt(9.5), bold=True, color=ACCENT)
    textbox(s, desc, x+0.1, 4.32, 1.78, 0.38, font_size=Pt(8.5), color=GRAY)

# messaging panel
rect(s, 7.2, 1.2, 5.5, 5.0, WHITE, LIGHT_LINE, Pt(1))
textbox(s, "💬  Real-Time Messaging", 7.45, 1.28, 5, 0.35, font_size=Pt(12), bold=True, color=DARK)
msg_feats = [
    ("1-to-1 chat", "Private conversations between any two users"),
    ("Transaction-linked", "Auto-opened thread when a transaction is accepted — context always available"),
    ("Real-time", "Server-Sent Events for instant delivery without websockets overhead"),
    ("Reportable", "Any message can be reported; lands in the moderation queue"),
    ("Persistent", "Full message history stored; no ephemeral-only conversations"),
]
for i, (title, desc) in enumerate(msg_feats):
    y = 1.78 + i * 0.82
    rect(s, 7.4, y, 5.1, 0.7, ACCENT2)
    textbox(s, title, 7.55, y+0.06, 4.8, 0.25, font_size=Pt(10), bold=True, color=DARK)
    textbox(s, desc, 7.55, y+0.33, 4.8, 0.34, font_size=Pt(9), color=GRAY, wrap=True)

footer(s, 8)

# ─────────────────────────────────────────────
# SLIDE 9 — Moderation & Safety
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Moderation & Platform Safety", "Trust at scale — structured, fair, transparent")

reportable = [("👤 User", 0.5), ("📦 Listing", 2.65), ("🔄 Transaction", 4.8), ("💬 Message", 6.95)]
textbox(s, "What Can Be Reported", 0.5, 1.2, 9, 0.32, font_size=Pt(11), bold=True, color=DARK)
for label, x in reportable:
    rect(s, x, 1.58, 1.92, 0.55, ACCENT2, LIGHT_LINE, Pt(0.75))
    textbox(s, label, x+0.12, 1.65, 1.72, 0.38, font_size=Pt(9.5), bold=True, color=DARK)

actions = [
    ("⚠️ Warn", "Formal warning on record — no immediate restriction"),
    ("🗑️ Remove Listing", "Listing taken down and flagged as removed"),
    ("🔇 Suspend User", "Account suspended; user cannot transact or post"),
    ("✅ Resolve Dispute", "Moderator rules on a disputed transaction"),
    ("❌ Reject Report", "Report found unfounded — no action taken"),
    ("🔄 Restore", "Suspension lifted or listing re-activated after review"),
]
textbox(s, "Moderator Actions", 0.5, 2.35, 9, 0.32, font_size=Pt(11), bold=True, color=DARK)
for i, (action, desc) in enumerate(actions):
    row, col = divmod(i, 3)
    x = 0.5 + col * 4.1
    y = 2.75 + row * 0.92
    rect(s, x, y, 3.85, 0.78, WHITE, LIGHT_LINE, Pt(1))
    textbox(s, action, x+0.12, y+0.08, 3.6, 0.28, font_size=Pt(10), bold=True, color=DARK)
    textbox(s, desc, x+0.12, y+0.38, 3.6, 0.35, font_size=Pt(9), color=GRAY, wrap=True)

rect(s, 0.5, 4.72, 12.3, 0.04, LIGHT_LINE)

# right — auto & admin
rect(s, 0.5, 4.85, 5.85, 1.2, RGBColor(0xFF,0xF7,0xED), RGBColor(0xFB,0xBF,0x24), Pt(1))
textbox(s, "🤖  Auto-Flagging", 0.7, 4.92, 5.5, 0.3, font_size=Pt(10), bold=True, color=RGBColor(0x78,0x35,0x00))
textbox(s, "System can auto-generate reports on suspicious patterns — feeding the queue without requiring manual user reports.",
        0.7, 5.22, 5.65, 0.75, font_size=Pt(9), color=DARK, wrap=True)

rect(s, 6.7, 4.85, 5.85, 1.2, RGBColor(0xF0,0xFD,0xF4), RGBColor(0x34,0xD3,0x99), Pt(1))
textbox(s, "🛡️  Admin Dashboard", 6.9, 4.92, 5.5, 0.3, font_size=Pt(10), bold=True, color=RGBColor(0x06,0x5F,0x46))
textbox(s, "Separate /admin route for system-level oversight: manage users, listings, and moderation across all communities.",
        6.9, 5.22, 5.65, 0.75, font_size=Pt(9), color=DARK, wrap=True)

footer(s, 9)

# ─────────────────────────────────────────────
# SLIDE 10 — User Roles & Access
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "User Roles & Access Hierarchy", "Clear permissions at every level")

roles = [
    ("👁️", "Guest", "No account required",
     ["Browse public listings", "Full-text search", "View public profiles", "Explore page access"], GRAY),
    ("👤", "User", "Verified account",
     ["Create & manage listings", "Request transactions", "Rate & review", "Follow users", "Join communities", "Send messages"], ACCENT),
    ("🔧", "Moderator", "Elevated community role",
     ["All User capabilities", "Access moderation queue", "Take moderation actions", "Award community points", "Pin/unpin posts", "Promote/demote members"], RGBColor(0x09,0x7C,0x4E)),
    ("⚙️", "Admin", "Platform-wide authority",
     ["Full platform access", "Admin dashboard", "System-level moderation", "User management", "Override any action"], RGBColor(0x7C,0x27,0x08)),
]
for i, (icon, role, subtitle, perms, color) in enumerate(roles):
    x = 0.4 + i * 3.15
    rect(s, x, 1.2, 2.95, 5.1, WHITE, LIGHT_LINE, Pt(1))
    rect(s, x, 1.2, 2.95, 0.65, color)
    textbox(s, f"{icon}  {role}", x+0.15, 1.27, 2.7, 0.35, font_size=Pt(13), bold=True, color=WHITE)
    textbox(s, subtitle, x+0.15, 1.62, 2.7, 0.28, font_size=Pt(8.5), color=RGBColor(0xD1,0xD5,0xDB), italic=True)
    for j, perm in enumerate(perms):
        rect(s, x+0.15, 2.0 + j*0.5, 0.08, 0.24, color)
        textbox(s, perm, x+0.32, 1.98 + j*0.5, 2.55, 0.38, font_size=Pt(9), color=DARK, wrap=True)

# account lifecycle note
rect(s, 0.4, 6.45, 12.5, 0.6, ACCENT2, ACCENT, Pt(1))
textbox(s, "Account lifecycle:  Pending Verification  →  Active  →  Suspended / Deleted  ·  Email verification required before account becomes Active.",
        0.65, 6.5, 12.0, 0.5, font_size=Pt(9.5), color=DARK, bold=False, wrap=True)

footer(s, 10)

# ─────────────────────────────────────────────
# SLIDE 11 — Business Model & Market Fit
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
slide_header(s, "Business Model & Market Fit", "How BookBridge creates and captures value")

# positioning
rect(s, 0.5, 1.2, 12.3, 0.72, ACCENT)
textbox(s, "Positioning:  BookBridge is a social trust network for book circulation — not a marketplace.",
        0.75, 1.32, 11.8, 0.5, font_size=Pt(12), bold=True, color=WHITE, wrap=True)

pillars = [
    ("🏫", "Campus-First GTM",
     "Launch within universities — where textbook pain is sharpest. Partner with student unions and library programs for distribution."),
    ("🤝", "Network Effect",
     "Value compounds with users: more givers → more receivers → more ratings → stronger trust → more engagement."),
    ("💰", "Monetization Paths",
     "Premium verified-seller badge · University partnership licensing · Local business listing sponsorships · API for library integrations."),
    ("🔒", "Defensible Moat",
     "Reputation equity. A user's ReputationScore and community standing can't be transferred — switching cost grows with every transaction."),
]
for i, (icon, title, body) in enumerate(pillars):
    row, col = divmod(i, 2)
    x = 0.5 + col * 6.25
    y = 2.15 + row * 1.75
    rect(s, x, y, 5.95, 1.55, WHITE, LIGHT_LINE, Pt(1))
    rect(s, x, y, 0.07, 1.55, ACCENT)
    textbox(s, icon, x+0.2, y+0.14, 0.55, 0.55, font_size=Pt(22))
    textbox(s, title, x+0.85, y+0.14, 4.9, 0.32, font_size=Pt(11), bold=True, color=DARK)
    textbox(s, body, x+0.85, y+0.5, 4.9, 0.95, font_size=Pt(9.5), color=GRAY, wrap=True)

rect(s, 0.5, 5.8, 12.3, 0.04, LIGHT_LINE)
textbox(s, "Price cap enforcement is a deliberate design choice — it keeps the platform non-commercial, builds social goodwill, and differentiates BookBridge from generic resale apps.",
        0.5, 5.95, 12.3, 0.6, font_size=Pt(9.5), color=GRAY, italic=True, wrap=True)

footer(s, 11)

# ─────────────────────────────────────────────
# SLIDE 12 — Summary & Next Steps
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.08, ACCENT)
rect(s, 0, 0.08, 13.33, 3.0, DARK)

textbox(s, "BookBridge", 0.7, 0.35, 10, 0.75, font_size=Pt(40), bold=True, color=WHITE)
textbox(s, "Share books. Build trust. Keep stories moving.", 0.7, 1.15, 10, 0.45,
        font_size=Pt(15), color=RGBColor(0x93,0xC5,0xFD), italic=True)

summary = [
    "Community-driven platform for second-hand book circulation",
    "End-to-end transaction state machine with full audit trail",
    "Reputation system with atomic score events and tiers",
    "Sub-communities scoped by university, location, and genre",
    "Real-time feed, messaging, and 9-kind notification system",
    "Structured moderation queue with admin oversight",
]
for i, item in enumerate(summary):
    x = 0.7 + (i % 2) * 5.9
    y = 1.75 + (i // 2) * 0.38
    textbox(s, "✓  " + item, x, y, 5.7, 0.34, font_size=Pt(9.5), color=RGBColor(0xBF,0xD7,0xFF), wrap=True)

rect(s, 0, 3.08, 13.33, 0.04, ACCENT)

textbox(s, "Next Steps", 0.7, 3.25, 12, 0.38, font_size=Pt(14), bold=True, color=DARK)
nexts = [
    ("🚀", "Production Launch", "Deploy to Vercel with production database migration; smoke-test all critical flows"),
    ("🔔", "Notification Emails", "Wire email dispatcher to SendGrid/Resend for immediate, digest, and off preferences"),
    ("📊", "Analytics", "Add lightweight event tracking (page views, listing clicks, conversion funnel)"),
    ("💳", "Monetization Pilot", "Run university partnership pilot — measure engagement before adding premium features"),
]
for i, (icon, title, body) in enumerate(nexts):
    x = 0.5 + i * 3.15
    rect(s, x, 3.75, 2.95, 2.35, WHITE, LIGHT_LINE, Pt(1))
    rect(s, x, 3.75, 2.95, 0.42, ACCENT2)
    textbox(s, f"{icon}  {title}", x+0.12, 3.8, 2.75, 0.3, font_size=Pt(10), bold=True, color=DARK)
    textbox(s, body, x+0.12, 4.25, 2.75, 1.75, font_size=Pt(9), color=GRAY, wrap=True)

textbox(s, "Built with Next.js 14 · PostgreSQL · Prisma · Vercel  ·  Team Zootopia — 2025 Software Engineering Capstone",
        0.5, 6.45, 12.3, 0.4, font_size=Pt(8.5), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

footer(s, 12)

out = r"C:\Users\ADMIN\Desktop\SE\SE_book_bridge\BookBridge_Business_Overview.pptx"
prs.save(out)
print("Saved:", out)

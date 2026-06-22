from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────
W   = RGBColor(0xFF,0xFF,0xFF)  # white
BG  = RGBColor(0xF8,0xFA,0xFF)  # near-white page bg
B1  = RGBColor(0x1A,0x56,0xDB)  # primary blue
B2  = RGBColor(0x36,0x78,0xF0)  # mid blue (accent)
B3  = RGBColor(0xDB,0xE9,0xFF)  # light blue tint
B4  = RGBColor(0xEF,0xF5,0xFF)  # very light blue
DK  = RGBColor(0x0F,0x17,0x2A)  # near-black
GR  = RGBColor(0x64,0x74,0x8B)  # gray text
LN  = RGBColor(0xCF,0xDB,0xF0)  # border line
GN  = RGBColor(0x05,0x96,0x69)  # green (success)
RD  = RGBColor(0xDC,0x26,0x26)  # red (warning)

blank = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(blank)
    f = s.background.fill; f.solid(); f.fore_color.rgb = W
    return s

def box(s, l,t,w,h, fill, border=None, bw=Pt(1)):
    sh = s.shapes.add_shape(1, Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb = border; sh.line.width = bw
    else: sh.line.fill.background()
    return sh

def tx(s, text, l,t,w,h, size=Pt(11), bold=False, color=DK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size=size; r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color
    return tb

def header(s, title, sub=None):
    box(s, 0,0,13.33,0.06, B1)                          # top rule
    tx(s, title, 0.55,0.18,11,0.58, Pt(28),True,DK)
    if sub:
        tx(s, sub,  0.55,0.74,11,0.3, Pt(11),False,GR,italic=True)
    box(s, 0.55,1.0,11.7,0.018, LN)                      # divider

def footer(s, n, total=12):
    box(s, 0,7.18,13.33,0.32, DK)
    tx(s,"BookBridge  ·  Business Overview",0.45,7.2,9,0.28,Pt(8),color=RGBColor(0x8A,0x9B,0xC0))
    tx(s,f"{n} / {total}",11.9,7.2,1.1,0.28,Pt(8),color=RGBColor(0x8A,0x9B,0xC0),align=PP_ALIGN.RIGHT)

def chip(s, text, l,t, fill=B3, color=B1):
    box(s, l,t,len(text)*0.095+0.25,0.28, fill)
    tx(s, text, l+0.1,t+0.02, len(text)*0.095+0.1,0.24, Pt(8.5),True,color)

def card(s, l,t,w,h, title, body, icon=None, title_color=DK, accent=B1):
    box(s, l,t,w,h, W, LN, Pt(0.75))
    box(s, l,t,0.04,h, accent)                           # left accent stripe
    yoff = 0.12
    if icon:
        tx(s, icon, l+0.15,t+yoff,0.4,0.38, Pt(16))
        tx(s, title, l+0.6,t+yoff,w-0.75,0.32, Pt(10.5),True,title_color)
    else:
        tx(s, title, l+0.2,t+yoff,w-0.3,0.32, Pt(10.5),True,title_color)
    tx(s, body, l+0.2,t+0.48,w-0.3,h-0.55, Pt(9),color=GR,wrap=True)

# ═══════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════
s = slide()
box(s, 0,0,13.33,4.0, B1)
box(s, 0,3.97,13.33,0.06, B2)

# decorative right panel
box(s, 9.8,0,3.53,4.0, B2)
box(s, 10.5,0.4,2.2,3.2, B3)

tx(s,"📚",0.7,0.45,0.9,0.9, Pt(42))
tx(s,"BookBridge",1.6,0.45,8,0.9, Pt(46),True,W)
tx(s,"Share books.  Build trust.  Keep stories moving.",
   1.6,1.45,8.5,0.5, Pt(16),False,RGBColor(0xBE,0xD8,0xFF),italic=True)
tx(s,"A community-driven platform for second-hand book sharing\namong students and local readers — built on reputation, not transaction fees.",
   1.6,2.1,8,0.9, Pt(11.5),False,RGBColor(0xDB,0xEA,0xFF),wrap=True)
tx(s,"Team Zootopia  ·  SE Capstone 2025",1.6,3.3,7,0.38,Pt(9),False,RGBColor(0x93,0xBB,0xFF))

# stat strip
stats = [("3","Transaction\nTypes"),("8+","Core\nFeatures"),("4","User\nRoles"),("∞","Stories\nShared")]
for i,(n,l) in enumerate(stats):
    x = 0.55 + i*2.9
    box(s, x,4.25,2.55,1.6, B4, LN, Pt(0.75))
    tx(s, n, x+0.15,4.35,2.3,0.75, Pt(34),True,B1,align=PP_ALIGN.CENTER)
    tx(s, l, x+0.1,5.08,2.35,0.65, Pt(9),False,GR,align=PP_ALIGN.CENTER,wrap=True)

tx(s,"🟢  Live on Vercel",0.55,6.15,3,0.3,Pt(9),False,GN)
tx(s,"Next.js 14  ·  PostgreSQL  ·  Prisma",0.55,6.5,5,0.3,Pt(8.5),False,GR,italic=True)
footer(s,1)

# ═══════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════
s = slide()
header(s,"The Problem We Solve","Why BookBridge exists")

problems = [
    ("📦","Books go to waste","Millions of textbooks sit unused after each semester. Students have no trusted channel to pass them on."),
    ("💸","Learning costs too much","New textbooks are prohibitively expensive. Students need Gift, Exchange, or low-cost Sell options."),
    ("🤝","No trust layer","Generic marketplaces lack the reputation infrastructure needed for safe peer-to-peer book sharing."),
    ("🏘️","Communities are disconnected","University and neighbourhood reading groups have no dedicated space to share and discover books together."),
]
for i,(ico,title,body) in enumerate(problems):
    x = 0.45 + i*3.12
    box(s, x,1.12,2.92,3.5, W,LN,Pt(0.75))
    box(s, x,1.12,2.92,0.06, B1)
    tx(s, ico,  x+0.2,1.25,0.55,0.55, Pt(22))
    tx(s, title,x+0.15,1.85,2.65,0.38, Pt(11),True,DK)
    tx(s, body, x+0.15,2.3, 2.65,2.1,  Pt(9.5),False,GR,wrap=True)

box(s,0.45,4.82,12.38,1.0, B4,B1,Pt(1))
tx(s,"💡  BookBridge answer",0.7,4.9,4,0.3,Pt(10),True,B1)
tx(s,"A trusted social network for book circulation — scoped to universities, locations, and reading communities. "
    "Value comes from reputation equity and community lock-in, not listing fees.",
    0.7,5.22,11.9,0.52, Pt(10.5),False,DK,wrap=True)
footer(s,2)

# ═══════════════════════════════════════════════
# SLIDE 3 — Product Overview (feature grid)
# ═══════════════════════════════════════════════
s = slide()
header(s,"Product Overview","8 features, one coherent platform")

features = [
    ("📖","Listings","Gift · Exchange · Sell\nCondition, photos, price cap, community scoping"),
    ("🔍","Discovery","Full-text search + personalized\nfeed via social follows"),
    ("🔄","Transactions","Full state machine from request\nto completion with audit trail"),
    ("💬","Messaging","Real-time 1-to-1 chat\nlinked to transaction threads"),
    ("⭐","Reputation","Atomic score events: completions +,\ndisputes −, time decay −"),
    ("🏘️","Communities","University · Location · Genre\nPosts, likes, leaderboards"),
    ("🔔","Notifications","9 kinds: listings, transactions,\nmessages, community activity"),
    ("🛡️","Moderation","Report queue with structured\nactions and admin oversight"),
]
for i,(ico,title,body) in enumerate(features):
    row,col = divmod(i,4)
    x = 0.4 + col*3.22
    y = 1.15 + row*2.3
    box(s, x,y,3.05,2.1, W,LN,Pt(0.75))
    box(s, x,y,3.05,0.42, B4)
    tx(s, f"{ico}  {title}", x+0.14,y+0.07,2.8,0.3, Pt(10.5),True,DK)
    tx(s, body, x+0.14,y+0.5,2.8,1.52, Pt(9.5),False,GR,wrap=True)
footer(s,3)

# ═══════════════════════════════════════════════
# SLIDE 4 — Transaction Workflow
# ═══════════════════════════════════════════════
s = slide()
header(s,"Transaction Workflow","How a book moves from one reader to the next")

steps = [
    ("1","List","Type · Condition\nPhotos · Price"),
    ("2","Discover","Search / Feed\nCommunity browse"),
    ("3","Request","Buyer sends\nrequest → Pending"),
    ("4","Accept","Seller accepts\n→ Reserved"),
    ("5","Chat","Auto thread\nopened"),
    ("6","Deliver","In-person or\npostal + tracking"),
    ("7","Complete","Buyer confirms\nreceipt"),
    ("8","Rate","Mutual stars\nupdate reputation"),
]
for i,(num,title,body) in enumerate(steps):
    x = 0.3 + i*1.595
    # step box
    box(s, x,1.18,1.42,0.48, B1)
    tx(s, f"{num}. {title}", x+0.07,1.23,1.3,0.38, Pt(9.5),True,W,align=PP_ALIGN.CENTER)
    # arrow
    if i<7:
        tx(s,"›", x+1.42,1.3,0.18,0.3, Pt(14),True,B1,align=PP_ALIGN.CENTER)
    # body card
    box(s, x,1.75,1.42,1.6, B4,LN,Pt(0.5))
    tx(s, body, x+0.07,1.82,1.3,1.45, Pt(8.5),False,GR,wrap=True)

box(s,0.3,3.5,12.65,0.025,LN)

# two note boxes
box(s,0.3,3.65,6.15,1.2, RGBColor(0xFF,0xF8,0xED),RGBColor(0xFB,0xBF,0x24),Pt(1))
tx(s,"⚠️  Dispute Flow",0.52,3.72,5.5,0.3,Pt(10),True,RGBColor(0x78,0x35,0x00))
tx(s,"Either party can raise a dispute. A moderator reviews and resolves or rejects via a structured ModerationAction.",
   0.52,4.02,5.9,0.76,Pt(9),False,DK,wrap=True)

box(s,6.78,3.65,6.15,1.2, RGBColor(0xF0,0xFD,0xF4),RGBColor(0x34,0xD3,0x99),Pt(1))
tx(s,"✅  Reputation Impact",7.0,3.72,5.5,0.3,Pt(10),True,RGBColor(0x06,0x5F,0x46))
tx(s,"Completion awards positive delta to both sides. Disputes and cancellations carry negative deltas — keeping quality high.",
   7.0,4.02,5.9,0.76,Pt(9),False,DK,wrap=True)
footer(s,4)

# ═══════════════════════════════════════════════
# SLIDE 5 — Reputation System
# ═══════════════════════════════════════════════
s = slide()
header(s,"Trust & Reputation System","The engine that makes peer-to-peer sharing safe")

# left: score events
box(s,0.45,1.12,6.0,5.45,W,LN,Pt(0.75))
box(s,0.45,1.12,6.0,0.4,B1)
tx(s,"Score Events",0.7,1.18,5,0.28,Pt(11),True,W)

pos_events = [
    ("＋","Transaction completed","Every successful handoff"),
    ("＋","Rating received","1–5 stars from counterparty"),
    ("＋","Community contribution","Awarded by moderator"),
]
neg_events = [
    ("－","Dispute upheld","Against the reported user"),
    ("－","Cancellation","Backing out of accepted deal"),
    ("－","Report upheld","Confirmed bad behaviour"),
    ("－","Time decay","Periodic reduction for inactivity"),
]
for i,(sign,label,note) in enumerate(pos_events):
    y = 1.65+i*0.72
    box(s,0.6,y,0.38,0.52,RGBColor(0xD1,0xFA,0xE5))
    tx(s,sign,0.63,y+0.05,0.32,0.38,Pt(14),True,GN,align=PP_ALIGN.CENTER)
    tx(s,label,1.1,y+0.04,3.5,0.28,Pt(10),True,DK)
    tx(s,note, 1.1,y+0.3,3.5,0.22,Pt(8.5),False,GR,italic=True)

for i,(sign,label,note) in enumerate(neg_events):
    y = 3.9+i*0.62
    box(s,0.6,y,0.38,0.48,RGBColor(0xFE,0xE2,0xE2))
    tx(s,sign,0.63,y+0.05,0.32,0.35,Pt(14),True,RD,align=PP_ALIGN.CENTER)
    tx(s,label,1.1,y+0.04,3.5,0.26,Pt(10),True,DK)
    tx(s,note, 1.1,y+0.28,3.5,0.22,Pt(8.5),False,GR,italic=True)

# right: tiers
box(s,6.9,1.12,5.95,5.45,W,LN,Pt(0.75))
box(s,6.9,1.12,5.95,0.4,B1)
tx(s,"Reputation Tiers",7.15,1.18,5,0.28,Pt(11),True,W)

tiers=[
    ("🌱","New","Starting tier — building first impressions",GR),
    ("🔵","Trusted","Consistent completions, positive ratings",B1),
    ("⭐","Reliable","Strong track record, low dispute rate",RGBColor(0xD9,0x77,0x06)),
    ("🏆","Veteran","Top tier — community pillar, unlocks trust signals",RGBColor(0x09,0x7C,0x4E)),
]
for i,(ico,tier,desc,c) in enumerate(tiers):
    y=1.68+i*1.15
    box(s,7.05,y,5.65,1.0,B4,LN,Pt(0.5))
    box(s,7.05,y,0.05,1.0,c)
    tx(s,ico,7.18,y+0.2,0.45,0.55,Pt(20))
    tx(s,tier,7.72,y+0.1,4,0.32,Pt(11),True,DK)
    tx(s,desc,7.72,y+0.45,4.85,0.46,Pt(9),False,GR,wrap=True)
footer(s,5)

# ═══════════════════════════════════════════════
# SLIDE 6 — Communities
# ═══════════════════════════════════════════════
s = slide()
header(s,"Communities","Scoped sub-networks — universities, locations, genres")

types=[
    ("🎓","University","HUST, UET, RMIT…\nClosed academic circles for textbook sharing within a campus"),
    ("📍","Location","District, city, or neighbourhood\n— local book swaps made easy"),
    ("📚","Genre / Interest","Non-fiction, Sci-fi, Vietnamese Lit\n— passion-driven reading circles"),
]
for i,(ico,t,d) in enumerate(types):
    x=0.45+i*4.25
    box(s,x,1.12,4.0,1.82,B4,LN,Pt(0.75))
    box(s,x,1.12,4.0,0.06,B1)
    tx(s,ico,x+0.2,1.28,0.6,0.55,Pt(24))
    tx(s,t,  x+0.9,1.3,3.0,0.38,Pt(12),True,DK)
    tx(s,d,  x+0.2,1.75,3.65,1.1,Pt(9.5),False,GR,wrap=True)

feats=[
    ("Public & Private","Private communities need an invite code — closed book clubs, department groups"),
    ("Community Posts","Members post discussions; liked, commented, pinned by moderators"),
    ("Scoped Listings","Listings visible only to community members — reduces noise"),
    ("Points Leaderboard","communityPoints earned through likes, comments, and activity"),
    ("Moderator Hierarchy","Owner → Moderator → Member with promote/demote controls"),
    ("Activity Notifications","All members notified of new posts; authors alerted on likes & comments"),
]
for i,(title,desc) in enumerate(feats):
    row,col=divmod(i,3)
    x=0.45+col*4.25
    y=3.18+row*1.42
    box(s,x,y,4.0,1.28,W,LN,Pt(0.75))
    box(s,x,y,0.05,1.28,B1)
    tx(s,title,x+0.2,y+0.12,3.65,0.3,Pt(10),True,DK)
    tx(s,desc, x+0.2,y+0.46,3.65,0.75,Pt(9),False,GR,wrap=True)
footer(s,6)

# ═══════════════════════════════════════════════
# SLIDE 7 — Discovery & Feed
# ═══════════════════════════════════════════════
s = slide()
header(s,"Discovery & Social Feed","How readers find the right book at the right time")

left=[
    ("🔍","Full-Text Search","Title, author, or ISBN. Filter by genre, condition, and transaction type (Gift / Exchange / Sell)."),
    ("🌐","Explore Page","Public browse — no login required. Great for acquisition and first impressions."),
    ("📰","Personalized Feed","Follow users. New listings fan-out into your feed in real time via Server-Sent Events."),
    ("🏘️","Community Browse","Listings and posts scoped to a group — less noise, higher relevance."),
]
for i,(ico,title,body) in enumerate(left):
    y=1.12+i*1.38
    box(s,0.45,y,6.3,1.24,W,LN,Pt(0.75))
    box(s,0.45,y,0.05,1.24,B1)
    tx(s,f"{ico}  {title}",0.65,y+0.1,5.9,0.32,Pt(10.5),True,DK)
    tx(s,body,0.65,y+0.46,5.9,0.7,Pt(9.5),False,GR,wrap=True)

# right: social graph details
box(s,7.2,1.12,5.65,5.45,B4,LN,Pt(0.75))
box(s,7.2,1.12,5.65,0.42,B1)
tx(s,"Social Graph",7.45,1.18,5,0.28,Pt(11),True,W)
social=[
    "Follow / unfollow any user",
    "Follower + following counts on profile",
    "Feed populated on publish (fan-out write)",
    "Real-time delivery via SSE (no polling)",
    "New listing notification to all followers",
    "Profile shows tier, bio, genres, district",
    "Public profiles — no login to view",
]
for i,item in enumerate(social):
    y=1.72+i*0.62
    box(s,7.38,y+0.08,0.26,0.26,B3)
    tx(s,"›",7.39,y+0.04,0.22,0.3,Pt(11),True,B1,align=PP_ALIGN.CENTER)
    tx(s,item,7.72,y+0.06,5.0,0.48,Pt(10),False,DK,wrap=True)
footer(s,7)

# ═══════════════════════════════════════════════
# SLIDE 8 — Notifications & Messaging
# ═══════════════════════════════════════════════
s = slide()
header(s,"Notifications & Messaging","Keeping every participant in the loop, in real time")

kinds=[
    ("📦","New listing from followed user"),
    ("🔄","Transaction status changed"),
    ("💬","New direct message"),
    ("📢","Community announcement"),
    ("⭐","Reputation tier changed"),
    ("🛡️","Moderation action taken"),
    ("📝","New post in your community"),
    ("❤️","Your post was liked"),
    ("🗨️","New comment on your post"),
]
tx(s,"9 Notification Kinds",0.45,1.12,6.5,0.35,Pt(11),True,DK)
for i,(ico,lbl) in enumerate(kinds):
    row,col=divmod(i,3)
    x=0.45+col*2.1
    y=1.55+row*0.62
    box(s,x,y,1.97,0.5,B4,LN,Pt(0.5))
    tx(s,f"{ico}  {lbl}",x+0.1,y+0.07,1.82,0.38,Pt(8.5),False,DK,wrap=True)

box(s,0.45,3.5,6.3,0.025,LN)

tx(s,"Email Preferences",0.45,3.62,6,0.3,Pt(10),True,DK)
for lbl,desc,x in [
    ("Immediate","Every event triggers an email",0.45),
    ("Daily Digest","Bundled once per day",2.65),
    ("Off","In-app only",4.85),
]:
    box(s,x,4.02,1.92,0.85,B4,LN,Pt(0.5))
    tx(s,lbl,x+0.12,4.08,1.72,0.28,Pt(9.5),True,B1)
    tx(s,desc,x+0.12,4.36,1.72,0.38,Pt(8.5),False,GR)

# right: messaging
box(s,7.2,1.12,5.65,4.75,W,LN,Pt(0.75))
box(s,7.2,1.12,5.65,0.42,B1)
tx(s,"💬  Real-Time Messaging",7.45,1.18,5,0.28,Pt(11),True,W)
msgs=[
    ("1-to-1 chat","Private conversations between any two users"),
    ("Transaction-linked","Thread auto-opens when a deal is accepted — full context always visible"),
    ("Server-Sent Events","Instant delivery with no WebSocket overhead"),
    ("Reportable","Any message can be reported directly to moderation queue"),
    ("Persistent history","Full conversation stored; no ephemeral-only sessions"),
]
for i,(title,desc) in enumerate(msgs):
    y=1.72+i*0.88
    box(s,7.38,y,5.28,0.75,B4)
    tx(s,title,7.55,y+0.06,4.9,0.26,Pt(10),True,DK)
    tx(s,desc, 7.55,y+0.36,4.9,0.34,Pt(9),False,GR,wrap=True)
footer(s,8)

# ═══════════════════════════════════════════════
# SLIDE 9 — Moderation & Safety
# ═══════════════════════════════════════════════
s = slide()
header(s,"Moderation & Platform Safety","Structured, fair, transparent — at scale")

tx(s,"What Can Be Reported",0.45,1.12,9,0.3,Pt(10),True,DK)
for lbl,x in [("👤  User",0.45),("📦  Listing",2.6),("🔄  Transaction",4.75),("💬  Message",6.9)]:
    box(s,x,1.5,1.92,0.52,B4,LN,Pt(0.5))
    tx(s,lbl,x+0.12,1.58,1.72,0.34,Pt(9.5),True,DK)

tx(s,"Moderator Actions",0.45,2.2,9,0.3,Pt(10),True,DK)
actions=[
    ("⚠️  Warn","Formal warning recorded — no immediate restriction"),
    ("🗑️  Remove Listing","Listing taken down and flagged"),
    ("🔇  Suspend User","Account suspended; cannot transact or post"),
    ("✅  Resolve Dispute","Moderator rules on disputed transaction"),
    ("❌  Reject Report","Report found unfounded — no action"),
    ("🔄  Restore","Lift suspension or re-activate listing"),
]
for i,(action,desc) in enumerate(actions):
    row,col=divmod(i,3)
    x=0.45+col*4.25
    y=2.58+row*1.0
    box(s,x,y,4.0,0.86,W,LN,Pt(0.75))
    box(s,x,y,0.05,0.86,B1)
    tx(s,action,x+0.2,y+0.08,3.65,0.28,Pt(10),True,DK)
    tx(s,desc,  x+0.2,y+0.4, 3.65,0.38,Pt(8.5),False,GR,wrap=True)

box(s,0.45,4.68,12.38,0.025,LN)

box(s,0.45,4.82,5.95,1.25,RGBColor(0xFF,0xF8,0xED),RGBColor(0xFB,0xBF,0x24),Pt(1))
tx(s,"🤖  Auto-Flagging",0.65,4.9,5.5,0.28,Pt(10),True,RGBColor(0x78,0x35,0x00))
tx(s,"System can auto-generate reports on suspicious patterns — feeds the queue without requiring manual user reports.",
   0.65,5.22,5.7,0.75,Pt(9),False,DK,wrap=True)

box(s,6.75,4.82,5.95,1.25,RGBColor(0xF0,0xFD,0xF4),RGBColor(0x34,0xD3,0x99),Pt(1))
tx(s,"🛡️  Admin Dashboard",6.95,4.9,5.5,0.28,Pt(10),True,RGBColor(0x06,0x5F,0x46))
tx(s,"Separate /admin route for platform-wide oversight: manage users, listings, and communities at system level.",
   6.95,5.22,5.7,0.75,Pt(9),False,DK,wrap=True)
footer(s,9)

# ═══════════════════════════════════════════════
# SLIDE 10 — User Roles
# ═══════════════════════════════════════════════
s = slide()
header(s,"User Roles & Access Hierarchy","Clear permissions at every level")

roles=[
    ("👁️","Guest","No account needed",
     ["Browse public listings","Full-text search","View public profiles","Explore page"],GR),
    ("👤","User","Verified account",
     ["Create & manage listings","Request transactions","Rate counterparty","Follow users","Join communities","Send messages"],B1),
    ("🔧","Moderator","Elevated community role",
     ["All User capabilities","Moderation queue access","Take moderation actions","Award community points","Pin posts · Promote members"],RGBColor(0x09,0x7C,0x4E)),
    ("⚙️","Admin","Platform authority",
     ["Full platform access","Admin dashboard","System moderation","User management","Override any action"],RGBColor(0x7C,0x27,0x08)),
]
for i,(ico,role,sub,perms,color) in enumerate(roles):
    x=0.38+i*3.15
    box(s,x,1.12,2.95,5.22,W,LN,Pt(0.75))
    box(s,x,1.12,2.95,0.62,color)
    tx(s,f"{ico}  {role}",x+0.15,1.18,2.7,0.35,Pt(13),True,W)
    tx(s,sub,x+0.15,1.56,2.7,0.25,Pt(8.5),False,RGBColor(0xD1,0xD5,0xDB),italic=True)
    for j,perm in enumerate(perms):
        box(s,x+0.18,1.95+j*0.52,0.06,0.24,color)
        tx(s,perm,x+0.32,1.93+j*0.52,2.55,0.38,Pt(9),False,DK,wrap=True)

box(s,0.38,6.45,12.5,0.55,B4,B1,Pt(0.75))
tx(s,"Account lifecycle:  Pending Verification  →  Active  →  Suspended / Deleted  ·  Email verification required before Active.",
   0.65,6.5,12.0,0.45,Pt(9.5),False,DK,wrap=True)
footer(s,10)

# ═══════════════════════════════════════════════
# SLIDE 11 — Business Model
# ═══════════════════════════════════════════════
s = slide()
header(s,"Business Model & Market Fit","How BookBridge creates and captures value")

box(s,0.45,1.12,12.38,0.68,B1)
tx(s,"Positioning — BookBridge is a social trust network for book circulation, not a marketplace.",
   0.7,1.22,11.9,0.5,Pt(12),True,W,wrap=True)

pillars=[
    ("🏫","Campus-First GTM","Launch within universities — where textbook pain is sharpest. Partner with student unions and library programs for distribution and trust bootstrap."),
    ("🔗","Network Effect","Value compounds: more givers → more receivers → more ratings → stronger trust → more engagement. Each user makes the platform better for everyone."),
    ("💰","Monetization Paths","Premium verified-seller badge · University partnership licensing · Local business listing sponsorships · API for library integrations."),
    ("🔒","Defensible Moat","Reputation equity. A user's ReputationScore and community standing can't be transferred to another platform — switching cost grows with every transaction."),
]
for i,(ico,title,body) in enumerate(pillars):
    row,col=divmod(i,2)
    x=0.45+col*6.3
    y=2.02+row*1.82
    box(s,x,y,5.98,1.65,W,LN,Pt(0.75))
    box(s,x,y,0.05,1.65,B1)
    tx(s,ico,x+0.2,y+0.18,0.55,0.55,Pt(20))
    tx(s,title,x+0.85,y+0.18,4.9,0.3,Pt(11),True,DK)
    tx(s,body,x+0.85,y+0.55,4.9,1.0,Pt(9.5),False,GR,wrap=True)

box(s,0.45,5.73,12.38,0.025,LN)
tx(s,"Price cap enforcement is intentional — it keeps BookBridge non-commercial, builds social goodwill, and differentiates from generic resale apps.",
   0.45,5.85,12.38,0.55,Pt(9.5),False,GR,italic=True,wrap=True)
footer(s,11)

# ═══════════════════════════════════════════════
# SLIDE 12 — Summary & Next Steps
# ═══════════════════════════════════════════════
s = slide()
box(s,0,0,13.33,0.06,B1)
box(s,0,0.06,13.33,3.1,DK)

tx(s,"BookBridge",0.65,0.28,10,0.85,Pt(44),True,W)
tx(s,"Share books.  Build trust.  Keep stories moving.",
   0.65,1.18,10,0.45,Pt(15),False,RGBColor(0x93,0xC5,0xFD),italic=True)

summary=[
    "Community-driven platform for second-hand book circulation",
    "End-to-end transaction state machine with full audit trail",
    "Reputation system with atomic score events and tiers",
    "Sub-communities: university, location, and genre scoped",
    "Real-time feed, messaging, and 9-kind notifications",
    "Structured moderation queue with admin oversight",
]
for i,item in enumerate(summary):
    x=0.65+(i%2)*5.9
    y=1.78+(i//2)*0.38
    tx(s,f"✓  {item}",x,y,5.7,0.34,Pt(9.5),False,RGBColor(0xBE,0xD8,0xFF),wrap=True)

box(s,0,3.16,13.33,0.05,B1)

tx(s,"Next Steps",0.65,3.32,12,0.38,Pt(14),True,DK)
nexts=[
    ("🚀","Production Deploy","Run prisma migrate deploy on Vercel production DB; smoke-test all critical flows end-to-end"),
    ("🔔","Email Notifications","Wire dispatcher to SendGrid/Resend for immediate, digest, and off preferences"),
    ("📊","Analytics","Lightweight event tracking: page views, listing clicks, conversion funnel"),
    ("💳","Monetization Pilot","University partnership pilot — measure engagement before adding premium features"),
]
for i,(ico,title,body) in enumerate(nexts):
    x=0.45+i*3.15
    box(s,x,3.82,2.98,2.42,W,LN,Pt(0.75))
    box(s,x,3.82,2.98,0.42,B4)
    box(s,x,3.82,0.05,2.42,B1)
    tx(s,f"{ico}  {title}",x+0.2,3.88,2.65,0.3,Pt(10),True,DK)
    tx(s,body,x+0.2,4.32,2.68,1.82,Pt(9),False,GR,wrap=True)

tx(s,"Next.js 14  ·  PostgreSQL  ·  Prisma  ·  Vercel  ·  Team Zootopia — SE Capstone 2025",
   0.45,6.5,12.3,0.38,Pt(8.5),False,GR,align=PP_ALIGN.CENTER,italic=True)
footer(s,12)

# ── Save ──────────────────────────────────────
out = r"C:\Users\ADMIN\Desktop\SE\SE_book_bridge\BookBridge_v2.pptx"
prs.save(out)
print("Saved:", out)

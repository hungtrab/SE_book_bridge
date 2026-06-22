from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ──────────────────────────────────────────────
W   = RGBColor(0xFF,0xFF,0xFF)
BG  = RGBColor(0xF8,0xFA,0xFF)   # near-white page bg
B   = RGBColor(0x18,0x55,0xE8)   # primary blue
BL  = RGBColor(0xE8,0xF0,0xFF)   # light blue fill
BD  = RGBColor(0x0D,0x35,0x9A)   # dark blue
INK = RGBColor(0x0F,0x17,0x2A)   # near-black
MID = RGBColor(0x47,0x55,0x69)   # body gray
SUB = RGBColor(0x94,0xA3,0xB8)   # subtitle / caption
DIV = RGBColor(0xE2,0xE8,0xF0)   # divider / border
GRN = RGBColor(0x05,0x96,0x69)   # green accent
RED = RGBColor(0xDC,0x26,0x26)   # red accent
AMB = RGBColor(0xD9,0x77,0x06)   # amber

blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(blank)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = W
    return s

def R(s, l, t, w, h, fc, lc=None, lw=Pt(0)):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if lc:
        sh.line.color.rgb = lc; sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def T(s, text, l, t, w, h, size=Pt(11), bold=False, color=INK,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def footer(s, n, total=12):
    R(s, 0, 7.3, 13.33, 0.2, DIV)
    T(s, "BookBridge  ·  Business Overview", 0.5, 7.32, 9, 0.16,
      size=Pt(7.5), color=SUB)
    T(s, f"{n} / {total}", 12.0, 7.32, 1.1, 0.16,
      size=Pt(7.5), color=SUB, align=PP_ALIGN.RIGHT)

def label_chip(s, text, l, t, color=B):
    R(s, l, t, len(text)*0.085+0.3, 0.28, BL)
    T(s, text, l+0.12, t+0.03, len(text)*0.085+0.1, 0.22,
      size=Pt(8), bold=True, color=color)

# ═══════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════
s = slide()
# full-bleed blue left panel
R(s, 0, 0, 5.8, 7.5, B)
# decorative circle top-right of panel (lighter blue)
sh = s.shapes.add_shape(9, Inches(3.2), Inches(-1.5), Inches(4.5), Inches(4.5))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x22,0x65,0xF0)
sh.line.fill.background()
sh = s.shapes.add_shape(9, Inches(0.2), Inches(5.2), Inches(3.0), Inches(3.0))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x14,0x46,0xCB)
sh.line.fill.background()

T(s, "📚", 0.55, 1.1, 1.2, 1.2, size=Pt(52))
T(s, "BookBridge", 0.55, 2.05, 5.0, 0.95,
  size=Pt(44), bold=True, color=W)
T(s, "Share books. Build trust.\nKeep stories moving.", 0.55, 3.05, 5.0, 0.9,
  size=Pt(15), color=RGBColor(0xBB,0xD2,0xFF), italic=True)
T(s, "Business Overview", 0.55, 4.1, 4.5, 0.38,
  size=Pt(10), color=RGBColor(0x93,0xB9,0xFF), bold=True)

# right side — white
T(s, "What is BookBridge?", 6.3, 0.7, 6.5, 0.45,
  size=Pt(20), bold=True, color=INK)
R(s, 6.3, 1.2, 0.06, 0.0, B)  # thin rule

desc = (
    "A community-driven platform for second-hand book sharing "
    "among students and local readers.\n\n"
    "Users gift, exchange, or sell pre-owned books at symbolic prices — "
    "underpinned by a trust-driven reputation system and scoped "
    "sub-communities for universities, neighbourhoods, and reading circles."
)
T(s, desc, 6.3, 1.28, 6.6, 2.0, size=Pt(11.5), color=MID, wrap=True)

# 4 stat pills
stats = [("3", "Transaction\nModes"), ("8+", "Core\nFeatures"),
         ("4", "User\nRoles"), ("∞", "Stories\nShared")]
for i, (num, lbl) in enumerate(stats):
    x = 6.3 + i * 1.68
    R(s, x, 3.55, 1.5, 1.6, BL)
    R(s, x, 3.55, 1.5, 0.04, B)
    T(s, num, x, 3.65, 1.5, 0.65,
      size=Pt(32), bold=True, color=B, align=PP_ALIGN.CENTER)
    T(s, lbl, x, 4.3, 1.5, 0.75,
      size=Pt(8.5), color=MID, align=PP_ALIGN.CENTER)

T(s, "Built with Next.js 14 · PostgreSQL · Prisma · Vercel\nTeam Zootopia — 2025 Software Engineering Capstone",
  6.3, 5.45, 6.6, 0.7, size=Pt(9), color=SUB, italic=True)

footer(s, 1)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "The Problem", 0.6, 0.2, 8, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Why students and readers need something better", 0.6, 0.72, 10, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

problems = [
    ("📦", "Books Go to Waste",
     "Millions of textbooks sit unused after each semester — students have no trusted, friction-free channel to pass them on."),
    ("💸", "Cost of Learning",
     "New textbooks are prohibitively expensive. Students desperately need gift, exchange, or low-cost resale options."),
    ("🤝", "No Trust Layer",
     "Generic marketplaces lack the reputation infrastructure needed for safe peer-to-peer sharing within academic circles."),
    ("🏘️", "Fragmented Communities",
     "University and neighbourhood reading groups have no dedicated space to share books, coordinate swaps, and build culture."),
]
for i, (icon, title, body) in enumerate(problems):
    x = 0.55 + i * 3.08
    R(s, x, 1.25, 2.85, 4.5, W, DIV, Pt(1))
    R(s, x, 1.25, 2.85, 0.05, B)
    T(s, icon, x+0.2, 1.45, 0.7, 0.65, size=Pt(28))
    T(s, title, x+0.2, 2.18, 2.5, 0.38, size=Pt(12), bold=True, color=INK)
    T(s, body,  x+0.2, 2.62, 2.5, 2.8,  size=Pt(10), color=MID, wrap=True)

R(s, 0.55, 6.05, 12.15, 0.78, BL)
T(s, "💡  BookBridge is the answer — a trusted social network for book circulation, scoped to the communities that matter most.",
  0.8, 6.12, 11.7, 0.62, size=Pt(11.5), bold=True, color=BD, wrap=True)

footer(s, 2)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — PRODUCT AT A GLANCE
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Product at a Glance", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Eight integrated features, one coherent experience", 0.6, 0.72, 10, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

feats = [
    ("📖", "Book Listings",    "Gift · Exchange · Sell\nCondition ratings + photos\nCommunity-scoped or public"),
    ("🔍", "Discovery",        "Full-text search by title,\nauthor, or ISBN\nPersonalised social feed"),
    ("🔄", "Transactions",     "8-step state machine\nFull audit trail\nIn-person or postal delivery"),
    ("💬", "Messaging",        "Real-time 1-to-1 chat\nTransaction-linked threads\nServer-Sent Events"),
    ("⭐", "Reputation",       "Atomic score events\nReputation tiers\nTrust visible on every profile"),
    ("🏘️", "Communities",     "University · Location · Genre\nPosts, likes, comments\nPoints leaderboard"),
    ("🔔", "Notifications",    "9 notification kinds\nIn-app + email preferences\nReal-time delivery"),
    ("🛡️", "Moderation",      "Structured report queue\nWarn · Remove · Suspend\nAdmin dashboard"),
]
for i, (icon, title, body) in enumerate(feats):
    col = i % 4; row = i // 4
    x = 0.45 + col * 3.2
    y = 1.3  + row * 2.6
    R(s, x, y, 3.0, 2.35, W, DIV, Pt(1))
    R(s, x, y, 3.0, 0.42, BL)
    T(s, f"{icon}  {title}", x+0.15, y+0.08, 2.75, 0.3,
      size=Pt(11), bold=True, color=INK)
    T(s, body, x+0.15, y+0.52, 2.75, 1.7,
      size=Pt(9.5), color=MID, wrap=True)

footer(s, 3)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — TRANSACTION WORKFLOW
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "How a Book Changes Hands", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "End-to-end transaction state machine — from listing to reputation update", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

steps = [
    ("List",     "Seller creates listing:\ncondition, photos, type"),
    ("Discover", "Buyer finds via search,\nfeed, or community"),
    ("Request",  "Buyer sends request\n→ status: Pending"),
    ("Accept",   "Seller accepts\n→ listing: Reserved"),
    ("Chat",     "Auto thread opens\nfor coordination"),
    ("Deliver",  "Ship + tracking\nor in-person handoff"),
    ("Complete", "Buyer confirms\n→ status: Completed"),
    ("Rate",     "Mutual star ratings\nupdate reputation"),
]
sw = 1.44
for i, (title, body) in enumerate(steps):
    x = 0.42 + i * 1.56
    # number bubble
    sh = s.shapes.add_shape(9, Inches(x+0.47), Inches(1.25), Inches(0.5), Inches(0.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = B; sh.line.fill.background()
    T(s, str(i+1), x+0.47, 1.25, 0.5, 0.5,
      size=Pt(10), bold=True, color=W, align=PP_ALIGN.CENTER)
    # step box
    R(s, x, 1.88, sw, 1.85, BL)
    T(s, title, x+0.08, 1.94, sw-0.16, 0.3, size=Pt(10), bold=True, color=BD)
    T(s, body,  x+0.08, 2.28, sw-0.16, 1.4, size=Pt(8.5), color=MID, wrap=True)
    # arrow
    if i < 7:
        T(s, "›", x+sw+0.04, 2.6, 0.18, 0.35, size=Pt(18), bold=True, color=B)

# bottom callouts
R(s, 0.42, 3.95, 5.85, 1.08, RGBColor(0xFF,0xF8,0xEC), RGBColor(0xFB,0xBF,0x24), Pt(1.5))
T(s, "⚠️  Dispute Flow",
  0.62, 4.02, 5.5, 0.3, size=Pt(10), bold=True, color=AMB)
T(s, "Either party can raise a dispute at any point. A moderator reviews and resolves or rejects via a structured moderation action.",
  0.62, 4.36, 5.55, 0.6, size=Pt(9.5), color=MID, wrap=True)

R(s, 7.0, 3.95, 5.85, 1.08, RGBColor(0xF0,0xFD,0xF4), RGBColor(0x34,0xD3,0x99), Pt(1.5))
T(s, "✅  Reputation Impact",
  7.2, 4.02, 5.5, 0.3, size=Pt(10), bold=True, color=GRN)
T(s, "Completion awards positive score to both buyer and seller. Disputes and cancellations carry negative deltas — keeping incentives aligned.",
  7.2, 4.36, 5.55, 0.6, size=Pt(9.5), color=MID, wrap=True)

# waitlist note
T(s, "Seller can also Waitlist or Decline a request. Accepted request auto-opens a conversation thread and locks the listing.",
  0.42, 5.25, 12.4, 0.4, size=Pt(9), color=SUB, italic=True, wrap=True)

footer(s, 4)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — REPUTATION & TRUST
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Trust & Reputation Engine", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "The layer that makes peer-to-peer sharing safe at scale", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

# left — score events
T(s, "Score Events", 0.6, 1.25, 5.5, 0.35, size=Pt(13), bold=True, color=INK)
events = [
    (GRN, "+", "Transaction completed"),
    (GRN, "+", "Positive rating received"),
    (GRN, "+", "Community contribution (mod-awarded)"),
    (RED, "−", "Dispute upheld against user"),
    (RED, "−", "Cancellation penalty"),
    (RED, "−", "Report upheld by moderator"),
    (RED, "−", "Periodic time decay"),
]
for i, (clr, sign, lbl) in enumerate(events):
    y = 1.72 + i * 0.52
    R(s, 0.6, y+0.04, 0.32, 0.32, clr)
    T(s, sign, 0.6, y+0.04, 0.32, 0.32, size=Pt(14), bold=True, color=W, align=PP_ALIGN.CENTER)
    T(s, lbl, 1.05, y+0.06, 5.1, 0.38, size=Pt(10.5), color=INK)

# right — tiers
T(s, "Reputation Tiers", 7.2, 1.25, 5.5, 0.35, size=Pt(13), bold=True, color=INK)
tiers = [
    (B, "🌱", "New",      "Starting tier — building first impressions"),
    (RGBColor(0x03,0x69,0xBB), "🔵", "Trusted",  "Consistent completions, positive ratings"),
    (RGBColor(0x07,0x52,0x85), "⭐", "Reliable", "Strong track record, low dispute rate"),
    (BD,                       "🏆", "Veteran",  "Top tier — unlocks trust signals, community pillar"),
]
for i, (clr, icon, tier, desc) in enumerate(tiers):
    y = 1.72 + i * 1.12
    R(s, 7.2, y, 5.6, 0.95, W, DIV, Pt(1))
    R(s, 7.2, y, 0.06, 0.95, clr)
    T(s, icon, 7.35, y+0.18, 0.5, 0.5, size=Pt(20))
    T(s, tier, 7.88, y+0.1, 2.5, 0.35, size=Pt(12), bold=True, color=INK)
    T(s, desc, 7.88, y+0.46, 4.8, 0.42, size=Pt(9.5), color=MID, wrap=True)

footer(s, 5)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — COMMUNITIES
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Communities", 0.6, 0.2, 8, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Scoped sub-networks — universities, neighbourhoods, reading circles", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

# 3 type cards — top row
types = [
    ("🎓", "University",  "Closed academic circles for textbook sharing within a campus (e.g. HUST, UET, RMIT)"),
    ("📍", "Location",    "District, city, or neighbourhood groups — ideal for local book swaps"),
    ("📚", "Genre / Interest", "Passion-driven reading circles: Non-fiction, Sci-fi, Vietnamese Lit…"),
]
for i, (icon, t, d) in enumerate(types):
    x = 0.55 + i * 4.18
    R(s, x, 1.28, 3.9, 1.72, BL)
    R(s, x, 1.28, 3.9, 0.06, B)
    T(s, icon, x+0.2, 1.45, 0.7, 0.6, size=Pt(26))
    T(s, t,    x+1.0, 1.48, 2.8, 0.38, size=Pt(13), bold=True, color=INK)
    T(s, d,    x+0.2, 1.88, 3.58, 1.0, size=Pt(9.5), color=MID, wrap=True)

# 6 feature bullets — two columns
feats = [
    ("🔒", "Public & Private modes",   "Private groups require an invite code — great for closed reading clubs"),
    ("📝", "Community Posts",          "Members post, like, comment; moderators can pin important announcements"),
    ("📦", "Scoped Listings",          "Listings published to a community are only visible to its members"),
    ("🏅", "Points Leaderboard",       "communityPoints earned via likes, comments, and moderated contributions"),
    ("👑", "Moderator Hierarchy",      "Owner → Moderator → Member with promote/demote controls"),
    ("🔔", "Activity Notifications",   "All members notified on new posts; authors alerted on likes and comments"),
]
for i, (icon, title, desc) in enumerate(feats):
    col = i % 2; row = i // 2
    x = 0.55 + col * 6.3
    y = 3.22 + row * 1.2
    R(s, x, y, 6.0, 1.08, W, DIV, Pt(1))
    T(s, f"{icon}  {title}", x+0.2, y+0.1, 5.6, 0.32, size=Pt(10.5), bold=True, color=INK)
    T(s, desc, x+0.2, y+0.46, 5.65, 0.55, size=Pt(9.5), color=MID, wrap=True)

footer(s, 6)

# ═══════════════════════════════════════════════════════
# SLIDE 7 — DISCOVERY & SOCIAL
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Discovery & Social Feed", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "How readers find the right books at the right time", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

channels = [
    ("🔍", "Full-Text Search",
     "Search by title, author, or ISBN.\nFilter by genre, condition, and transaction type (Gift / Exchange / Sell)."),
    ("🌐", "Explore Page",
     "Public browse — no login required.\nLowers the acquisition barrier; first impressions before sign-up."),
    ("📰", "Personalised Feed",
     "Follow any user. Every listing they publish fans into your feed in real time via Server-Sent Events."),
    ("🏘️", "Community Browse",
     "Within a community, members see listings and posts scoped only to that group — less noise, more relevance."),
]
for i, (icon, title, body) in enumerate(channels):
    y = 1.28 + i * 1.38
    R(s, 0.55, y, 0.06, 1.15, B)
    R(s, 0.65, y, 6.05, 1.15, W, DIV, Pt(1))
    T(s, icon, 0.82, y+0.12, 0.6, 0.6, size=Pt(22))
    T(s, title, 1.52, y+0.14, 5.1, 0.32, size=Pt(11), bold=True, color=INK)
    T(s, body,  1.52, y+0.5,  5.05, 0.6, size=Pt(9.5), color=MID, wrap=True)

# right panel — social graph
R(s, 7.3, 1.28, 5.5, 5.25, BL)
T(s, "Social Graph", 7.55, 1.36, 5.1, 0.38, size=Pt(13), bold=True, color=INK)
social = [
    "Follow / unfollow any user",
    "Follower + following counts on every profile",
    "Feed populated on publish — fan-out write",
    "Real-time delivery via Server-Sent Events",
    "New listing notification sent to all followers",
    "Profiles show reputation tier, bio, genres, district",
]
for i, item in enumerate(social):
    y = 1.9 + i * 0.68
    R(s, 7.55, y+0.1, 0.06, 0.38, B)
    T(s, item, 7.75, y+0.06, 4.85, 0.55, size=Pt(10), color=INK, wrap=True)

footer(s, 7)

# ═══════════════════════════════════════════════════════
# SLIDE 8 — NOTIFICATIONS & MESSAGING
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Notifications & Messaging", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Every participant stays in context — instantly", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

# left — notifications
T(s, "9 Notification Kinds", 0.55, 1.25, 6, 0.35, size=Pt(12), bold=True, color=INK)
notifs = [
    ("📦", "New listing from followed user"),
    ("🔄", "Transaction status changed"),
    ("💬", "New direct message"),
    ("📢", "Community announcement"),
    ("⭐", "Reputation tier changed"),
    ("🛡️", "Moderation action taken"),
    ("📝", "New post in your community"),
    ("❤️", "Your post was liked"),
    ("🗨️", "New comment on your post"),
]
for i, (icon, lbl) in enumerate(notifs):
    col = i % 3; row = i // 3
    x = 0.55 + col * 2.08
    y = 1.72 + row * 0.56
    R(s, x, y, 1.95, 0.46, W, DIV, Pt(0.75))
    T(s, f"{icon}  {lbl}", x+0.1, y+0.06, 1.8, 0.36, size=Pt(8.5), color=INK, wrap=True)

R(s, 0.55, 3.44, 6.25, 0.025, DIV)
T(s, "Email Preferences", 0.55, 3.58, 5.5, 0.3, size=Pt(11), bold=True, color=INK)
prefs = [("Immediate", "Every event"), ("Daily Digest", "Bundled daily"), ("Off", "In-app only")]
for i, (lbl, desc) in enumerate(prefs):
    x = 0.55 + i * 2.1
    R(s, x, 3.98, 1.92, 0.82, BL)
    T(s, lbl,  x+0.12, 4.05, 1.7, 0.28, size=Pt(10), bold=True, color=B)
    T(s, desc, x+0.12, 4.36, 1.7, 0.38, size=Pt(9),  color=MID)

# right panel — messaging
R(s, 7.3, 1.25, 5.5, 5.1, W, DIV, Pt(1))
R(s, 7.3, 1.25, 5.5, 0.06, B)
T(s, "💬  Real-Time Messaging", 7.5, 1.38, 5.1, 0.38, size=Pt(13), bold=True, color=INK)
msg = [
    ("1-to-1 Chat",         "Private conversations between any two users"),
    ("Transaction-Linked",  "Auto-opened thread on acceptance — context always at hand"),
    ("Server-Sent Events",  "Real-time delivery without WebSocket complexity"),
    ("Reportable Messages", "Any message can be flagged; lands in the moderation queue"),
    ("Persistent History",  "Full message history stored — no ephemeral-only threads"),
]
for i, (title, desc) in enumerate(msg):
    y = 1.88 + i * 0.85
    R(s, 7.5, y, 4.95, 0.72, BL)
    T(s, title, 7.65, y+0.07, 4.6, 0.28, size=Pt(10), bold=True, color=INK)
    T(s, desc,  7.65, y+0.38, 4.6, 0.3,  size=Pt(9),  color=MID, wrap=True)

footer(s, 8)

# ═══════════════════════════════════════════════════════
# SLIDE 9 — MODERATION
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Moderation & Platform Safety", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Structured, fair, transparent — trust at scale", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

T(s, "Reportable Objects", 0.55, 1.2, 7, 0.32, size=Pt(11), bold=True, color=INK)
for icon_lbl, x in [("👤  User", 0.55), ("📦  Listing", 2.7), ("🔄  Transaction", 4.85), ("💬  Message", 7.3)]:
    R(s, x, 1.58, 1.88, 0.48, BL)
    T(s, icon_lbl, x+0.12, 1.64, 1.65, 0.34, size=Pt(10), bold=True, color=BD)

T(s, "Moderator Actions", 0.55, 2.28, 7, 0.32, size=Pt(11), bold=True, color=INK)
actions = [
    ("⚠️  Warn",           "Formal warning — no immediate restriction"),
    ("🗑️  Remove Listing",  "Listing taken down and flagged"),
    ("🔇  Suspend User",    "Account locked; can't transact or post"),
    ("✅  Resolve Dispute", "Moderator rules on a disputed transaction"),
    ("❌  Reject Report",   "Report unfounded — no action taken"),
    ("🔄  Restore",         "Suspension lifted or listing re-activated"),
]
for i, (act, desc) in enumerate(actions):
    col = i % 3; row = i // 2
    x = 0.55 + col * 4.12
    y = 2.72 + (i // 3) * 1.0
    R(s, x, y, 3.88, 0.84, W, DIV, Pt(1))
    T(s, act,  x+0.15, y+0.08, 3.6, 0.3,  size=Pt(10), bold=True, color=INK)
    T(s, desc, x+0.15, y+0.42, 3.6, 0.36, size=Pt(9),  color=MID, wrap=True)

R(s, 0.55, 4.88, 5.75, 1.12, RGBColor(0xFF,0xF9,0xEC), RGBColor(0xFB,0xBF,0x24), Pt(1.5))
T(s, "🤖  Auto-Flagging",  0.75, 4.95, 5.4, 0.28, size=Pt(10), bold=True, color=AMB)
T(s, "System can auto-generate reports on suspicious patterns — feeding the queue without requiring manual user reports.",
  0.75, 5.27, 5.4, 0.65, size=Pt(9.5), color=MID, wrap=True)

R(s, 6.75, 4.88, 5.75, 1.12, RGBColor(0xF0,0xFD,0xF4), RGBColor(0x34,0xD3,0x99), Pt(1.5))
T(s, "🛡️  Admin Dashboard", 6.95, 4.95, 5.4, 0.28, size=Pt(10), bold=True, color=GRN)
T(s, "Separate /admin route for system-level oversight — manage users, listings, and moderation across all communities.",
  6.95, 5.27, 5.4, 0.65, size=Pt(9.5), color=MID, wrap=True)

footer(s, 9)

# ═══════════════════════════════════════════════════════
# SLIDE 10 — USER ROLES
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "User Roles & Access Hierarchy", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "Four tiers of access — clear permissions at every level", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

roles = [
    (SUB,                        "👁️",  "Guest",     "No account required",
     ["Browse public listings", "Full-text search", "View public profiles", "Explore page"]),
    (B,                          "👤",  "User",      "Verified account",
     ["Create & manage listings", "Request transactions", "Rate & review", "Follow users", "Join communities", "1-to-1 messaging"]),
    (RGBColor(0x05,0x80,0x56),   "🔧", "Moderator", "Elevated community role",
     ["All User capabilities", "Moderation queue access", "Warn / remove / suspend", "Award community points", "Pin posts", "Promote / demote members"]),
    (RGBColor(0x7C,0x27,0x08),   "⚙️", "Admin",     "Platform-wide authority",
     ["Full platform access", "Admin dashboard", "System-level moderation", "User management", "Override any action"]),
]
for i, (clr, icon, role, sub, perms) in enumerate(roles):
    x = 0.45 + i * 3.15
    R(s, x, 1.28, 2.98, 5.2, W, DIV, Pt(1))
    R(s, x, 1.28, 2.98, 0.72, clr)
    T(s, f"{icon}  {role}", x+0.15, 1.33, 2.75, 0.38, size=Pt(14), bold=True, color=W)
    T(s, sub,               x+0.15, 1.72, 2.75, 0.26, size=Pt(8.5), color=RGBColor(0xD1,0xD5,0xDB), italic=True)
    for j, perm in enumerate(perms):
        R(s, x+0.15, 2.14+j*0.5, 0.07, 0.28, clr)
        T(s, perm, x+0.3, 2.12+j*0.5, 2.55, 0.38, size=Pt(9), color=INK, wrap=True)

R(s, 0.45, 6.52, 12.4, 0.5, BL)
T(s, "Account lifecycle:  Pending Verification  →  Active  →  Suspended / Deleted  ·  Email verification required before account becomes Active",
  0.7, 6.58, 11.9, 0.38, size=Pt(9), color=BD, wrap=True)

footer(s, 10)

# ═══════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════
s = slide()
R(s, 0, 0, 13.33, 0.06, B)
T(s, "Business Model & Market Fit", 0.6, 0.2, 9, 0.5, size=Pt(28), bold=True, color=INK)
T(s, "How BookBridge creates and captures value", 0.6, 0.72, 11, 0.3,
  size=Pt(12), color=SUB, italic=True)
R(s, 0.6, 1.08, 12.0, 0.025, DIV)

R(s, 0.55, 1.25, 12.2, 0.65, B)
T(s, "Positioning:  BookBridge is a social trust network for book circulation — not a marketplace.",
  0.8, 1.34, 11.8, 0.48, size=Pt(11.5), bold=True, color=W, wrap=True)

pillars = [
    ("🏫", "Campus-First Go-to-Market",
     "Launch within universities — where textbook pain is sharpest and community density is highest. "
     "Partner with student unions and library programs for rapid distribution."),
    ("🤝", "Network Effect",
     "Value compounds with every user: more givers → more receivers → more ratings → "
     "stronger trust signals → deeper engagement. Classic two-sided network."),
    ("💰", "Monetisation Paths",
     "Premium verified-seller badge · University partnership licensing · "
     "Local business listing sponsorships · API for library system integrations."),
    ("🔒", "Defensible Moat",
     "Reputation equity. A user's score and community standing can't be transferred to a competitor — "
     "switching cost grows with every completed transaction."),
]
for i, (icon, title, body) in enumerate(pillars):
    row = i // 2; col = i % 2
    x = 0.55 + col * 6.3
    y = 2.12 + row * 1.88
    R(s, x, y, 6.0, 1.72, W, DIV, Pt(1))
    R(s, x, y, 0.07, 1.72, B)
    T(s, icon,  x+0.25, y+0.18, 0.65, 0.65, size=Pt(24))
    T(s, title, x+1.0,  y+0.15, 4.85, 0.35, size=Pt(11), bold=True, color=INK)
    T(s, body,  x+1.0,  y+0.56, 4.85, 1.08, size=Pt(9.5), color=MID, wrap=True)

T(s, "Price cap enforcement is deliberate — it keeps the platform non-commercial, builds social goodwill, and differentiates BookBridge from generic resale apps.",
  0.55, 6.1, 12.2, 0.45, size=Pt(9), color=SUB, italic=True, wrap=True)

footer(s, 11)

# ═══════════════════════════════════════════════════════
# SLIDE 12 — CLOSING / NEXT STEPS
# ═══════════════════════════════════════════════════════
s = slide()
# full dark header
R(s, 0, 0, 13.33, 3.2, INK)
sh = s.shapes.add_shape(9, Inches(9.5), Inches(-2.0), Inches(6.0), Inches(6.0))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x1A,0x30,0x5C)
sh.line.fill.background()
sh = s.shapes.add_shape(9, Inches(10.8), Inches(0.8), Inches(4.0), Inches(4.0))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x14,0x27,0x4E)
sh.line.fill.background()

T(s, "BookBridge", 0.7, 0.3, 10, 0.85, size=Pt(46), bold=True, color=W)
T(s, "Share books. Build trust. Keep stories moving.", 0.7, 1.2, 10, 0.45,
  size=Pt(15), color=RGBColor(0x93,0xC5,0xFD), italic=True)

summary = [
    "Community-driven platform for second-hand book circulation",
    "End-to-end transaction state machine with full audit trail",
    "Atomic reputation system with tiers and trust signals",
    "Sub-communities scoped by university, location, and genre",
    "Real-time feed, SSE messaging, and 9-kind notifications",
    "Structured moderation queue with admin oversight",
]
for i, item in enumerate(summary):
    col = i % 2; row = i // 2
    x = 0.7 + col * 5.9
    y = 1.82 + row * 0.38
    T(s, "✓  " + item, x, y, 5.7, 0.34,
      size=Pt(9.5), color=RGBColor(0xBF,0xD7,0xFF), wrap=True)

R(s, 0, 3.2, 13.33, 0.05, B)

T(s, "Next Steps", 0.65, 3.42, 12, 0.42, size=Pt(16), bold=True, color=INK)
nexts = [
    ("🚀", "Production Deploy",      "Run migration on production DB; smoke-test all critical transaction flows on Vercel"),
    ("📧", "Email Notifications",    "Wire dispatcher to SendGrid/Resend for immediate, digest, and off preferences"),
    ("📊", "Analytics",              "Add event tracking: page views, listing clicks, transaction conversion funnel"),
    ("💳", "Monetisation Pilot",     "University partnership pilot — measure engagement before introducing premium features"),
]
for i, (icon, title, body) in enumerate(nexts):
    x = 0.45 + i * 3.18
    R(s, x, 3.98, 3.0, 2.52, W, DIV, Pt(1))
    R(s, x, 3.98, 3.0, 0.06, B)
    T(s, f"{icon}  {title}", x+0.15, 4.1, 2.8, 0.35, size=Pt(10.5), bold=True, color=INK)
    T(s, body,               x+0.15, 4.52, 2.8, 1.85, size=Pt(9.5), color=MID, wrap=True)

T(s, "Built with Next.js 14 · PostgreSQL · Prisma · Vercel  ·  Team Zootopia — 2025 Software Engineering Capstone",
  0.5, 6.62, 12.3, 0.38, size=Pt(8.5), color=SUB, align=PP_ALIGN.CENTER, italic=True)

footer(s, 12)

# ── Save ─────────────────────────────────────────────
out = r"C:\Users\ADMIN\Desktop\SE\SE_book_bridge\BookBridge_v3.pptx"
prs.save(out)
print("Saved:", out)
